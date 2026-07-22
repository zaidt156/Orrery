import io
import random
import subprocess
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

import pytest
from docx import Document
from docx.shared import Pt
from openpyxl import Workbook

from backend.features import filepreview
from backend.features import files as file_library


def _docx_bytes(text: str = "Safe Office fallback") -> bytes:
    document = Document()
    document.add_paragraph(text)
    stream = io.BytesIO()
    document.save(stream)
    return stream.getvalue()


def _pdf_bytes(text: str = "Valid PDF") -> bytes:
    from reportlab.pdfgen import canvas

    stream = io.BytesIO()
    document = canvas.Canvas(stream)
    document.drawString(72, 720, text)
    document.showPage()
    document.save()
    return stream.getvalue()


def test_qt_pdf_renderer_releases_its_input_and_returns_real_pngs():
    pytest.importorskip("PySide6.QtPdf")
    from reportlab.pdfgen import canvas

    source = io.BytesIO()
    document = canvas.Canvas(source)
    document.drawString(72, 720, "Webview-safe PDF preview")
    document.showPage()
    document.save()

    rendered = filepreview._render_pdf_pngs(source.getvalue())

    assert rendered is not None
    assert len(rendered.pages) == 1
    assert rendered.pages[0].startswith(b"\x89PNG\r\n\x1a\n")
    assert rendered.complete is True


def test_pdf_preview_uses_rendered_pages_instead_of_the_embedded_pdf_viewer(monkeypatch):
    monkeypatch.setattr(
        filepreview,
        "_render_pdf_pngs",
        lambda _data: filepreview._PdfRender((b"first page", b"second page"), 2, True),
    )

    content, media = filepreview.to_preview("resume.pdf", "application/pdf", b"%PDF-valid")

    html = content.decode("utf-8")
    assert media == "text/html; charset=utf-8"
    assert 'data-renderer="qt-pdf"' in html
    assert html.count("data:image/png;base64,") == 2
    assert "Page 1" in html
    assert "Page 2" in html


def test_pdf_preview_returns_bounded_html_when_page_rendering_is_unavailable(monkeypatch):
    monkeypatch.setattr(filepreview, "_render_pdf_pngs", lambda _data: None)

    content, media = filepreview.to_preview("resume.pdf", "application/pdf", b"%PDF-valid")

    assert media == "text/html; charset=utf-8"
    assert b"%PDF-valid" not in content
    assert b"PDF preview is unavailable" in content
    assert len(content) <= filepreview._MAX_PREVIEW_OUTPUT_BYTES


def test_converted_office_pdf_never_falls_back_to_native_pdf(monkeypatch):
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: b"%PDF-office")
    monkeypatch.setattr(filepreview, "_render_pdf_pngs", lambda _data: None)

    content, media = filepreview.to_preview(
        "resume.docx",
        "application/vnd.openxmlformats",
        _docx_bytes("Original Office fallback content"),
    )

    assert media == "text/html; charset=utf-8"
    assert b"%PDF-office" not in content
    assert b"Original Office fallback content" in content
    assert b"PDF preview is unavailable" not in content


def test_pdf_page_render_adapts_resolution_to_remaining_budget():
    attempted = []

    def high_entropy_png(width):
        attempted.append(width)
        return b"x" * (width * 2)

    rendered = filepreview._bounded_pdf_page_png(high_entropy_png, remaining=1_000)

    assert attempted[0] == 1000
    assert attempted[-1] < attempted[0]
    assert rendered is not None
    assert len(rendered) <= 1_000


def test_high_entropy_pdf_never_returns_application_pdf_when_png_budget_is_exhausted(monkeypatch):
    pytest.importorskip("PySide6.QtPdf")
    from PIL import Image
    from reportlab.pdfgen import canvas

    noise = random.Random(7).randbytes(600 * 600 * 3)
    image = Image.frombytes("RGB", (600, 600), noise)
    source = io.BytesIO()
    document = canvas.Canvas(source, pagesize=(600, 600))
    document.drawInlineImage(image, 0, 0, width=600, height=600)
    document.showPage()
    document.save()
    monkeypatch.setattr(filepreview, "_MAX_PDF_PREVIEW_PNG_BYTES", 8_000)

    content, media = filepreview.to_preview("noise.pdf", "application/pdf", source.getvalue())

    assert media == "text/html; charset=utf-8"
    assert len(content) <= filepreview._MAX_PREVIEW_OUTPUT_BYTES
    assert b"%PDF" not in content


def test_converted_office_pdf_uses_the_webview_safe_page_preview(monkeypatch):
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: b"%PDF-office")
    monkeypatch.setattr(
        filepreview,
        "_render_pdf_pngs",
        lambda _data: filepreview._PdfRender((b"office page",), 1, True),
    )

    content, media = filepreview.to_preview(
        "resume.docx",
        "application/vnd.openxmlformats",
        _docx_bytes(),
    )

    assert media == "text/html; charset=utf-8"
    assert b'data-renderer="qt-pdf"' in content


def test_pdf_page_limit_emits_partial_metadata_and_explicit_notice(monkeypatch):
    monkeypatch.setattr(
        filepreview,
        "_render_pdf_pngs",
        lambda _data: filepreview._PdfRender((b"first",), 3, False, "page limit"),
    )

    content, media = filepreview.to_preview("long.pdf", "application/pdf", b"%PDF-valid")

    assert media == "text/html; charset=utf-8"
    assert b'data-preview-complete="false"' in content
    assert b"1 of 3" in content
    assert b"Preview truncated for safety" in content


def test_pdf_byte_limit_with_zero_pages_still_emits_partial_metadata(monkeypatch):
    monkeypatch.setattr(
        filepreview,
        "_render_pdf_pngs",
        lambda _data: filepreview._PdfRender((), 2, False, "byte limit"),
    )

    content, media = filepreview.to_preview("dense.pdf", "application/pdf", b"%PDF-valid")

    assert media == "text/html; charset=utf-8"
    assert b'data-preview-complete="false"' in content
    assert b'data-preview-truncation-reason="byte-limit"' in content
    assert b"0 of 2" in content
    assert b"Preview truncated for safety" in content


def test_preview_status_requires_both_libreoffice_and_pdf_renderer(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "/opt/libreoffice/soffice")
    monkeypatch.setattr(filepreview, "pdf_renderer_available", lambda: False)

    status = filepreview.office_preview_status()

    assert status["available"] is False
    assert status["officePreview"] == "html"
    assert status["pdfRendererAvailable"] is False
    assert status["canInstall"] is False
    assert "renderer" in status["message"].lower()


def test_docx_basic_fallback_infers_direct_format_headings_and_bullets(monkeypatch):
    document = Document()
    title = document.add_paragraph()
    title_run = title.add_run("ETHAN J. WHITFIELD")
    title_run.bold = True
    title_run.font.size = Pt(25)
    section = document.add_paragraph()
    section_run = section.add_run("PROFESSIONAL SUMMARY")
    section_run.bold = True
    section_run.font.size = Pt(12.5)
    document.add_paragraph("A senior data engineer.")
    document.add_paragraph("▪   Built reliable data pipelines.")
    skills = document.add_table(rows=2, cols=2)
    skills.cell(0, 0).text = "Category"
    skills.cell(0, 1).text = "Skills"
    skills.cell(1, 0).text = "Languages"
    skills.cell(1, 1).text = "Python, SQL"
    experience = document.add_paragraph()
    experience_run = experience.add_run("PROFESSIONAL EXPERIENCE")
    experience_run.bold = True
    experience_run.font.size = Pt(12.5)
    stream = io.BytesIO()
    document.save(stream)
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: None)

    content, media = filepreview.to_preview("resume.docx", "application/vnd.openxmlformats", stream.getvalue())

    html = content.decode("utf-8")
    assert media == "text/html; charset=utf-8"
    assert "<h1>ETHAN J. WHITFIELD</h1>" in html
    assert "<h2>PROFESSIONAL SUMMARY</h2>" in html
    assert "<li>Built reliable data pipelines.</li>" in html
    assert html.index("<table>") < html.index("<h2>PROFESSIONAL EXPERIENCE</h2>")


def test_malformed_office_xhtml_payload_becomes_inert_escaped_html(monkeypatch):
    payload = b'<html xmlns="http://www.w3.org/1999/xhtml"><script>alert(1)</script><p>visible</p></html>'
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: None)

    content, media = filepreview.to_preview("broken.docx", "application/xhtml+xml", payload)

    text = content.decode("utf-8")
    assert media == "text/html; charset=utf-8"
    assert 'data-renderer="inert-fallback"' in text
    assert "<script>" not in text
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in text
    assert "Preview unavailable" in text


def test_active_xhtml_source_is_never_returned_raw():
    payload = b'<script>window.top.location="https://example.invalid"</script>'

    content, media = filepreview.to_preview("attack.xhtml", "application/xhtml+xml", payload)

    assert media == "text/html; charset=utf-8"
    assert content != payload
    assert b"&lt;script&gt;" in content
    assert b"<script>" not in content


def test_compressed_office_bomb_is_rejected_before_libreoffice_or_parser(monkeypatch):
    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as package:
        package.writestr("word/document.xml", b"A" * (filepreview._MAX_OFFICE_MEMBER_BYTES + 1))
    calls = []
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: calls.append(True))

    content, media = filepreview.to_preview(
        "compressed.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        archive.getvalue(),
    )

    assert media == "text/html; charset=utf-8"
    assert calls == []
    assert b"Preview truncated for safety" in content
    assert len(content) <= filepreview._MAX_PREVIEW_OUTPUT_BYTES


def test_source_fallback_bounds_input_characters_and_escapes_scripts(monkeypatch):
    monkeypatch.setattr(filepreview, "_MAX_PREVIEW_CHARS", 80)
    payload = b"<script>alert(1)</script>" + (b"A" * 5_000)

    content, media = filepreview.to_preview("large.tex", "text/plain", payload)

    assert media == "text/html; charset=utf-8"
    assert b"<script>" not in content
    assert b"&lt;script&gt;" in content
    assert b"Preview truncated for safety" in content
    assert len(content) <= filepreview._MAX_PREVIEW_OUTPUT_BYTES


def test_xlsx_fallback_bounds_cells_and_shows_truncation(monkeypatch):
    workbook = Workbook()
    sheet = workbook.active
    for row in range(1, 6):
        for column in range(1, 6):
            sheet.cell(row, column, f"cell-{row}-{column}")
    stream = io.BytesIO()
    workbook.save(stream)
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(filepreview, "_MAX_OFFICE_CELLS", 4)

    content, media = filepreview.to_preview("wide.xlsx", "application/vnd.openxmlformats", stream.getvalue())

    assert media == "text/html; charset=utf-8"
    assert b"Preview truncated for safety" in content
    assert content.count(b"<td>") + content.count(b"<th>") <= 4


def test_empty_xlsx_fallback_does_not_expand_a_synthetic_maximum_grid(monkeypatch):
    workbook = Workbook()
    stream = io.BytesIO()
    workbook.save(stream)
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: None)

    content, media = filepreview.to_preview("empty.xlsx", "application/vnd.openxmlformats", stream.getvalue())

    assert media == "text/html; charset=utf-8"
    assert b"Preview truncated for safety" not in content
    assert content.count(b"<td>") + content.count(b"<th>") <= 1


def test_docx_fallback_bounds_nodes_and_output(monkeypatch):
    document = Document()
    for index in range(12):
        document.add_paragraph(f"Paragraph {index} " + ("x" * 400))
    stream = io.BytesIO()
    document.save(stream)
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(filepreview, "_MAX_OFFICE_NODES", 3)
    monkeypatch.setattr(filepreview, "_MAX_PREVIEW_OUTPUT_BYTES", 6_000)

    content, media = filepreview.to_preview("large.docx", "application/vnd.openxmlformats", stream.getvalue())

    assert media == "text/html; charset=utf-8"
    assert b"Preview truncated for safety" in content
    assert content.count(b"<p>") <= 3
    assert len(content) <= 6_000


def test_office_preview_status_reports_pdf_when_libreoffice_is_available(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "/opt/libreoffice/soffice")
    monkeypatch.setattr(filepreview, "pdf_renderer_available", lambda: True)

    assert filepreview.office_preview_status() == {
        "available": True,
        "engine": "libreoffice",
        "officePreview": "pdf",
        "pdfRendererAvailable": True,
        "canInstall": False,
        "message": "Faithful Office previews are available.",
    }


def test_office_preview_status_reports_html_fallback_when_unavailable(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: None)
    monkeypatch.setattr(filepreview, "_installer_command", lambda: None)
    monkeypatch.setattr(filepreview, "pdf_renderer_available", lambda: True)

    assert filepreview.office_preview_status() == {
        "available": False,
        "engine": "libreoffice",
        "officePreview": "html",
        "pdfRendererAvailable": True,
        "canInstall": False,
        "message": "LibreOffice is not installed; Office files use the HTML fallback.",
    }


def test_office_preview_status_enables_install_when_package_manager_exists(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: None)
    monkeypatch.setattr(filepreview, "_installer_command", lambda: ["fixed-installer"])
    monkeypatch.setattr(filepreview, "pdf_renderer_available", lambda: True)

    assert filepreview.office_preview_status()["canInstall"] is True


def test_office_preview_status_never_offers_install_when_renderer_is_missing(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: None)
    monkeypatch.setattr(filepreview, "_installer_command", lambda: ["fixed-installer"])
    monkeypatch.setattr(filepreview, "pdf_renderer_available", lambda: False)

    assert filepreview.office_preview_status()["canInstall"] is False


def test_libreoffice_installer_uses_fixed_windows_package_id(monkeypatch):
    monkeypatch.setattr(filepreview.platform, "system", lambda: "Windows")
    monkeypatch.setattr(
        filepreview.shutil,
        "which",
        lambda name: r"C:\\Windows\\winget.exe" if name == "winget" else None,
    )

    command = filepreview._installer_command()

    assert command == [
        r"C:\\Windows\\winget.exe",
        "install",
        "--id",
        "TheDocumentFoundation.LibreOffice",
        "--exact",
        "--source",
        "winget",
        "--accept-package-agreements",
        "--accept-source-agreements",
        "--disable-interactivity",
        "--silent",
    ]


def test_libreoffice_installer_uses_fixed_macos_brew_cask(monkeypatch):
    monkeypatch.setattr(filepreview.platform, "system", lambda: "Darwin")
    monkeypatch.setattr(
        filepreview.proc,
        "find_executable",
        lambda name: "/opt/homebrew/bin/brew" if name == "brew" else None,
    )

    assert filepreview._installer_command() == [
        "/opt/homebrew/bin/brew",
        "install",
        "--cask",
        "libreoffice",
    ]


def test_install_libreoffice_requires_consent(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: None)

    with pytest.raises(ValueError, match="Confirm the LibreOffice installation"):
        filepreview.install_office_preview(False)


def test_install_libreoffice_refuses_when_packaged_renderer_is_missing(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: None)
    monkeypatch.setattr(filepreview, "pdf_renderer_available", lambda: False)

    with pytest.raises(ValueError, match="PDF preview renderer"):
        filepreview.install_office_preview(True)


def test_install_libreoffice_runs_fixed_command_then_reprobes(monkeypatch):
    installed = False
    command = ["winget.exe", "install", "--id", "TheDocumentFoundation.LibreOffice"]
    calls = []

    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "soffice.exe" if installed else None)
    monkeypatch.setattr(filepreview, "_installer_command", lambda: command)

    def fake_run(argv, **kwargs):
        nonlocal installed
        calls.append((argv, kwargs))
        installed = True
        return subprocess.CompletedProcess(argv, 0, "", "")

    monkeypatch.setattr(filepreview.proc, "run", fake_run)

    status = filepreview.install_office_preview(True)

    assert calls[0][0] == command
    assert calls[0][1]["check"] is False
    assert status["available"] is True
    assert status["officePreview"] == "pdf"


def test_install_libreoffice_fails_closed_without_supported_package_manager(monkeypatch):
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: None)
    monkeypatch.setattr(filepreview, "_installer_command", lambda: None)

    with pytest.raises(ValueError, match="WinGet on Windows or Homebrew on macOS"):
        filepreview.install_office_preview(True)


def test_office_pdf_reuses_cached_conversion(monkeypatch, tmp_path):
    cache_path = tmp_path / "artifact.digest.preview.pdf"
    calls = []

    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "soffice")

    valid_pdf = _pdf_bytes("cached")

    def fake_run(argv, **_kwargs):
        calls.append(argv)
        source = Path(argv[-1])
        source.with_suffix(".pdf").write_bytes(valid_pdf)

    monkeypatch.setattr(filepreview.proc, "run", fake_run)

    assert filepreview._office_pdf("deck.pptx", b"pptx", cache_path=cache_path) == valid_pdf
    assert cache_path.read_bytes() == valid_pdf
    assert filepreview._office_pdf("deck.pptx", b"pptx", cache_path=cache_path) == valid_pdf
    assert len(calls) == 1


def test_office_pdf_uses_a_unique_libreoffice_profile_per_conversion(monkeypatch):
    profiles = []
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "soffice")

    valid_pdf = _pdf_bytes("profile")

    def fake_run(argv, **_kwargs):
        profiles.append(next(arg for arg in argv if arg.startswith("-env:UserInstallation=")))
        source = Path(argv[-1])
        source.with_suffix(".pdf").write_bytes(valid_pdf)

    monkeypatch.setattr(filepreview.proc, "run", fake_run)

    assert filepreview._office_pdf("first.docx", b"docx") == valid_pdf
    assert filepreview._office_pdf("second.docx", b"docx") == valid_pdf
    assert len(set(profiles)) == 2
    assert all(profile.startswith("-env:UserInstallation=file:") for profile in profiles)


def test_failed_office_conversion_does_not_poison_cache(monkeypatch, tmp_path):
    cache_path = tmp_path / "artifact.digest.preview.pdf"
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "soffice")

    def failed_run(*_args, **_kwargs):
        raise RuntimeError("conversion failed")

    monkeypatch.setattr(filepreview.proc, "run", failed_run)

    assert filepreview._office_pdf("document.docx", b"docx", cache_path=cache_path) is None
    assert not cache_path.exists()


def test_oversized_fresh_office_pdf_is_rejected_before_read_or_cache(monkeypatch, tmp_path):
    cache_path = tmp_path / "artifact.digest.preview.pdf"
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "soffice")
    monkeypatch.setattr(filepreview, "_MAX_CACHED_OFFICE_PDF_BYTES", 4)

    def fake_run(argv, **_kwargs):
        source = Path(argv[-1])
        source.with_suffix(".pdf").write_bytes(b"12345")

    monkeypatch.setattr(filepreview.proc, "run", fake_run)

    assert filepreview._office_pdf("deck.pptx", b"pptx", cache_path=cache_path) is None
    assert not cache_path.exists()


def test_oversized_cached_office_pdf_is_rejected_before_read(monkeypatch, tmp_path):
    cache_path = tmp_path / "artifact.digest.preview.pdf"
    cache_path.write_bytes(b"x" * 5)
    monkeypatch.setattr(filepreview, "_MAX_CACHED_OFFICE_PDF_BYTES", 4)
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: None)

    assert filepreview._office_pdf("deck.pptx", b"pptx", cache_path=cache_path) is None


def test_invalid_pdf_signature_or_structure_is_not_cached(monkeypatch, tmp_path):
    cache_path = tmp_path / "artifact.digest.preview.pdf"
    monkeypatch.setattr(filepreview, "_find_soffice", lambda: "soffice")

    def fake_run(argv, **_kwargs):
        Path(argv[-1]).with_suffix(".pdf").write_bytes(b"%PDF-1.7\nnot a real document\n%%EOF")

    monkeypatch.setattr(filepreview.proc, "run", fake_run)

    assert filepreview._office_pdf("deck.pptx", b"pptx", cache_path=cache_path) is None
    assert not cache_path.exists()


def test_preview_jobs_are_process_serialized_on_one_cpu(monkeypatch):
    monkeypatch.setattr(filepreview, "_preview_slot", filepreview.threading.BoundedSemaphore(1))
    active = 0
    maximum = 0
    lock = filepreview.threading.Lock()

    def worker(key):
        def job():
            nonlocal active, maximum
            with lock:
                active += 1
                maximum = max(maximum, active)
            time.sleep(0.04)
            with lock:
                active -= 1
            return key

        return filepreview._run_preview_job(key, job)

    with ThreadPoolExecutor(max_workers=3) as pool:
        results = list(pool.map(worker, ("one", "two", "three")))

    assert results == ["one", "two", "three"]
    assert maximum == 1


def test_preview_jobs_singleflight_identical_sources(monkeypatch):
    monkeypatch.setattr(filepreview, "_preview_slot", filepreview.threading.BoundedSemaphore(1))
    calls = 0

    def worker():
        def job():
            nonlocal calls
            calls += 1
            time.sleep(0.05)
            return b"shared result"

        return filepreview._run_preview_job("same-source", job)

    with ThreadPoolExecutor(max_workers=4) as pool:
        results = list(pool.map(lambda _index: worker(), range(4)))

    assert results == [b"shared result"] * 4
    assert calls == 1


@pytest.mark.parametrize(
    ("name", "members", "expected"),
    [
        (
            "external.docx",
            {
                "[Content_Types].xml": b"<Types/>",
                "_rels/.rels": b'<Relationships><Relationship TargetMode="External" Target="https://evil.invalid"/></Relationships>',
            },
            b"external relationship",
        ),
        (
            "macro.docx",
            {"[Content_Types].xml": b"<Types/>", "word/vbaProject.bin": b"macro"},
            b"macro",
        ),
        (
            "encrypted.docx",
            {"EncryptionInfo": b"encrypted", "EncryptedPackage": b"payload"},
            b"encrypted",
        ),
    ],
)
def test_active_or_encrypted_ooxml_is_rejected_before_libreoffice(monkeypatch, name, members, expected):
    package = io.BytesIO()
    with zipfile.ZipFile(package, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for member, value in members.items():
            archive.writestr(member, value)
    calls = []
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: calls.append(True))

    content, media = filepreview.to_preview(name, "application/vnd.openxmlformats", package.getvalue())

    assert media == "text/html; charset=utf-8"
    assert calls == []
    assert expected in content.lower()


def test_xlsm_is_rejected_before_libreoffice_even_without_vba_part(monkeypatch):
    calls = []
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: calls.append(True))

    content, media = filepreview.to_preview("macro.xlsm", "application/vnd.ms-excel.sheet.macroenabled.12", _docx_bytes())

    assert media == "text/html; charset=utf-8"
    assert calls == []
    assert b"macro-enabled" in content.lower()


def test_ole_encrypted_package_is_rejected_before_libreoffice(monkeypatch):
    calls = []
    payload = bytes.fromhex("D0CF11E0A1B11AE1") + (b"encrypted" * 20)
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *_args, **_kwargs: calls.append(True))

    content, media = filepreview.to_preview("encrypted.docx", "application/vnd.openxmlformats", payload)

    assert media == "text/html; charset=utf-8"
    assert calls == []
    assert b"encrypted" in content.lower()


def test_preview_cache_path_is_content_addressed_and_removes_stale_variants(monkeypatch, tmp_path):
    file_id = "a" * 32
    monkeypatch.setattr(file_library, "_DIR", tmp_path)

    first = file_library.office_preview_cache_path(file_id, b"first version")
    first.write_bytes(b"old preview")
    second = file_library.office_preview_cache_path(file_id, b"second version")

    assert first != second
    assert not first.exists()
    assert second.parent == tmp_path
    assert second.name.startswith(f"{file_id}.")
    assert second.name.endswith(".preview.pdf")


def test_preview_cache_path_prunes_oldest_entries_to_a_fixed_limit(monkeypatch, tmp_path):
    monkeypatch.setattr(file_library, "_DIR", tmp_path)
    monkeypatch.setattr(file_library, "_MAX_OFFICE_PREVIEW_CACHE_ITEMS", 2)
    old = tmp_path / f"{'1' * 32}.old.preview.pdf"
    newer = tmp_path / f"{'2' * 32}.new.preview.pdf"
    old.write_bytes(b"old")
    newer.write_bytes(b"newer")
    old.touch()
    newer.touch()
    old_mtime = old.stat().st_mtime - 10
    import os

    os.utime(old, (old_mtime, old_mtime))

    target = file_library.office_preview_cache_path("3" * 32, b"new")
    target.write_bytes(b"new")

    assert not old.exists()
    assert newer.exists()
    assert target.exists()
    assert len(list(tmp_path.glob("*.preview.pdf"))) == 2


# --- faithful Python-library previews (Workstream 2: documents render, not text dumps) ---

import base64 as _b64

_PNG_1PX = _b64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _preview_html(name, mime, data):
    content, media = filepreview.to_preview(name, mime, data)
    assert media.startswith("text/html")
    return content.decode("utf-8")


def test_csv_previews_as_a_table_not_raw_text():
    page = _preview_html("results.csv", "text/csv", b"name,score\nAda,99\nGrace,97\n")
    assert "<table>" in page
    assert "<pre>" not in page
    assert "<th>name</th>" in page
    assert "<td>Ada</td>" in page


def test_tsv_previews_as_a_table():
    page = _preview_html("data.tsv", "text/tab-separated-values", b"a\tb\n1\t2\n")
    assert "<th>a</th>" in page and "<td>2</td>" in page


def test_markdown_renders_with_raw_html_escaped():
    data = "# Title\n\nSome **bold** text.\n\n<script>alert(1)</script>\n".encode()
    page = _preview_html("notes.md", "text/markdown", data)
    assert "<h1>Title</h1>" in page
    assert "<strong>bold</strong>" in page
    assert "<script>alert" not in page  # raw HTML is escaped, never executed


def test_markdown_never_fetches_remote_images():
    page = _preview_html("readme.md", "text/markdown", b"![logo](https://evil.example/x.png)\n")
    assert 'src="https://' not in page


def test_xlsx_fallback_keeps_merges_styles_and_widths(monkeypatch):
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *a, **k: None)  # force the Python path
    from openpyxl import Workbook as _Workbook
    from openpyxl.styles import Font, PatternFill

    workbook = _Workbook()
    sheet = workbook.active
    sheet.title = "Report"
    sheet.merge_cells("A1:B1")
    header = sheet["A1"]
    header.value = "Quarter"
    header.font = Font(bold=True, color="FFFFFFFF")
    header.fill = PatternFill(fill_type="solid", start_color="FF26314F")
    sheet["A2"] = "North"
    sheet["B2"] = 1234.5
    sheet.column_dimensions["A"].width = 24
    stream = io.BytesIO()
    workbook.save(stream)

    page = _preview_html("report.xlsx", "application/octet-stream", stream.getvalue())
    assert 'colspan="2"' in page
    assert "font-weight:700" in page
    assert "background:#26314F" in page
    assert "1234.5" in page
    assert "<col style=" in page


def test_docx_fallback_keeps_run_formatting_and_images(monkeypatch):
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *a, **k: None)
    from docx.shared import RGBColor

    document = Document()
    paragraph = document.add_paragraph()
    bold_run = paragraph.add_run("Important")
    bold_run.bold = True
    color_run = paragraph.add_run(" note")
    color_run.font.color.rgb = RGBColor(0x33, 0x66, 0xFF)
    document.add_picture(io.BytesIO(_PNG_1PX), width=Pt(72))
    stream = io.BytesIO()
    document.save(stream)

    page = _preview_html("memo.docx", "application/octet-stream", stream.getvalue())
    assert "<strong>Important</strong>" in page
    assert "color:#3366FF" in page
    assert '<img src="data:image/png;base64,' in page


def test_pptx_fallback_positions_shapes_and_inlines_pictures(monkeypatch):
    monkeypatch.setattr(filepreview, "_office_pdf", lambda *a, **k: None)
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.util import Pt as PptxPt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Hello deck"
    run.font.bold = True
    run.font.size = PptxPt(32)
    slide.shapes.add_picture(io.BytesIO(_PNG_1PX), Inches(6), Inches(2), Inches(1), Inches(1))
    stream = io.BytesIO()
    prs.save(stream)

    page = _preview_html("deck.pptx", "application/octet-stream", stream.getvalue())
    assert "<strong>Hello deck</strong>" in page
    assert 'class="shape"' in page
    assert "left:" in page and "top:" in page  # shapes keep their slide position
    assert '<img src="data:image/png;base64,' in page
