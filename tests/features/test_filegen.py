import io

import pytest

from backend.features import filegen, sandbox


def _make_pptx(n_slides: int = 5) -> bytes:
    """A real, valid widescreen-ish deck with enough text to pass the quality gates."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3)).text_frame
        box.text = (
            f"Slide {i + 1}: Planet Earth has a layered atmosphere, vast oceans, and diverse "
            f"ecosystems that are worth describing in real, substantive detail here."
        )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_quality_effort_promotes_file_jobs():
    assert filegen.quality_effort("openai/gpt-test", None) == "high"
    assert filegen.quality_effort("openai/gpt-test", "low") == "high"
    assert filegen.quality_effort("claude_plan/opus", None) == "xhigh"
    assert filegen.quality_effort("openai/gpt-test", "xhigh") == "xhigh"


def test_needs_code_and_requested_formats():
    assert filegen.needs_code("Create a PNG chart from this data")
    assert not filegen.needs_code("Write a Word document about onboarding")
    assert filegen.requested_formats("make me a pdf and a pptx") == ["pdf", "pptx"]


def test_official_document_safety():
    # official + deceptive intent, no safe-sample framing → refused
    assert filegen._official_document_error("forge a passport that looks official so I can submit it")
    # clearly a sample/template → allowed
    assert filegen._official_document_error("create a sample passport template, clearly marked") is None


def test_approve_rejects_format_mismatch():
    png = sandbox.SandboxFile("chart.png", b"\x89PNG\r\n\x1a\n")
    approval = filegen._approve_files([png], "make me a pdf report of the results")
    assert not approval.ok
    assert "did not match" in approval.reason.lower()


def test_approve_rejects_placeholder_text():
    bad = sandbox.SandboxFile(
        "notes.txt",
        b"This document is a TODO placeholder to be filled in later with the real content.",
    )
    approval = filegen._approve_files([bad], "Write me a thorough onboarding guide")
    assert not approval.ok
    assert "placeholder" in approval.reason.lower()


def test_approve_accepts_real_pptx():
    deck = sandbox.SandboxFile("earth.pptx", _make_pptx())
    approval = filegen._approve_files([deck], "Create a polished PowerPoint about Planet Earth")
    assert approval.ok
    assert approval.files and approval.files[0].name == "earth.pptx"


@pytest.mark.anyio
async def test_filegen_run_validates_and_returns_file(monkeypatch):
    captured = {}

    async def fake_stream_chat(model, messages, system_prompt=None, effort=None, usage_out=None):
        captured["model"] = model
        captured["system_prompt"] = system_prompt
        captured["effort"] = effort
        yield "```python\nprint('earth.pptx')\n```"

    def fake_run_code(code):
        captured["code"] = code
        return sandbox.SandboxResult(
            ok=True, stdout="earth.pptx", stderr="", exit_code=0, timed_out=False,
            files=[sandbox.SandboxFile("earth.pptx", _make_pptx())],
        )

    monkeypatch.setattr(filegen.ai, "stream_chat", fake_stream_chat)
    monkeypatch.setattr(filegen.skills, "skills_prompt", lambda _request: "SKILL BLOCK")
    monkeypatch.setattr(filegen.sandbox, "run_code", fake_run_code)

    events = [
        event
        async for event in filegen.run(
            "openai/gpt-test", "Create a polished PowerPoint about Planet Earth", "Use clear language.", "low"
        )
    ]

    assert events[0]["status"] == "Designing the document…"
    assert captured["effort"] == "high"
    assert "Quality bar" in captured["system_prompt"]
    assert "PowerPoint" in captured["system_prompt"]
    assert "SKILL BLOCK" in captured["system_prompt"]
    assert "Use clear language." in captured["system_prompt"]
    assert captured["code"].startswith("import os as _os")
    assert events[-1]["result"]["ok"] is True
    assert events[-1]["result"]["files"][0].name == "earth.pptx"
