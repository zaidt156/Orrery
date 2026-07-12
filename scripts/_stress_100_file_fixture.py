"""Real 100-file FileGen + project-context lifecycle stress test.

This is intentionally a manual test: it needs Docker, the Orrery sandbox image, and the local
Postgres/pgvector database.  It never calls an AI provider.  Instead, deterministic Python runs
inside the same locked-down sandbox used by FileGen, then exercises validation, local storage,
previewing, project ingestion, retrieval, isolation, deletion, and cleanup.

Run:  .venv/Scripts/python scripts/stress_100_file_lifecycle.py
"""
from __future__ import annotations

import asyncio
import pathlib
import sys

if sys.platform == "win32":
    asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))

from backend.core import database
from backend.features import filegen, filepreview, projects, rag, sandbox, taskrouter
from backend.features import files as file_library
from backend.features.chat import retrieval

PREFIX = "stress-100-file-lifecycle-"
FILES_PER_FORMAT = 5
FORMATS = (
    "pdf", "docx", "xlsx", "pptx", "csv", "tex", "png", "jpg", "gif", "webp",
    "svg", "wav", "mp3", "mp4", "webm", "zip", "html", "md", "txt", "json",
)
TEXT_FORMATS = {"csv", "tex", "html", "md", "txt", "json"}
IMAGE_FORMATS = {"png", "jpg", "gif", "webp", "svg"}
MIME = {
    "pdf": "application/pdf", "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "csv": "text/csv", "tex": "application/x-tex", "png": "image/png", "jpg": "image/jpeg",
    "gif": "image/gif", "webp": "image/webp", "svg": "image/svg+xml", "wav": "audio/wav",
    "mp3": "audio/mpeg", "mp4": "video/mp4", "webm": "video/webm", "zip": "application/zip",
    "html": "text/html", "md": "text/markdown", "txt": "text/plain", "json": "application/json",
}

_GENERATOR = r'''
from pathlib import Path
import csv, json, math, struct, subprocess, wave, zipfile

out = Path("/work/out")
ext = FMT
filler = (
    "This artifact records a deterministic Orrery lifecycle fact. "
    "Its contents are intentionally substantive enough for validation and retrieval. "
)

for offset in range(COUNT):
    idx = START + offset
    marker = f"{FMT.upper()}-{1000 + idx}"
    fact = f"The {FMT} stress marker for unit {idx} is {marker}."
    path = out / f"stress-{FMT}-{idx:03d}.{ext}"

    if FMT == "pdf":
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(path)); c.setFont("Helvetica", 12)
        for line_no, line in enumerate((fact, filler, filler, "Verified by the Orrery 100-file lifecycle run.")):
            c.drawString(54, 760 - line_no * 28, line)
        c.save()
    elif FMT == "docx":
        from docx import Document
        doc = Document(); doc.add_heading("Orrery lifecycle verification", 0)
        doc.add_paragraph(fact); doc.add_paragraph(filler); doc.add_paragraph(filler)
        doc.save(path)
    elif FMT == "xlsx":
        from openpyxl import Workbook
        wb = Workbook(); ws = wb.active; ws.title = "Verification"
        ws.append(["unit", "marker", "description"]); ws.append([idx, marker, fact]); ws.append([idx, marker, filler])
        wb.save(path)
    elif FMT == "pptx":
        from pptx import Presentation
        prs = Presentation()
        for slide_no in range(4):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Orrery verification {slide_no + 1}"
            slide.placeholders[1].text = fact + " " + filler
        prs.save(path)
    elif FMT == "csv":
        with path.open("w", newline="", encoding="utf-8") as handle:
            writer = csv.writer(handle); writer.writerow(["unit", "marker", "description"])
            writer.writerow([idx, marker, fact]); writer.writerow([idx, marker, filler])
    elif FMT == "tex":
        path.write_text(r"\documentclass{article}\begin{document}\section{Orrery Verification}" + fact + " " + filler + r"\end{document}", encoding="utf-8")
    elif FMT in {"png", "jpg", "gif", "webp"}:
        from PIL import Image, ImageDraw
        image = Image.new("RGB", (720, 240), (16, 27, 54)); draw = ImageDraw.Draw(image)
        draw.text((30, 40), fact, fill=(240, 205, 120)); draw.text((30, 100), "Orrery 100-file lifecycle", fill=(180, 210, 255))
        image.save(path, format={"jpg": "JPEG"}.get(FMT, FMT.upper()))
    elif FMT == "svg":
        path.write_text(f'<svg xmlns="http://www.w3.org/2000/svg" width="720" height="240"><rect width="100%" height="100%" fill="#101b36"/><text x="30" y="80" fill="#f0cd78">{fact}</text><text x="30" y="130" fill="#b4d2ff">Orrery lifecycle verification</text></svg>', encoding="utf-8")
    elif FMT == "wav":
        with wave.open(str(path), "wb") as audio:
            rate = 8000; audio.setnchannels(1); audio.setsampwidth(2); audio.setframerate(rate)
            audio.writeframes(b"".join(struct.pack("<h", int(9000 * math.sin(2 * math.pi * (300 + idx) * n / rate))) for n in range(rate // 2)))
    elif FMT == "mp3":
        subprocess.run(["ffmpeg", "-loglevel", "error", "-f", "lavfi", "-i", f"sine=frequency={400 + idx}:duration=0.4", "-metadata", f"comment={fact}", "-y", str(path)], check=True)
    elif FMT in {"mp4", "webm"}:
        codec = "mpeg4" if FMT == "mp4" else "libvpx"
        source = "color=c=0x101b36:s=160x120:d=0.4" if FMT == "mp4" else "testsrc=size=320x240:rate=12:duration=1"
        subprocess.run(["ffmpeg", "-loglevel", "error", "-f", "lavfi", "-i", source, "-c:v", codec, "-metadata", f"comment={fact}", "-y", str(path)], check=True)
    elif FMT == "zip":
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
            archive.writestr(f"facts/unit-{idx}.txt", fact + " " + filler)
            archive.writestr("README.md", "Orrery lifecycle verification archive.")
    elif FMT == "html":
        path.write_text(f'<!doctype html><html><head><style>body{{font-family:sans-serif;background:#101b36;color:#eef}}</style></head><body><main><h1>Orrery lifecycle</h1><p>{fact} {filler}</p><button onclick="document.body.dataset.checked=\'yes\'">Verify</button></main></body></html>', encoding="utf-8")
    elif FMT == "md":
        path.write_text(f"# Orrery lifecycle verification\n\n{fact}\n\n{filler}{filler}", encoding="utf-8")
    elif FMT == "txt":
        path.write_text(f"Orrery lifecycle verification\n{fact}\n{filler}{filler}", encoding="utf-8")
    elif FMT == "json":
        path.write_text(json.dumps({"unit": idx, "marker": marker, "fact": fact, "description": filler}), encoding="utf-8")
'''
