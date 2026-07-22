"""file.generate — produce any file by having the model WRITE CODE that we run in the sandbox.

The loop: select skills → ask the model for one Python program (guided by the skill) → run it in
the locked-down sandbox → validate the generated files in the backend → if it errors, produces no
file, or fails quality checks, feed the failure reason back and retry (bounded) → return approved
files from ./out.

This is the open-source "code interpreter" mechanism; no model code ever runs in the backend
(see backend/features/sandbox.py and ARCHITECTURE.md section 9).
"""

from __future__ import annotations

import ast
import asyncio
import csv
import io
import json
import posixpath
import re
import wave
import zipfile
from collections.abc import AsyncIterator
from dataclasses import dataclass, field
from pathlib import PurePosixPath
from urllib.parse import unquote, urlsplit

from backend.features import events as stream_events
from backend.features import files as file_library
from backend.features import reasoning, sandbox, skills
from backend.features.prompting import FILE_SYSTEM_PROMPT, build_system_prompt
from backend.features.reasoning_trace import ThinkStream, reasoning_event
from backend.providers import ai

MAX_ATTEMPTS = 3  # default repair budget; the active reasoning mode can widen it (see reasoning.file_retries)
QUALITY_FILE_EFFORT = "high"
QUALITY_CLAUDE_PLAN_EFFORT = "xhigh"
_CODE_FENCE = re.compile(r"```(?:python|py)?\s*\n([\s\S]*?)```", re.IGNORECASE)

# strong, precise file intent: a creation verb is not required — naming a concrete artifact is enough
_FILE_INTENT = re.compile(
    r"\b(pdf|docx|word\s+doc(?:ument)?|excel|xlsx|spreadsheet|workbook|powerpoint|pptx|"
    r"presentation|slide\s*deck|slides?|deck|csv|chart|graph|plot|diagram|infographic|"
    r"invoice|resume|cv|brochure|flyer|certificate|tex|latex|html|web\s?page|webpage|website|"
    r"landing page|single[-\s]?page app|audio|sound|sound file|sound effect|voiceover|"
    r"voice-over|narration|text[-\s]?to[-\s]?speech|tts|speech|video|movie|animation)\b"
    # a bare ".ext" mention ("five .png files"): \b never matches between a space and a dot,
    # so the extension alternative must live OUTSIDE the word-boundary group above
    r"|\.(?:pdf|docx?|xlsx?|pptx?|csv|tex|png|jpe?g|gif|webp|svg|zip|wav|mp3|mp4|webm|html?|md|txt|json)\b",
    re.IGNORECASE,
)
_CREATE_VERB = re.compile(
    r"\b(create|make|generate|build|write|compose|give\s+me|need|want|produce|export|draft|design|prepare|"
    r"put\s+together|write\s+me|turn\s+.*\binto|as\s+an?)\b",
    re.IGNORECASE,
)

_APP_INTENT = re.compile(
    r"\b(?:small|mini|tiny|quick|throwaway|one[-\s]?(?:time|off))\s+(?:web\s+)?"
    r"(?:[a-z0-9][a-z0-9_-]*\s+){0,3}app\b|"
    r"\b(?:quick|one[-\s]?(?:time|off))\s+(?:tool|calculator|tracker|planner|converter|utility)\b",
    re.IGNORECASE,
)

# File requests that benefit from the richer code-execution path. PRESENTATIONS go here too:
# the model designs the deck freely (varied layouts, visuals) instead of the fixed docgen template —
# docgen remains the fast fallback. Plain Word/Excel/PDF docs still route to docgen first.
_NEEDS_CODE = re.compile(
    r"\b(powerpoint|pptx|presentation|slide\s*deck|slides?|deck|"
    r"tex|latex|"
    r"html|web\s?page|webpage|website|landing page|single[-\s]?page app|interactive|"
    r"chart|graph|plot|diagram|figure|visuali[sz]|infographic|image|picture|photo|logo|icon|"
    r"video|movie|animation|mp4|webm|"
    r"audio|sound|soundtrack|sound effect|sfx|tone|beep|voiceover|voice-over|narration|wav|mp3|"
    r"calculat|comput|analy[sz]|statistic|regression|simulat|forecast|matplotlib|seaborn|"
    r"\.(tex|png|jpe?g|gif|webp|svg|zip|tar|html?|mp4|webm))\b",
    re.IGNORECASE,
)

_FORMAT_PATTERNS: tuple[tuple[str, re.Pattern[str]], ...] = (
    # Deliverable-first: a web app/page ask often *mentions* other formats as features inside it
    # ("… and a Download CSV button") — the page is the deliverable, not the button's format.
    ("html", re.compile(r"\b(web\s?app|web\s?page|webpage|website|landing page|single[-\s]?page app)\b", re.I)),
    ("pdf", re.compile(r"\bpdf\b|\.pdf\b", re.I)),
    ("docx", re.compile(r"\b(word\s+doc(?:ument)?|docx?)\b|\.docx?\b", re.I)),
    ("xlsx", re.compile(r"\b(excel|xlsx?|spreadsheet|workbook|sheet)\b|\.xlsx?\b", re.I)),
    ("pptx", re.compile(r"\b(powerpoint|pptx?|presentation|slide\s*deck|slides?|deck)\b|\.pptx?\b", re.I)),
    ("csv", re.compile(r"\bcsv\b|\.csv\b", re.I)),
    ("tex", re.compile(r"\b(tex|latex|latex\s+source|latex\s+document|latex\s+template)\b|\.tex\b", re.I)),
    ("png", re.compile(r"\bpng\b|\.png\b", re.I)),
    ("jpg", re.compile(r"\bjpe?g\b|\.jpe?g\b", re.I)),
    ("gif", re.compile(r"\bgif\b|\.gif\b", re.I)),
    ("webp", re.compile(r"\bwebp\b|\.webp\b", re.I)),
    ("svg", re.compile(r"\bsvg\b|\.svg\b", re.I)),
    ("wav", re.compile(r"\b(wav|sound effect|sfx|tone|beep)\b|\.wav\b", re.I)),
    ("mp3", re.compile(r"\bmp3\b|\.mp3\b", re.I)),
    ("mp4", re.compile(r"\b(mp4|video|movie)\b|\.mp4\b", re.I)),
    ("webm", re.compile(r"\bwebm\b|\.webm\b", re.I)),
    ("zip", re.compile(r"\bzip\b|\.zip\b", re.I)),
    ("html", re.compile(r"\b(html|web\s?page|webpage|website|landing page|single[-\s]?page app)\b|\.html?\b", re.I)),
    ("md", re.compile(r"\b(markdown|md)\b|\.md\b", re.I)),
    ("txt", re.compile(r"\b(text file|txt)\b|\.txt\b", re.I)),
    ("json", re.compile(r"\bjson\b|\.json\b", re.I)),
)

_EXTENSION_TO_FORMAT = {
    "pdf": "pdf",
    "doc": "docx",
    "docx": "docx",
    "xls": "xlsx",
    "xlsx": "xlsx",
    "xlsm": "xlsx",
    "ppt": "pptx",
    "pptx": "pptx",
    "csv": "csv",
    "tex": "tex",
    "png": "png",
    "jpg": "jpg",
    "jpeg": "jpg",
    "gif": "gif",
    "webp": "webp",
    "svg": "svg",
    "wav": "wav",
    "mp3": "mp3",
    "mp4": "mp4",
    "webm": "webm",
    "zip": "zip",
    "md": "md",
    "txt": "txt",
    "html": "html",
    "json": "json",
}

_PLACEHOLDER_RE = re.compile(
    r"\b(todo|lorem ipsum|placeholder|insert text here|sample text|tbd|your name|company name|"
    r"\[title\]|\[date\]|\[name\]|\[company\])\b",
    re.IGNORECASE,
)
_REMOTE_HTML_REF_RE = re.compile(
    r"""(?:src|href)\s*=\s*['"]\s*(?:https?:|//|file:|javascript:)""",
    re.IGNORECASE,
)
_HTML_SCRIPT_REMOTE_RE = re.compile(r"<script\b[^>]*\bsrc\s*=", re.IGNORECASE)
_APP_NETWORK_SCHEME_RE = re.compile(r"\b(?:https?|wss?|ftp|file):", re.IGNORECASE)
_APP_JS_NETWORK_RE = re.compile(
    r"\b(?:fetch|XMLHttpRequest|WebSocket|EventSource|importScripts)\s*\(|"
    r"\bnavigator\s*\.\s*sendBeacon\s*\(|\bwindow\s*\.\s*open\s*\(",
    re.IGNORECASE,
)
_APP_JS_IMPORT_RE = re.compile(
    r"(?:\b(?:import|export)\s+(?:[^;\n]*?\s+from\s+)?|\bimport\s*\()\s*['\"]([^'\"]+)['\"]",
    re.IGNORECASE,
)
_APP_CSS_URL_RE = re.compile(r"url\(\s*(['\"]?)(.*?)\1\s*\)", re.IGNORECASE)
_APP_CSS_IMPORT_RE = re.compile(r"@import\s+(?:url\(\s*)?['\"]([^'\"]+)['\"]", re.IGNORECASE)
_APP_HTML_SINGLE_URL_ATTRS = (
    "background",
    "cite",
    "classid",
    "code",
    "codebase",
    "data",
    "dynsrc",
    "href",
    "icon",
    "longdesc",
    "lowsrc",
    "manifest",
    "poster",
    "profile",
    "src",
    "usemap",
    "xlink:href",
)
_APP_HTML_URL_LIST_ATTRS = ("archive", "ping")
_APP_HTML_SRCSET_ATTRS = ("imagesrcset", "srcset")
_APP_HTML_CSS_VALUE_ATTRS = (
    "clip-path",
    "color-profile",
    "cursor",
    "fill",
    "filter",
    "marker",
    "marker-end",
    "marker-mid",
    "marker-start",
    "mask",
    "stroke",
)
_APP_ALLOWED_EXTENSIONS = {
    "css", "gif", "html", "ico", "jpeg", "jpg", "js", "json", "mjs", "mp3", "mp4",
    "otf", "png", "svg", "ttf", "txt", "wav", "webm", "webp", "woff", "woff2",
}
_APP_BUNDLE_PROMPT = """

STATIC APP BUNDLE MODE (this overrides the single-file HTML rule above):
- Build a client-side-only app in ./out/<short-app-name>/.
- The directory must contain index.html, at least one local .js file, and at least one local .css file.
- Put any images/fonts/media in nested local asset directories. Do not create the ZIP yourself.
- index.html must load the local JS/CSS files with relative paths. Put no JavaScript inline and use
  no inline event-handler attributes; bind events from the local JS file.
- Do not use fetch, XMLHttpRequest, WebSocket, EventSource, service workers, external URLs, CDNs,
  remote fonts, form actions, browser storage, or dependencies that need a server/build step.
- The app must work by opening index.html from Orrery's read-only preview route.
"""
_LATEX_STRUCTURE_RE = re.compile(
    r"\\(?:documentclass|begin\s*\{\s*document\s*\}|section|subsection|title|author|"
    r"usepackage|begin\s*\{\s*(?:equation|align|tabular|itemize|enumerate)\s*\})",
    re.IGNORECASE,
)
_LATEX_UNSAFE_RE = re.compile(
    r"\\(?:write18|openout|read|input|include|includegraphics)\s*(?:\{|\s+)(?:/|[A-Za-z]:|\\.\\.|~)",
    re.IGNORECASE,
)

_SMALL_ARTIFACT_OK = re.compile(
    r"\b(one\s+page|1\s+page|one\s+slide|1\s+slide|single\s+slide|thumbnail|icon|logo|"
    r"template|blank|sample|example|draft|short|simple)\b",
    re.IGNORECASE,
)

_OFFICIAL_DOCUMENT_RE = re.compile(
    r"\b(medical certificate|doctor'?s?\s+note|sick note|health certificate|degree certificate|"
    r"diploma|transcript|residence permit|visa|passport|identity card|id card|driver'?s?\s+license|"
    r"bank statement|payslip|pay slip|tax document|government form|immigration document|"
    r"employment contract|signed letter|stamped letter)\b",
    re.IGNORECASE,
)
_DECEPTIVE_DOCUMENT_RE = re.compile(
    r"\b(fake|forge|forgery|counterfeit|backdate|change|replace|edit|modify|copy|clone|"
    r"exactly\s+same|look\s+same|make\s+it\s+look|signature|signed|stamp|stamped|official|submit|use\s+it)\b",
    re.IGNORECASE,
)
_SAFE_SAMPLE_RE = re.compile(
    r"\b(blank template|template|sample|fictional|mock|training|example|clearly marked|not for use|"
    r"watermark|draft)\b",
    re.IGNORECASE,
)



@dataclass
class FileCheck:
    name: str
    format: str
    size: int
    ok: bool
    checks: list[str] = field(default_factory=list)
    issues: list[str] = field(default_factory=list)


@dataclass
class Approval:
    ok: bool
    files: list[sandbox.SandboxFile]
    manifest: list[dict]
    reason: str = ""
    kind: str = "files"
    bundle_name: str | None = None


def wants_file(text: str) -> bool:
    """True when the user is asking for a downloadable file (vs. an in-chat answer)."""
    if not text:
        return False
    return wants_app(text) or bool(_FILE_INTENT.search(text) and (_CREATE_VERB.search(text) or "." in text))


def wants_app(text: str) -> bool:
    """True for an explicit request to build a small, one-off client-side app."""
    return bool(text and _APP_INTENT.search(text) and _CREATE_VERB.search(text))


def needs_code(text: str) -> bool:
    return wants_app(text) or bool(text and _NEEDS_CODE.search(text))


def requested_formats(text: str) -> list[str]:
    """Return explicit output formats requested by the user, preserving priority order."""
    found: list[str] = []
    for fmt, pattern in _FORMAT_PATTERNS:
        if pattern.search(text or "") and fmt not in found:
            found.append(fmt)
    return found


def quality_effort(model: str, effort: str | None) -> str:
    """File jobs should not inherit low/auto chat effort; they need deliberate planning."""
    if effort in {"high", "xhigh", "max"}:
        return effort
    if (model or "").startswith("claude_plan/"):
        return QUALITY_CLAUDE_PLAN_EFFORT
    return QUALITY_FILE_EFFORT


def _extension(name: str) -> str:
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


def _format_for_name(name: str) -> str:
    return _EXTENSION_TO_FORMAT.get(_extension(name), "unknown")


def _valid_python(code: str) -> bool:
    try:
        ast.parse(code)
        return True
    except SyntaxError:
        return False


def _extract_code(text: str) -> str:
    match = _CODE_FENCE.search(text or "")
    if match:
        return match.group(1).strip()

    # Only accept unfenced output if the entire response is valid Python.
    # This avoids accidentally running prose that merely contains "import" or "def".
    stripped = (text or "").strip()
    return stripped if stripped and _valid_python(stripped) else ""


def _guard(code: str) -> str:
    return "import os as _os\n_os.makedirs('out', exist_ok=True)\n" + code


def _summary(files: list[sandbox.SandboxFile]) -> str:
    if len(files) == 1:
        return f"Here is your file: **{files[0].name}**."
    names = ", ".join(f"**{f.name}**" for f in files)
    return f"Done — created {len(files)} files: {names}."


def _official_document_error(request: str) -> str | None:
    text = request or ""
    official = bool(_OFFICIAL_DOCUMENT_RE.search(text))
    deceptive = bool(_DECEPTIVE_DOCUMENT_RE.search(text))
    safe_sample = bool(_SAFE_SAMPLE_RE.search(text))

    if official and deceptive and not safe_sample:
        return (
            "I can't generate or modify official, medical, academic, legal, banking, employment, "
            "immigration, or identity documents in a way that could deceive. I can help create a "
            "clearly marked sample/template, a checklist, or an explanation of what a legitimate "
            "document should contain."
        )
    return None


def _requested_file_filter(files: list[sandbox.SandboxFile], request: str) -> tuple[list[sandbox.SandboxFile], str | None]:
    wanted = requested_formats(request)
    if not wanted:
        return files, None

    kept = [f for f in files if _format_for_name(f.name) in wanted]
    if kept:
        return kept, None

    actual = sorted({_format_for_name(f.name) for f in files})
    return [], (
        "Generated files did not match the requested format. "
        f"Requested: {', '.join(wanted)}. Generated: {', '.join(actual) or 'none'}."
    )


def _decode_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_docx_text(data: bytes) -> tuple[str, list[str]]:
    from docx import Document

    document = Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells if cell.text.strip())
    checks = ["opens_as_docx"]
    if parts:
        checks.append("contains_text")
    if document.tables:
        checks.append("contains_tables")
    return "\n".join(parts), checks


def _extract_pptx_text(data: bytes) -> tuple[str, list[str], int]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
    checks = ["opens_as_pptx", f"slide_count:{len(prs.slides)}"]
    if parts:
        checks.append("contains_text")
    return "\n".join(parts), checks, len(prs.slides)


def _extract_xlsx_text(data: bytes) -> tuple[str, list[str], int, int]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=False)
    try:
        parts: list[str] = []
        row_count = 0
        non_empty_cells = 0
        for sheet in workbook.worksheets:
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_index > 500:
                    break
                values = ["" if v is None else str(v) for v in row]
                if any(v.strip() for v in values):
                    row_count += 1
                    non_empty_cells += sum(1 for v in values if v.strip())
                    parts.append("\t".join(values))
        checks = ["opens_as_xlsx", f"sheet_count:{len(workbook.worksheets)}"]
        if row_count:
            checks.append(f"non_empty_rows:{row_count}")
        return "\n".join(parts), checks, row_count, non_empty_cells
    finally:
        workbook.close()


def _extract_pdf_text(data: bytes) -> tuple[str, list[str], int]:
    checks = ["has_pdf_header"] if data.startswith(b"%PDF") else []
    if not data.startswith(b"%PDF"):
        raise ValueError("PDF does not start with a valid %PDF header.")

    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        page_count = len(reader.pages)
        text = "\n".join((page.extract_text() or "").strip() for page in reader.pages)
        checks.append(f"page_count:{page_count}")
        if text.strip():
            checks.append("contains_extractable_text")
        return text, checks, page_count
    except Exception:  # noqa: BLE001
        # Header is valid, but extraction failed. Treat as weakly valid so image-based PDFs can pass.
        checks.append("pdf_text_extraction_unavailable")
        return "", checks, 0


def _extract_csv_text(data: bytes) -> tuple[str, list[str], int, int]:
    text = _decode_text(data)
    rows = list(csv.reader(io.StringIO(text)))
    row_count = len(rows)
    width = max((len(row) for row in rows), default=0)
    checks = ["parses_as_csv", f"row_count:{row_count}", f"max_columns:{width}"]
    return text, checks, row_count, width


def _validate_image(data: bytes, fmt: str) -> tuple[str, list[str]]:
    from PIL import Image

    with Image.open(io.BytesIO(data)) as image:
        width, height = image.size
        image.verify()
    return "", [f"opens_as_image:{fmt}", f"dimensions:{width}x{height}"]


def _validate_html(data: bytes) -> tuple[str, list[str]]:
    from html.parser import HTMLParser

    text = _decode_text(data)
    tags: set[str] = set()
    visible_parts: list[str] = []
    has_inline_style = False
    has_interaction = False

    class PageParser(HTMLParser):
        def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
            nonlocal has_inline_style, has_interaction
            tag = tag.lower()
            tags.add(tag)
            attrs_dict = {name.lower(): value or "" for name, value in attrs}
            if tag == "style" or "style" in attrs_dict:
                has_inline_style = True
            if tag in {"script", "button", "input", "canvas", "select", "textarea"}:
                has_interaction = True
            if any(name.startswith("on") for name in attrs_dict):
                has_interaction = True

        def handle_data(self, data: str) -> None:
            if data.strip():
                visible_parts.append(data)

    parser = PageParser()
    parser.feed(text)
    checks = ["parses_as_html"]
    body_text = " ".join(" ".join(visible_parts).split())

    if not (tags & {"html", "body", "main", "section"}):
        raise ValueError("HTML does not contain a recognizable page structure.")
    if len(body_text) < 60 and not has_interaction:
        raise ValueError("HTML page has too little visible content.")
    if _REMOTE_HTML_REF_RE.search(text) or _HTML_SCRIPT_REMOTE_RE.search(text):
        raise ValueError("HTML includes external or unsafe references; generated pages must be self-contained.")
    if body_text:
        checks.append(f"text_chars:{len(body_text)}")
    if has_inline_style:
        checks.append("has_styles")
    if has_interaction:
        checks.append("has_interaction")
    return text, checks


def _validate_json(data: bytes) -> tuple[str, list[str]]:
    text = _decode_text(data)
    json.loads(text)
    return text, ["parses_as_json"]


def _validate_tex(data: bytes) -> tuple[str, list[str]]:
    if len(data) > 1_000_000:
        raise ValueError("LaTeX source is too large to preview safely.")
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ValueError("LaTeX source must be valid UTF-8 text.") from exc
    if "\x00" in text:
        raise ValueError("LaTeX source contains binary data.")
    if not _LATEX_STRUCTURE_RE.search(text):
        raise ValueError("LaTeX source does not contain recognizable document structure.")
    if _LATEX_UNSAFE_RE.search(text):
        raise ValueError("LaTeX source references unsafe host paths or shell-style file access.")
    return text, ["decodes_as_tex", "has_latex_structure"]


def _parse_svg(text: str):
    """Parse untrusted model SVG via defusedxml (stdlib ElementTree expands entities — a
    billion-laughs payload balloons past any size cap). All parse failures become ValueError, the
    type every caller's guard expects (ParseError is a SyntaxError, not a ValueError)."""
    from defusedxml import ElementTree as SafeElementTree

    try:
        return SafeElementTree.fromstring(text)
    except Exception as exc:  # noqa: BLE001 — ParseError/EntitiesForbidden/DTDForbidden all mean "bad SVG"
        raise ValueError(f"SVG could not be parsed safely: {str(exc)[:200]}") from exc


def _validate_svg(data: bytes) -> tuple[str, list[str]]:
    text = _decode_text(data)
    root = _parse_svg(text)
    tag = str(root.tag).rsplit("}", 1)[-1]
    if tag != "svg":
        raise ValueError("SVG root element is not <svg>.")
    return text, ["parses_as_svg"]


def _validate_wav(data: bytes) -> tuple[str, list[str]]:
    with wave.open(io.BytesIO(data), "rb") as wav:
        channels = wav.getnchannels()
        sample_rate = wav.getframerate()
        frames = wav.getnframes()
        sample_width = wav.getsampwidth()

    if channels < 1 or sample_rate < 8_000 or frames < 1 or sample_width < 1:
        raise ValueError("WAV file has invalid or empty audio metadata.")

    duration = frames / float(sample_rate)
    if duration < 0.2:
        raise ValueError("WAV file is too short to be useful.")

    return "", [
        "opens_as_wav",
        f"channels:{channels}",
        f"sample_rate:{sample_rate}",
        f"duration_seconds:{duration:.2f}",
    ]


def _validate_mp3(data: bytes) -> tuple[str, list[str]]:
    if not (data.startswith(b"ID3") or data[:2] in {b"\xff\xfb", b"\xff\xf3", b"\xff\xf2"}):
        raise ValueError("MP3 file does not have a recognizable MP3 header.")
    return "", ["has_mp3_header"]


def _validate_video(data: bytes, fmt: str) -> tuple[str, list[str]]:
    if len(data) < 1_024:
        raise ValueError(f"{fmt.upper()} file is too small to be a useful video.")
    if fmt == "mp4":
        if b"ftyp" not in data[:64]:
            raise ValueError("MP4 file does not have a recognizable ftyp header.")
        return "", ["has_mp4_header", f"bytes:{len(data)}"]
    if fmt == "webm":
        if not data.startswith(b"\x1a\x45\xdf\xa3"):
            raise ValueError("WebM file does not have a recognizable EBML header.")
        return "", ["has_webm_header", f"bytes:{len(data)}"]
    raise ValueError(f"Unsupported video format: {fmt}")


def _validate_zip(data: bytes) -> tuple[str, list[str]]:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
        for name in names:
            path = PurePosixPath(name)
            if path.is_absolute() or ".." in path.parts:
                raise ValueError("ZIP contains an unsafe path traversal entry.")
        bad = zf.testzip()
        if bad:
            raise ValueError(f"ZIP contains a corrupt member: {bad}")
    return "\n".join(names), ["opens_as_zip", f"member_count:{len(names)}"]


def _safe_app_member(name: str) -> PurePosixPath:
    return file_library.safe_app_member_path(name)


def _normalize_app_bundle(
    files: list[sandbox.SandboxFile],
) -> tuple[list[sandbox.SandboxFile], str]:
    if not files:
        raise ValueError("The sandbox produced no app files.")
    if len(files) > file_library.MAX_APP_BUNDLE_FILES:
        raise ValueError("App bundle contains too many files.")
    if sum(len(file.data) for file in files) > file_library.MAX_APP_BUNDLE_BYTES:
        raise ValueError("App bundle exceeds the total size limit.")
    parsed = [(_safe_app_member(file.name), file) for file in files]
    folded = [path.as_posix().casefold() for path, _file in parsed]
    if len(set(folded)) != len(folded):
        raise ValueError("App bundle contains duplicate file paths.")

    entrypoints = [path for path, _file in parsed if path.name == "index.html"]
    if len(entrypoints) != 1:
        raise ValueError("App bundle must contain exactly one index.html entry point.")
    root = entrypoints[0].parent
    if any(not path.is_relative_to(root) for path, _file in parsed):
        raise ValueError("Every app file must live in the same bundle directory as index.html.")

    # "app.js" and "app.js/logo.png" are each valid alone, but one cannot be both a file and a
    # directory on disk. Caught here, it is a plain validation error the model can repair; left to
    # the host write it surfaces as FileExistsError and takes the whole turn down with it.
    folded_set = set(folded)
    for path, _file in parsed:
        for parent in path.parents:
            if parent.as_posix().casefold() in folded_set:
                raise ValueError(
                    f"App bundle uses '{parent.as_posix()}' as both a file and a folder."
                )

    normalized = [
        sandbox.SandboxFile(path.relative_to(root).as_posix(), file.data)
        for path, file in parsed
    ]
    extensions = {_extension(file.name) for file in normalized}
    if not (extensions & {"js", "mjs"}) or "css" not in extensions:
        raise ValueError("App bundle must include local JavaScript and CSS files.")
    if any(ext not in _APP_ALLOWED_EXTENSIONS for ext in extensions):
        unsupported = sorted(ext or "(no extension)" for ext in extensions if ext not in _APP_ALLOWED_EXTENSIONS)
        raise ValueError(f"App bundle contains unsupported file types: {', '.join(unsupported)}.")

    label = root.name if root != PurePosixPath(".") else "small-app"
    slug = re.sub(r"[^A-Za-z0-9._-]+", "-", label).strip(".-_")[:80] or "small-app"
    return normalized, f"{slug}.zip"


def _app_local_target(reference: str, source: str, members: set[str]) -> str | None:
    value = (reference or "").strip()
    if not value or value.startswith("#"):
        return None
    parsed = urlsplit(value)
    if parsed.scheme:
        raise ValueError(f"{source} contains an external or unsafe reference.")
    if parsed.netloc or value.startswith(("//", "\\")):
        raise ValueError(f"{source} contains an external or unsafe reference.")
    raw_path = unquote(parsed.path)
    if not raw_path:
        return None
    if "\\" in raw_path or raw_path.startswith("/"):
        raise ValueError(f"{source} contains an absolute or unsafe reference.")
    relative = PurePosixPath(raw_path)
    if ".." in relative.parts:
        raise ValueError(f"{source} contains a path traversal reference.")
    target = posixpath.normpath((PurePosixPath(source).parent / relative).as_posix())
    if target not in members:
        directory_entry = f"{target.rstrip('/')}/index.html"
        if directory_entry in members:
            return directory_entry
        raise ValueError(f"{source} references missing local file: {raw_path}.")
    return target


def _app_srcset_targets(value: str) -> list[str]:
    if value.strip().lower().startswith("data:"):
        return [value]
    return [part.strip().split()[0] for part in value.split(",") if part.strip()]


def _validate_app_css_references(text: str, source: str, members: set[str]) -> None:
    references = [match.group(2) for match in _APP_CSS_URL_RE.finditer(text)]
    references.extend(match.group(1) for match in _APP_CSS_IMPORT_RE.finditer(text))
    for reference in references:
        _app_local_target(reference, source, members)


def _validate_app_html(
    text: str,
    source: str,
    members: set[str],
) -> tuple[list[str], set[str], set[str]]:
    from html.parser import HTMLParser

    tags: set[str] = set()
    refs: list[str] = []
    script_refs: list[str] = []
    stylesheet_refs: list[str] = []
    style_chunks: list[str] = []
    issues: list[str] = []
    visible: list[str] = []
    ignored_depth = 0
    style_depth = 0

    class BundleParser(HTMLParser):
        def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
            nonlocal ignored_depth, style_depth
            tag = tag.lower()
            tags.add(tag)
            values: dict[str, str] = {}
            for raw_name, raw_value in attrs:
                name = raw_name.lower()
                if name in values:
                    issues.append(f'Duplicate "{name}" attributes are not allowed.')
                    continue
                values[name] = raw_value or ""
            if tag in {"script", "style"}:
                ignored_depth += 1
            if tag == "style":
                style_depth += 1
            if tag in {"applet", "base", "iframe", "frame", "object", "embed"}:
                issues.append(f"<{tag}> is not allowed in a static app bundle.")
            if tag == "script" and not values.get("src"):
                issues.append("Inline scripts are not allowed; put JavaScript in a local .js file.")
            if tag == "meta" and values.get("http-equiv", "").lower() == "refresh":
                issues.append("HTML refresh navigation is not allowed.")
            if any(name.startswith("on") for name in values):
                issues.append("Inline event handlers are not allowed; bind events from the local JavaScript file.")
            if values.get("action") or values.get("formaction"):
                issues.append("Form navigation is disabled; handle the form locally in JavaScript.")
            if values.get("style"):
                style_chunks.append(values["style"])
            for attr in _APP_HTML_SINGLE_URL_ATTRS:
                if values.get(attr):
                    refs.append(values[attr])
            for attr in _APP_HTML_URL_LIST_ATTRS:
                refs.extend(values.get(attr, "").split())
            for attr in _APP_HTML_SRCSET_ATTRS:
                if values.get(attr):
                    refs.extend(_app_srcset_targets(values[attr]))
            for attr in _APP_HTML_CSS_VALUE_ATTRS:
                if values.get(attr):
                    style_chunks.append(values[attr])
            if tag == "script" and values.get("src"):
                script_refs.append(values["src"])
            rel_values = {part.lower() for part in values.get("rel", "").split()}
            if tag == "link" and "stylesheet" in rel_values and values.get("href"):
                stylesheet_refs.append(values["href"])

        def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
            self.handle_starttag(tag, attrs)
            if tag.lower() in {"script", "style"}:
                self.handle_endtag(tag)

        def handle_endtag(self, tag: str) -> None:
            nonlocal ignored_depth, style_depth
            lowered = tag.lower()
            if lowered == "style" and style_depth:
                style_depth -= 1
            if lowered in {"script", "style"} and ignored_depth:
                ignored_depth -= 1

        def handle_data(self, data: str) -> None:
            if style_depth and data.strip():
                style_chunks.append(data)
            elif not ignored_depth and data.strip():
                visible.append(data.strip())

    parser = BundleParser()
    parser.feed(text)
    if not (tags & {"html", "body", "main", "section"}):
        issues.append("index.html does not contain a recognizable page structure.")
    if source == "index.html" and len(" ".join(visible)) < 20:
        issues.append("index.html has too little visible content.")
    if issues:
        raise ValueError(f"{source}: {' '.join(dict.fromkeys(issues))}")

    for reference in refs:
        _app_local_target(reference, source, members)
    for css in style_chunks:
        _validate_app_css_references(css, source, members)
    loaded_scripts = {
        target
        for reference in script_refs
        if (target := _app_local_target(reference, source, members)) is not None
    }
    loaded_stylesheets = {
        target
        for reference in stylesheet_refs
        if (target := _app_local_target(reference, source, members)) is not None
    }
    return ["parses_as_html", "local_references_only"], loaded_scripts, loaded_stylesheets


def _validate_app_text_file(file: sandbox.SandboxFile, members: set[str]) -> list[str]:
    try:
        text = file.data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ValueError(f"{file.name} must be valid UTF-8 text.") from exc
    ext = _extension(file.name)
    if ext == "html":
        checks, _scripts, _stylesheets = _validate_app_html(text, file.name, members)
        return checks
    if ext == "css":
        _validate_app_css_references(text, file.name, members)
        return ["decodes_as_css", "local_references_only"]
    if ext in {"js", "mjs"}:
        if _APP_NETWORK_SCHEME_RE.search(text) or _APP_JS_NETWORK_RE.search(text):
            raise ValueError(f"{file.name} contains an external reference or network API.")
        if re.search(r"['\"]\s*//", text):
            raise ValueError(f"{file.name} contains a protocol-relative external reference.")
        if re.search(
            r"\b(?:localStorage|sessionStorage|indexedDB|document\s*\.\s*cookie|"
            r"navigator\s*\.\s*serviceWorker|window\s*\.\s*(?:parent|top|opener))\b",
            text,
            re.IGNORECASE,
        ):
            raise ValueError(f"{file.name} requests browser capabilities unavailable to a sandboxed app.")
        for match in _APP_JS_IMPORT_RE.finditer(text):
            _app_local_target(match.group(1), file.name, members)
        return ["decodes_as_javascript", "no_network_apis"]
    if ext == "json":
        json.loads(text)
        return ["parses_as_json"]
    if ext == "svg":
        root = _parse_svg(text)
        for element in root.iter():
            tag = str(element.tag).rsplit("}", 1)[-1].lower()
            if tag in {"script", "foreignobject"}:
                raise ValueError(f"{file.name} contains unsafe SVG content.")
            if tag == "style":
                _validate_app_css_references(element.text or "", file.name, members)
            for attr, value in element.attrib.items():
                local_attr = str(attr).rsplit("}", 1)[-1].lower()
                if local_attr.startswith("on"):
                    raise ValueError(f"{file.name} contains an SVG event handler.")
                if local_attr == "style":
                    _validate_app_css_references(value, file.name, members)
                if local_attr in {"href", "src"}:
                    _app_local_target(value, file.name, members)
        return ["parses_as_svg", "local_references_only"]
    return [f"decodes_as_{ext}"]


def _approve_app_bundle(files: list[sandbox.SandboxFile]) -> Approval:
    try:
        normalized, bundle_name = _normalize_app_bundle(files)
        members = {file.name for file in normalized}
        entry_html = next(file for file in normalized if file.name == "index.html").data.decode("utf-8")
        _entry_checks, loaded_scripts, loaded_stylesheets = _validate_app_html(
            entry_html,
            "index.html",
            members,
        )
        if not any(_extension(target) in {"js", "mjs"} for target in loaded_scripts):
            raise ValueError("index.html must load the bundle's local JavaScript file.")
        if not any(_extension(target) == "css" for target in loaded_stylesheets):
            raise ValueError("index.html must load the bundle's local CSS file.")

        manifest = []
        for file in normalized:
            ext = _extension(file.name)
            checks = ["bundle_member"]
            if ext in {"html", "css", "js", "mjs", "json", "svg", "txt"}:
                checks.extend(_validate_app_text_file(file, members))
            else:
                checks.append(f"local_asset:{ext}")
            if file.name == "index.html":
                checks.append("app_entrypoint")
            manifest.append({
                "name": file.name,
                "format": ext,
                "size": len(file.data),
                "ok": True,
                "checks": checks,
                "issues": [],
            })
        return Approval(
            ok=True,
            files=normalized,
            manifest=manifest,
            kind="app",
            bundle_name=bundle_name,
        )
    except (ValueError, UnicodeDecodeError, json.JSONDecodeError) as exc:
        return Approval(
            ok=False,
            files=[],
            manifest=[],
            reason=f"App bundle validation failed: {str(exc)[:1000]}",
            kind="app",
        )


def _extract_and_validate(file: sandbox.SandboxFile, request: str) -> FileCheck:
    fmt = _format_for_name(file.name)
    check = FileCheck(name=file.name, format=fmt, size=len(file.data), ok=True)

    try:
        if fmt == "docx":
            text, checks = _extract_docx_text(file.data)
            check.checks.extend(checks)
            _check_text_quality(check, text, request, minimum_chars=120)
        elif fmt == "pptx":
            text, checks, slide_count = _extract_pptx_text(file.data)
            check.checks.extend(checks)
            _check_text_quality(check, text, request, minimum_chars=80)
            if slide_count < 4 and not _SMALL_ARTIFACT_OK.search(request or ""):
                check.issues.append("PowerPoint has fewer than 4 slides for a non-trivial deck request.")
        elif fmt == "xlsx":
            text, checks, row_count, non_empty_cells = _extract_xlsx_text(file.data)
            check.checks.extend(checks)
            if row_count < 1 or non_empty_cells < 1:
                check.issues.append("Spreadsheet has no non-empty rows/cells.")
            _check_text_quality(check, text, request, minimum_chars=20, require_text=False)
        elif fmt == "pdf":
            text, checks, page_count = _extract_pdf_text(file.data)
            check.checks.extend(checks)
            if page_count == 0:
                check.checks.append("page_count_unknown")
            _check_text_quality(check, text, request, minimum_chars=120, require_text=False)
        elif fmt == "csv":
            text, checks, row_count, width = _extract_csv_text(file.data)
            check.checks.extend(checks)
            if row_count < 1 or width < 1:
                check.issues.append("CSV has no usable rows/columns.")
            _check_text_quality(check, text, request, minimum_chars=20, require_text=False)
        elif fmt in {"png", "jpg", "gif", "webp"}:
            _text, checks = _validate_image(file.data, fmt)
            check.checks.extend(checks)
        elif fmt == "svg":
            text, checks = _validate_svg(file.data)
            check.checks.extend(checks)
            _check_text_quality(check, text, request, minimum_chars=20, require_text=False)
        elif fmt == "wav":
            _text, checks = _validate_wav(file.data)
            check.checks.extend(checks)
        elif fmt == "mp3":
            _text, checks = _validate_mp3(file.data)
            check.checks.extend(checks)
        elif fmt in {"mp4", "webm"}:
            _text, checks = _validate_video(file.data, fmt)
            check.checks.extend(checks)
        elif fmt == "zip":
            text, checks = _validate_zip(file.data)
            check.checks.extend(checks)
            if not text.strip():
                check.issues.append("ZIP archive is empty.")
        elif fmt == "html":
            text, checks = _validate_html(file.data)
            check.checks.extend(checks)
            _check_text_quality(check, text, request, minimum_chars=80)
        elif fmt == "json":
            text, checks = _validate_json(file.data)
            check.checks.extend(checks)
            _check_text_quality(check, text, request, minimum_chars=20, require_text=False)
        elif fmt == "tex":
            text, checks = _validate_tex(file.data)
            check.checks.extend(checks)
            _check_text_quality(check, text, request, minimum_chars=80)
        elif fmt in {"md", "txt"}:
            text = _decode_text(file.data)
            check.checks.append(f"decodes_as_{fmt}")
            _check_text_quality(check, text, request, minimum_chars=80)
        else:
            check.issues.append(f"Unsupported or unknown output format: {file.name}")
    except Exception as exc:  # noqa: BLE001
        check.issues.append(f"Could not validate {file.name}: {str(exc)[:180]}")

    check.ok = not check.issues
    return check


def _check_text_quality(
    check: FileCheck,
    text: str,
    request: str,
    *,
    minimum_chars: int,
    require_text: bool = True,
) -> None:
    clean = " ".join((text or "").split())
    if require_text and not clean:
        check.issues.append("Generated file contains no readable text.")
        return

    if clean:
        check.checks.append(f"text_chars:{len(clean)}")

    if clean and _PLACEHOLDER_RE.search(clean):
        check.issues.append("Generated file contains placeholder/TODO/sample text.")

    if (
        clean
        and len(clean) < minimum_chars
        and not _SMALL_ARTIFACT_OK.search(request or "")
        and len((request or "").split()) > 6
    ):
        check.issues.append("Generated content is too thin for the user's request.")


def _approve_files(files: list[sandbox.SandboxFile], request: str) -> Approval:
    if wants_app(request):
        return _approve_app_bundle(files)

    filtered, filter_error = _requested_file_filter(files, request)
    if filter_error:
        return Approval(ok=False, files=[], manifest=[], reason=filter_error)

    if not filtered:
        return Approval(ok=False, files=[], manifest=[], reason="The sandbox produced no files to approve.")

    checks = [_extract_and_validate(file, request) for file in filtered]
    manifest = [
        {
            "name": c.name,
            "format": c.format,
            "size": c.size,
            "ok": c.ok,
            "checks": c.checks,
            "issues": c.issues,
        }
        for c in checks
    ]

    failed = [c for c in checks if not c.ok]
    if failed:
        reasons = []
        for c in failed:
            reasons.append(f"{c.name}: " + "; ".join(c.issues))
        return Approval(ok=False, files=[], manifest=manifest, reason="\n".join(reasons)[:3000])

    return Approval(ok=True, files=filtered, manifest=manifest)


async def run(
    model: str,
    request: str,
    system_prompt: str | None,
    effort: str | None,
    untrusted_context: str | None = None,
    trusted_context: str | None = None,
) -> AsyncIterator[dict]:
    """Yield progress events and a final {'result': {...}} with approved files or an error."""
    safety_error = _official_document_error(request)
    if safety_error:
        yield stream_events.result({"ok": False, "error": safety_error})
        return

    app_request = wants_app(request)
    file_effort = quality_effort(model, effort)
    instructions = build_system_prompt(
        app_rules=FILE_SYSTEM_PROMPT,
        skills_block=skills.skills_prompt(request),
        user_preferences=system_prompt,
        trusted_context=trusted_context,
        untrusted_context=untrusted_context,
    )
    if app_request:
        instructions += _APP_BUNDLE_PROMPT
    convo: list[dict] = [{"role": "user", "content": request}]
    last_error = ""
    run_manifests: list[dict] = []
    max_attempts = reasoning.file_retries(effort)  # Quick=1 … Max=4 repair attempts

    for attempt in range(max_attempts):
        yield stream_events.status(
            ("Designing the app…" if app_request else "Designing the document…")
            if attempt == 0
            else f"Fixing the generated {'app' if app_request else 'file'} ({attempt + 1}/{max_attempts})…"
        )
        # Don't advertise a retry counter on the first pass — that reads as a canned "1/N" ladder even
        # when generation succeeds on the first try. Show the attempt number only once a repair actually
        # happens (attempt > 0), so the trace reflects what really occurred instead of a fixed script.
        yield reasoning_event(
            ("Writing the app bundle" if app_request else "Writing the file")
            if attempt == 0
            else ("Repairing the app bundle" if app_request else "Repairing the file"),
            (
                (
                    "Generating a program that builds the requested static app bundle in the sandbox."
                    if app_request
                    else "Generating a program that builds the requested artifact in the sandbox."
                )
                if attempt == 0
                else (
                    f"Retry {attempt + 1}: using the previous runtime or validation failure to fix "
                    f"the generated {'app bundle' if app_request else 'file'}."
                )
            ),
            kind="script",
            status="running",
            phase="generate",
            metadata={"attempt": attempt + 1, "max_attempts": max_attempts},
        )

        parts: list[str] = []
        think = ThinkStream()  # universal: separate reasoning channel OR inline <think>
        try:
            async for delta in ai.stream_chat(model, convo, instructions, file_effort):
                if isinstance(delta, ai.ReasoningDelta):
                    for ev in think.feed_reasoning(str(delta)):
                        yield ev
                    continue
                answer, events = think.feed(str(delta))
                for ev in events:
                    yield ev
                if answer:
                    parts.append(answer)
            tail, events = think.finish()
            for ev in events:
                yield ev
            if tail:
                parts.append(tail)
        except ai.MissingKeyError as exc:
            yield stream_events.result({"ok": False, "error": f"No API key for {exc.provider}. Add it in Settings."})
            return
        except Exception as exc:  # noqa: BLE001 — provider errors already sanitized upstream
            yield stream_events.result({"ok": False, "error": str(exc)})
            return

        reply = "".join(parts)
        code = _extract_code(reply)
        if not code:
            last_error = "The model did not return a valid Python program."
            convo += [
                {"role": "assistant", "content": reply},
                {
                    "role": "user",
                    "content": (
                        "Reply with ONLY one fenced ```python code block that writes the requested "
                        "file(s) into ./out. No prose before or after."
                    ),
                },
            ]
            continue

        yield stream_events.status("Building the app…" if app_request else "Building the file…")
        yield reasoning_event(
            "Running sandbox",
            "Executing the code in the locked-down offline sandbox and collecting output files.",
            kind="tool",
            status="running",
            phase="execute",
            metadata={"attempt": attempt + 1},
        )

        outcome = await asyncio.to_thread(sandbox.run_code, _guard(code))
        if outcome.manifest:
            run_manifests.append(outcome.manifest)
        if outcome.ok and outcome.files:
            yield stream_events.status("Checking the output…")
            yield reasoning_event(
                "Validating output",
                "Opening the generated files, checking requested formats, scanning for placeholders, and enforcing basic quality gates.",
                kind="validation",
                status="running",
                phase="validate",
                metadata={"attempt": attempt + 1},
            )
            # validation parses real Office/PDF/image files — do it off the event loop
            approval = await asyncio.to_thread(_approve_files, outcome.files, request)
            if approval.ok:
                yield stream_events.result({
                    "ok": True,
                    "files": approval.files,
                    "code": code,
                    "summary": (
                        f"Built **{approval.bundle_name}** as a self-contained app bundle."
                        if approval.kind == "app"
                        else _summary(approval.files)
                    ),
                    "kind": approval.kind,
                    "bundle_name": approval.bundle_name,
                    "manifest": approval.manifest,
                    "sandbox": outcome.manifest,
                    "sandbox_runs": run_manifests,
                })
                return

            last_error = approval.reason
            convo += [
                {"role": "assistant", "content": reply},
                {
                    "role": "user",
                    "content": (
                        "The code ran and produced file(s), but the backend rejected the output during "
                        "validation/quality checks. Fix the code and regenerate the file(s).\n\n"
                        f"Validation failure:\n{last_error}\n\n"
                        "Reply with ONLY the corrected Python code block."
                    ),
                },
            ]
            continue

        last_error = (
            outcome.stderr
            or outcome.stdout
            or "The code ran but wrote no files to ./out."
        ).strip()[:3000]
        convo += [
            {"role": "assistant", "content": reply},
            {
                "role": "user",
                "content": (
                    "That failed or produced no files in ./out. Fix it and reply with ONLY the corrected Python.\n\n"
                    f"Error / output:\n{last_error}"
                ),
            },
        ]

    yield stream_events.result({
        "ok": False,
        "error": "I couldn't build a production-quality file after several attempts.",
        "logs": last_error,
        "sandbox_runs": run_manifests,
    })
