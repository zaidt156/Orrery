"""Context building for chat: attachment handling, history serialization, token budgeting, and the
light intent/skill helpers. Pure functions (no DB, no provider calls) extracted from chat.py.
"""

from __future__ import annotations

import base64
import io
import json
import re

from backend.features.prompting import FORMAT_INSTRUCTIONS

CONTEXT_INPUT_SHARE = 0.75
DEFAULT_CONTEXT_WINDOW = 1_000_000


def _title_from(text: str) -> str:
    words = text.strip().split()
    title = " ".join(words[:8])
    return (title[:80] or "New chat").strip()


def _pdf_text(data_url: str) -> str:
    """Extract text from a base64 data-URL PDF locally (no cloud parsing)."""
    try:
        b64 = data_url.split(",", 1)[1] if "," in data_url else data_url
        raw = base64.b64decode(b64)
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(raw))
        pages = [(p.extract_text() or "").strip() for p in reader.pages]
        return "\n\n".join(t for t in pages if t).strip()
    except Exception:  # noqa: BLE001 — a bad/encrypted PDF shouldn't break the turn
        return ""


def _office_text(name: str, data_url: str) -> str:
    """Extract text from Office Open XML files locally, without sending the binary to the model."""
    try:
        b64 = data_url.split(",", 1)[1] if "," in data_url else data_url
        raw = io.BytesIO(base64.b64decode(b64))
        low = (name or "").lower()
        if low.endswith(".docx"):
            from docx import Document

            doc = Document(raw)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append("\t".join(cells))
            return "\n".join(parts).strip()
        if low.endswith((".xlsx", ".xlsm")):
            from openpyxl import load_workbook

            wb = load_workbook(raw, read_only=True, data_only=True)
            parts: list[str] = []
            try:
                for ws in wb.worksheets:
                    parts.append(f"# {ws.title}")
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(cell) for cell in row if cell is not None]
                        if cells:
                            parts.append("\t".join(cells))
            finally:
                wb.close()
            return "\n".join(parts).strip()
        if low.endswith(".pptx"):
            from pptx import Presentation

            prs = Presentation(raw)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts).strip()
    except Exception:  # noqa: BLE001 — unreadable Office docs should not break the chat turn
        return ""
    return ""


def _content_parts(text: str, attachments: list[dict]) -> tuple[str, list[dict]]:
    """Split a turn into (combined_text, image_blocks): text/PDF files are inlined as text,
    images become blocks and also leave a textual marker so later turns know they existed."""
    text_parts = [text] if text else []
    images = []
    for a in attachments:
        content = a.get("content")
        if not content:
            continue
        kind = a.get("kind")
        name = a.get("name", "file")
        if kind == "image":
            images.append({"type": "image_url", "image_url": {"url": content}})
            text_parts.append(f"\n\n[image attached: {name}]")
        elif kind == "pdf":
            extracted = _pdf_text(content)
            body = extracted or "[No extractable text — may be a scanned/image-only PDF.]"
            text_parts.append(f"\n\n--- {name} (PDF) ---\n{body}")
        elif kind == "file" and re.search(r"\.(docx|xlsx|xlsm|pptx)$", name, re.IGNORECASE):
            extracted = _office_text(name, content)
            body = extracted or "[No extractable text found in this Office document.]"
            text_parts.append(f"\n\n--- {name} ---\n{body}")
        else:  # text file — inline its contents
            text_parts.append(f"\n\n--- {name} ---\n{content}")
    return "".join(text_parts).strip(), images


def _build_user_content(text: str, attachments: list[dict]):
    """Multimodal content for this turn: images as blocks, text/PDF inlined; str when no images."""
    combined, images = _content_parts(text, attachments)
    if images:
        blocks = []
        if combined:
            blocks.append({"type": "text", "text": combined})
        return blocks + images
    return combined


def _history_text(text: str, attachments: list[dict]) -> str:
    """Full text persisted as a message's context, so file/PDF text survives later turns."""
    combined, _images = _content_parts(text, attachments)
    return combined


def _db_content(text: str, attachments: list[dict]) -> str:
    """What the bubble shows — the text plus a compact note of attachments (not the file dump)."""
    if not attachments:
        return text
    note = "📎 " + ", ".join(a.get("name", "file") for a in attachments)
    return f"{text}\n\n{note}".strip() if text else note


def _content_token_estimate(content) -> int:
    """Estimate tokens without sending content anywhere or adding a tokenizer dependency."""
    if isinstance(content, str):
        return max(1, (len(content) + 3) // 4)
    if isinstance(content, list):
        tokens = 0
        for block in content:
            if not isinstance(block, dict):
                continue
            if block.get("type") == "text":
                tokens += _content_token_estimate(block.get("text", ""))
            elif block.get("type") == "image_url":
                tokens += 1024
        return max(1, tokens)
    return 1


def _message_groups(messages: list[dict]) -> list[list[dict]]:
    """Group each user turn with the replies that follow it so trimming keeps turns intact."""
    groups: list[list[dict]] = []
    current: list[dict] = []
    for message in messages:
        if message.get("role") == "user" and current:
            groups.append(current)
            current = [message]
        else:
            current.append(message)
    if current:
        groups.append(current)
    return groups


def _limit_messages(
    messages: list[dict],
    context_window: int | None,
    system_prompt: str | None = None,
) -> list[dict]:
    """Keep the newest complete turns within the per-chat approximate token budget."""
    if context_window is None:
        return messages

    formatted_prompt = (
        f"{FORMAT_INSTRUCTIONS}\n\n{system_prompt}" if system_prompt else FORMAT_INSTRUCTIONS
    )
    input_budget = max(
        1,
        int(context_window * CONTEXT_INPUT_SHARE) - _content_token_estimate(formatted_prompt),
    )
    selected: list[list[dict]] = []
    used = 0
    for group in reversed(_message_groups(messages)):
        group_tokens = sum(
            4 + _content_token_estimate(message.get("content", "")) for message in group
        )
        if selected and used + group_tokens > input_budget:
            break
        selected.append(group)
        used += group_tokens

    return [message for group in reversed(selected) for message in group]


def _effective_context_window(context_window: int | None) -> int:
    return context_window or DEFAULT_CONTEXT_WINDOW


def _message_artifacts(raw: str | None) -> list[dict]:
    if not raw:
        return []
    try:
        value = json.loads(raw)
    except (TypeError, json.JSONDecodeError):
        return []
    return value if isinstance(value, list) else []


def _latest_user_text(messages: list[dict]) -> str:
    """The most recent user turn's text — used to select which skills to load."""
    for message in reversed(messages):
        if message.get("role") != "user":
            continue
        content = message.get("content")
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            return " ".join(
                block.get("text", "") for block in content
                if isinstance(block, dict) and block.get("type") == "text"
            )
    return ""


_HIGH_EFFORT_INTENT = re.compile(
    r"\b(code|coding|function|script|program|app|web\s*app|website|web\s*page|frontend|backend|"
    r"component|api|endpoint|algorithm|data\s*structure|implement|refactor|debug|optimi[sz]e|"
    r"svg|diagram|flow\s*chart|chart|infographic|mockup|prototype|architecture|schema|query|regex|"
    r"build\s+(a|an|me)|design\s+(a|an)|create\s+(a|an)|make\s+(a|an)|write\s+(a|an|me))\b",
    re.IGNORECASE,
)


def _wants_high_effort(text: str) -> bool:
    """Code/creative/build requests deserve deliberate, high-effort reasoning."""
    return bool(text and _HIGH_EFFORT_INTENT.search(text))
