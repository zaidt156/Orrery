from __future__ import annotations

import asyncio
import base64
import io
import logging
import uuid

from sqlalchemy import func, select, text, update

from backend.core.config import settings
from backend.core.database import get_sessionmaker
from backend.core.models import EMBED_DIM, Chunk, Collection
from backend.features import team

log = logging.getLogger("orrery.rag")

# Multilingual by default (paraphrase-multilingual-MiniLM-L12-v2 covers ~50 languages) and 384-dim,
# so it drops straight into the existing vector column. The old English-only bge-small model lives
# on inside already-indexed collections, which are queried with their own recorded model (see search).
FALLBACK_EMBED_MODEL = "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
_embedders: dict[str, object] = {}  # model name -> fastembed TextEmbedding (one per model, lazy)
_CURRENT_OWNER = object()


def default_embed_model() -> str:
    """The model NEW collections are built with (configurable; multilingual by default)."""
    return settings.embed_model or FALLBACK_EMBED_MODEL


async def _require_collection_access(
    cid: str, *, owner_id: str | None | object = _CURRENT_OWNER
) -> str | None:
    """Return the effective owner after proving this client can access the collection.

    Team-mode equality deliberately excludes legacy NULL-owner rows. Solo mode has no owner filter,
    preserving the existing local workspace behavior. Callers may pass an already authenticated
    owner to a background job so authorization does not depend on another machine's keychain.
    """
    owner = await team.current_owner_id() if owner_id is _CURRENT_OWNER else owner_id
    filters = [Collection.id == uuid.UUID(cid)]
    if owner is not None:
        filters.append(Collection.owner_id == owner)
    async with get_sessionmaker()() as s:
        found = (await s.execute(select(Collection.id).where(*filters))).scalar_one_or_none()
    if found is None:
        raise PermissionError("Collection not found.")
    return owner  # type: ignore[return-value]


def _model_dim(model_name: str) -> int | None:
    """Declared dimension of a fastembed model, or None if the name isn't recognized."""
    from fastembed import TextEmbedding

    for m in TextEmbedding.list_supported_models():
        if m.get("model") == model_name:
            return m.get("dim")
    return None


def _get_embedder(model_name: str):
    emb = _embedders.get(model_name)
    if emb is None:
        from fastembed import TextEmbedding  # heavy import, deferred to first use

        dim = _model_dim(model_name)
        if dim is not None and dim != EMBED_DIM:
            # The vector column is fixed at EMBED_DIM; a mismatched model would corrupt every row or
            # silently break search. Fail loudly instead (conventions: accuracy over assumption).
            raise ValueError(
                f"Embedding model {model_name} is {dim}-dim, but the store is {EMBED_DIM}-dim. "
                f"Choose a {EMBED_DIM}-dim model or migrate the vector column first."
            )
        emb = _embedders[model_name] = TextEmbedding(model_name=model_name)
    return emb


def _embed_docs(texts: list[str], model_name: str) -> list[list[float]]:
    return [v.tolist() for v in _get_embedder(model_name).embed(texts)]


def _embed_query(q: str, model_name: str) -> list[float]:
    return list(_get_embedder(model_name).query_embed(q))[0].tolist()


async def embed_docs(texts: list[str], model: str | None = None) -> list[list[float]]:
    return await asyncio.to_thread(_embed_docs, texts, model or default_embed_model())


async def embed_query(q: str, model: str | None = None) -> list[float]:
    return await asyncio.to_thread(_embed_query, q, model or default_embed_model())


async def collection_embed_model(cid: str) -> str:
    """The model a collection was built with (falls back to the current default if unknown)."""
    async with get_sessionmaker()() as s:
        m = (await s.execute(
            select(Collection.embed_model).where(Collection.id == uuid.UUID(cid))
        )).scalar_one_or_none()
    return m or default_embed_model()


async def embed_models(cids: list[str]) -> dict[str, str]:
    """Map collection id -> its embedding model, so retrieval embeds the query with the right model
    per collection (a legacy English collection and a new multilingual one both search correctly)."""
    ids = []
    for c in cids:
        try:
            ids.append(uuid.UUID(c))
        except (ValueError, TypeError, AttributeError):
            continue
    if not ids:
        return {}
    async with get_sessionmaker()() as s:
        rows = (await s.execute(
            select(Collection.id, Collection.embed_model).where(Collection.id.in_(ids))
        )).all()
    return {str(cid): model for cid, model in rows}


def _vec(v: list[float]) -> str:
    return "[" + ",".join(f"{x:.7g}" for x in v) + "]"


def _pdf_text(data_url: str) -> str:
    try:
        b64 = data_url.split(",", 1)[1] if "," in data_url else data_url
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(base64.b64decode(b64)))
        return "\n\n".join((p.extract_text() or "").strip() for p in reader.pages).strip()
    except Exception:  # noqa: BLE001
        return ""


def _office_text(name: str, data_url: str) -> str:
    """Extract text from Office files (docx/xlsx/pptx) so any file type can become RAG context."""
    try:
        b64 = data_url.split(",", 1)[1] if "," in data_url else data_url
        raw = io.BytesIO(base64.b64decode(b64))
        low = name.lower()
        if low.endswith(".docx"):
            from docx import Document
            doc = Document(raw)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    parts.append("\t".join(c.text for c in row.cells))
            return "\n".join(parts).strip()
        if low.endswith((".xlsx", ".xls", ".xlsm")):
            from openpyxl import load_workbook
            wb = load_workbook(raw, read_only=True, data_only=True)
            out: list[str] = []
            for ws in wb.worksheets:
                out.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        out.append("\t".join(cells))
            wb.close()
            return "\n".join(out).strip()
        if low.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(raw)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts).strip()
    except Exception:  # noqa: BLE001 — unreadable file → no text, skip rather than break upload
        return ""
    return ""


def _extract(f: dict) -> str:
    name = (f.get("name") or "").lower()
    content = f.get("content") or ""
    if f.get("kind") == "pdf" or name.endswith(".pdf"):
        return _pdf_text(content)
    if name.endswith((".docx", ".xlsx", ".xls", ".xlsm", ".pptx")):
        return _office_text(name, content)
    if f.get("kind") == "text":
        return content  # text file: content is the raw text
    return ""  # unknown binary (image, zip, …): no extractable text, skip rather than embed base64


def chunk_text(body: str, size: int = 900, overlap: int = 150) -> list[str]:
    body = (body or "").strip()
    if not body:
        return []
    out, i = [], 0
    while i < len(body):
        out.append(body[i:i + size])
        i += size - overlap
    return out


def _collection_dict(c: Collection, chunks: int) -> dict:
    return {
        "id": str(c.id), "name": c.name, "embed_model": c.embed_model, "chunks": int(chunks or 0),
        # True when built on an older model than the current default → the UI offers a one-click upgrade.
        "embed_outdated": bool(c.embed_model and c.embed_model != default_embed_model()),
        "kind": getattr(c, "kind", "collection"), "connected": bool(getattr(c, "connected", False)),
        "description": getattr(c, "description", None),
    }


async def list_collections(kind: str = "collection") -> list[dict]:
    """List collections of a given kind ('collection' for the Data tab, 'ontology' for the Ontology tab)."""
    owner = await team.current_owner_id()
    filters = [Collection.kind == kind]
    if owner is not None:
        filters.append(Collection.owner_id == owner)
    async with get_sessionmaker()() as s:
        rows = (await s.execute(
            select(Collection).where(*filters).order_by(Collection.created_at)
        )).scalars().all()
        # One grouped query for all chunk counts instead of a COUNT(*) per collection (N+1).
        counts = dict(
            (await s.execute(
                select(Chunk.collection_id, func.count())
                .where(Chunk.collection_id.in_([c.id for c in rows]))
                .group_by(Chunk.collection_id)
            )).all()
        ) if rows else {}
        return [_collection_dict(c, counts.get(c.id, 0)) for c in rows]


async def create_collection(name: str, kind: str = "collection", description: str | None = None) -> dict:
    owner = await team.current_owner_id()
    async with get_sessionmaker()() as s:
        c = Collection(
            name=(name.strip() or "documents"), embed_model=default_embed_model(), kind=kind,
            description=(description or None), owner_id=owner,
        )
        s.add(c)
        await s.commit()
        await s.refresh(c)
        return _collection_dict(c, 0)


async def set_connected(cid: str, connected: bool) -> bool:
    """Connect/disconnect an ontology so its knowledge is (or isn't) used as context in every chat."""
    owner = await team.current_owner_id()
    filters = [Collection.id == uuid.UUID(cid), Collection.kind == "ontology"]
    if owner is not None:
        filters.append(Collection.owner_id == owner)
    async with get_sessionmaker()() as s:
        c = (await s.execute(select(Collection).where(*filters))).scalar_one_or_none()
        if c is None:
            return False
        c.connected = bool(connected)
        await s.commit()
        return True


async def update_collection(cid: str, name: str | None = None, description: str | None = None) -> bool:
    owner = await team.current_owner_id()
    filters = [Collection.id == uuid.UUID(cid)]
    if owner is not None:
        filters.append(Collection.owner_id == owner)
    async with get_sessionmaker()() as s:
        c = (await s.execute(select(Collection).where(*filters))).scalar_one_or_none()
        if c is None:
            return False
        if name is not None and name.strip():
            c.name = name.strip()[:120]
        if description is not None:
            c.description = description or None
        await s.commit()
        return True


async def connected_collection_ids() -> list[str]:
    """The current owner's connected ontologies; solo mode remains workspace-wide."""
    owner = await team.current_owner_id()
    filters = [Collection.kind == "ontology", Collection.connected.is_(True)]
    if owner is not None:
        filters.append(Collection.owner_id == owner)
    async with get_sessionmaker()() as s:
        rows = (await s.execute(select(Collection.id).where(*filters))).scalars().all()
        return [str(r) for r in rows]


async def delete_collection(cid: str) -> bool:
    owner = await team.current_owner_id()
    filters = [Collection.id == uuid.UUID(cid)]
    if owner is not None:
        filters.append(Collection.owner_id == owner)
    async with get_sessionmaker()() as s:
        c = (await s.execute(select(Collection).where(*filters))).scalar_one_or_none()
        if c is None:
            return False
        await s.delete(c)
        await s.commit()
        return True


async def add_documents(
    cid: str, files: list[dict], *, owner_id: str | None | object = _CURRENT_OWNER
) -> int:
    from sqlalchemy import delete as sa_delete

    await _require_collection_access(cid, owner_id=owner_id)
    model = await collection_embed_model(cid)  # stay on the collection's own model — never mix spaces
    items: list[tuple[str, int, str]] = []
    for f in files:
        for i, ch in enumerate(chunk_text(_extract(f))):
            items.append((f.get("name", "file"), i, ch))
    if not items:
        return 0
    vecs = await embed_docs([c for _, _, c in items], model=model)
    replaced = {src for src, _, _ in items}
    async with get_sessionmaker()() as s:
        # Re-uploading a source REPLACES it — delete-then-insert in ONE transaction, so duplicates
        # can't accumulate and retrieval never sees a half-replaced source (plan Task 3).
        await s.execute(sa_delete(Chunk).where(
            Chunk.collection_id == uuid.UUID(cid), Chunk.source.in_(replaced)
        ))
        for (src, ordn, content), v in zip(items, vecs):
            s.add(Chunk(collection_id=uuid.UUID(cid), source=src, ordinal=ordn, content=content, embedding=v))
        await s.commit()
    return len(items)


# ── bulk ingestion as a durable queue job (large drops must never freeze the app) ──
_INGEST_PROGRESS: dict[str, dict] = {}  # collection id → {state,total_files,done_files,chunks,error}


def ingest_progress(cid: str) -> dict | None:
    return _INGEST_PROGRESS.get(str(uuid.UUID(cid)))


async def enqueue_ingest(cid: str, files: list[dict]) -> dict:
    """Spool payloads to disk and index them as a durable queue job (inline fallback).

    The upload request returns immediately; the UI polls ingest_progress. Payloads are spooled
    as a file because large base64 bodies don't belong in queue-job arguments."""
    import json as _json
    from pathlib import Path

    from backend.core.paths import user_data_dir

    safe_cid = str(uuid.UUID(cid))  # normalize before it touches a filename
    owner = await _require_collection_access(safe_cid)
    spool_dir = user_data_dir() / "tmp" / "ingest"
    spool_dir.mkdir(parents=True, exist_ok=True)
    spool = spool_dir / f"{safe_cid}-{uuid.uuid4().hex}.json"
    spool.write_text(_json.dumps(files), encoding="utf-8")
    _INGEST_PROGRESS[safe_cid] = {"state": "queued", "total_files": len(files),
                                  "done_files": 0, "chunks": 0, "error": None}
    from backend.core.queue import get_queue_app

    try:
        await get_queue_app().configure_task(name="ingest_documents").defer_async(
            cid=safe_cid, spool=str(spool), owner_id=owner)
    except Exception:  # noqa: BLE001 — queue down (tests/dev): index inline so uploads still work
        log.warning("ingest defer failed; running inline for %s", safe_cid)
        await run_ingest(safe_cid, str(spool), owner_id=owner)
    return {"queued": True, "files_queued": len(files)}


async def run_ingest(cid: str, spool: str, owner_id: str | None = None) -> None:
    """Index one spool file with per-file transactions (progress is visible file by file)."""
    import json as _json
    from pathlib import Path

    progress = _INGEST_PROGRESS.setdefault(cid, {"state": "queued", "total_files": 0,
                                                 "done_files": 0, "chunks": 0, "error": None})
    progress.update(state="running", error=None)
    try:
        files = _json.loads(Path(spool).read_text(encoding="utf-8"))
        progress["total_files"] = len(files)
        for payload in files:
            added = await add_documents(cid, [payload], owner_id=owner_id)
            progress["done_files"] += 1
            progress["chunks"] += added
        progress["state"] = "done"
    except Exception as exc:  # noqa: BLE001 — recorded for the UI, never kills the worker
        progress["state"] = "error"
        progress["error"] = str(exc)[:300]
        log.warning("ingest failed for %s: %s", cid, exc)
    finally:
        try:
            Path(spool).unlink(missing_ok=True)
        except OSError:
            pass


async def documents(cid: str) -> list[dict]:
    """Distinct source files in a collection, with their chunk counts."""
    await _require_collection_access(cid)
    async with get_sessionmaker()() as s:
        rows = (await s.execute(
            select(Chunk.source, func.count())
            .where(Chunk.collection_id == uuid.UUID(cid))
            .group_by(Chunk.source)
            .order_by(Chunk.source)
        )).all()
        return [{"source": src, "chunks": int(n or 0)} for src, n in rows]


async def document_text(cid: str, source: str, max_chars: int = 60_000) -> str:
    """The extracted text of one indexed file, re-joined from its chunks (attachment preview)."""
    await _require_collection_access(cid)
    async with get_sessionmaker()() as s:
        rows = (await s.execute(
            select(Chunk.content)
            .where(Chunk.collection_id == uuid.UUID(cid), Chunk.source == source)
            .order_by(Chunk.ordinal)
        )).scalars().all()
    return "\n".join(rows)[:max_chars]


async def delete_source(cid: str, source: str) -> int:
    """Remove all chunks for one source file from a collection."""
    from sqlalchemy import delete
    await _require_collection_access(cid)
    async with get_sessionmaker()() as s:
        result = await s.execute(
            delete(Chunk).where(Chunk.collection_id == uuid.UUID(cid), Chunk.source == source)
        )
        await s.commit()
        return result.rowcount or 0


async def reindex_collection(cid: str) -> int:
    """Re-embed every chunk with the current default model and record the switch — the opt-in way to
    upgrade a collection built on the old English-only model to the multilingual default, no re-upload.

    Chunk text is already stored, so nothing is re-read from source files. The whole collection moves
    to one model at once, so its vector space never ends up mixed."""
    await _require_collection_access(cid)
    model = default_embed_model()
    async with get_sessionmaker()() as s:
        rows = (await s.execute(
            select(Chunk.id, Chunk.content)
            .where(Chunk.collection_id == uuid.UUID(cid))
            .order_by(Chunk.ordinal)
        )).all()
    vecs = await embed_docs([content for _, content in rows], model=model) if rows else []
    async with get_sessionmaker()() as s:
        for (chunk_id, _), v in zip(rows, vecs):
            await s.execute(update(Chunk).where(Chunk.id == chunk_id).values(embedding=v))
        await s.execute(
            update(Collection).where(Collection.id == uuid.UUID(cid)).values(embed_model=model)
        )
        await s.commit()
    return len(rows)


# Above this cosine distance a chunk is judged unrelated to the question and dropped. Without this
# gate the vector arm ALWAYS returns k chunks, so files from earlier turns kept leaking into every
# answer ("why is it reading my Q2 report when I asked about pasta?").
MAX_COSINE_DISTANCE = 0.58


async def search(cid: str, query: str, k: int = 5, query_vector: list[float] | None = None) -> list[dict]:
    """Hybrid retrieval: vector (pgvector) + keyword (Postgres FTS), fused by RRF.

    Relevance-gated: vector hits past MAX_COSINE_DISTANCE are dropped; keyword hits always pass
    (the text literally matched). An empty result means "these files say nothing about this".

    Pass `query_vector` (from embed_query) to reuse one embedding across several collections in a
    single turn instead of re-embedding the same query per collection."""
    await _require_collection_access(cid)
    if query_vector is None:  # embed with THIS collection's model so legacy + new collections both match
        query_vector = await embed_query(query, await collection_embed_model(cid))
    qv = _vec(query_vector)
    async with get_sessionmaker()() as s:
        vec = (await s.execute(
            text("SELECT id::text AS id, source, content, (embedding <=> (:q)::vector) AS dist "
                 "FROM chunks WHERE collection_id = (:cid)::uuid "
                 "ORDER BY embedding <=> (:q)::vector LIMIT :k"),
            {"q": qv, "cid": cid, "k": k},
        )).mappings().all()
        kw = (await s.execute(
            # 'simple' (not 'english') keeps keyword search language-neutral so non-English terms match
            text("SELECT id::text AS id, source, content "
                 "FROM chunks WHERE collection_id = (:cid)::uuid "
                 "AND tsv @@ plainto_tsquery('simple', :q) "
                 "ORDER BY ts_rank(tsv, plainto_tsquery('simple', :q)) DESC LIMIT :k"),
            {"q": query, "cid": cid, "k": k},
        )).mappings().all()

    fused: dict[str, dict] = {}
    for ranked in (vec, kw):
        for rank, row in enumerate(ranked, 1):
            entry = fused.setdefault(row["id"], {"source": row["source"], "content": row["content"], "score": 0.0})
            entry["score"] += 1.0 / (60 + rank)
            if "dist" in row:
                entry["dist"] = float(row["dist"])
            else:
                entry["kw"] = True  # matched by actual words → relevant regardless of distance
    kept = [e for e in fused.values() if e.get("kw") or e.get("dist", 1.0) <= MAX_COSINE_DISTANCE]
    return sorted(kept, key=lambda r: r["score"], reverse=True)[:k]
