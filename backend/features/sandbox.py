"""Sandboxed code execution — the ChatGPT/Claude "code interpreter" path.

Model-written Python runs inside a throwaway Docker container that is locked down hard:
no network, capped memory/CPU/PIDs, read-only root filesystem, non-root user, all Linux
capabilities dropped, no privilege escalation, and a wall-clock timeout. Generated code and input
are read-only mounts; only per-run scratch and output directories are writable. We hand back files
from /work/out plus bounded stdout/stderr. Nothing can reach the user's normal host files.
"""

from __future__ import annotations

import io
import json
import os
import shutil
import stat
import subprocess
import tarfile
import tempfile
import time
import uuid
from collections.abc import Sequence
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from backend.core import proc
from backend.core.config import settings

IMAGE = "orrery-sandbox:latest"
SANDBOX_VERSION = "2"
TIMEOUT_SECONDS = settings.sandbox_timeout_seconds
_MEMORY = "640m"
_CPUS = "1.0"
_PIDS = "256"
_MAX_OUTPUT_FILES = 12
_MAX_TOTAL_OUTPUT_BYTES = 30_000_000
_MAX_FILE_BYTES = 25_000_000
_MAX_INPUT_FILE_BYTES = 50_000_000
_MAX_LOG_CHARS = 6_000
_MAX_CODE_CHARS = 200_000
_LAYOUT = {
    "root": "/work",
    "input": "/work/input",
    "workspace": "/work/workspace",
    "output": "/work/out",
}

_OFFICE_HTML_RENDERER = r"""
import sys
from pathlib import Path

# The renderer module travels in beside the document, on the read-only input mount.
sys.path.insert(0, "/work/input")
import office_render

source = next(
    path for path in sorted(Path("/work/input").iterdir()) if path.name.startswith("document.")
)
suffix = source.suffix.lower()
data = source.read_bytes()

if suffix == ".docx":
    html = office_render._docx_html(data)
elif suffix in (".xlsx", ".xlsm"):
    html = office_render._xlsx_html(data)
elif suffix == ".pptx":
    html = office_render._pptx_html(data)
else:
    raise SystemExit("unsupported Office type: " + suffix)

Path("/work/out/preview.html").write_bytes(html)
"""

_PDF_PAGE_RENDERER = r"""
import io
import json
import tarfile
from pathlib import Path

import pypdfium2 as pdfium

config = json.loads(Path("/work/input/render.json").read_text(encoding="utf-8"))
max_pages = config["max_pages"]
widths = config["widths"]
max_height = config["max_height"]
remaining = config["max_total_bytes"]

out = Path("/work/out")
source = pdfium.PdfDocument("/work/input/document.pdf")
total_pages = len(source)
reason = "page limit" if total_pages > max_pages else None
written = 0
rendered = []

for index in range(min(total_pages, max_pages)):
    page = source[index]
    page_width, page_height = page.get_size()
    if page_width <= 0 or page_height <= 0:
        reason = "page rendering failed"
        break
    encoded = None
    # Step down the resolution ladder until a page fits what is left of the budget.
    for width in widths:
        target = width
        height = max(1, round(target * page_height / page_width))
        if height > max_height:
            target = max(1, round(target * max_height / height))
        image = page.render(scale=target / page_width).to_pil()
        buffer = io.BytesIO()
        image.save(buffer, format="PNG")
        candidate = buffer.getvalue()
        if len(candidate) <= remaining:
            encoded = candidate
            break
    if encoded is None:
        reason = "byte limit"
        break
    rendered.append(("page-%03d.png" % (index + 1), encoded))
    remaining -= len(encoded)
    written += 1

# One archive, not one file per page: the sandbox caps output at a fixed number of files for every
# caller, and a dozen-page document would otherwise be refused outright rather than previewed.
with tarfile.open(out / "pages.tar", "w") as archive:
    for entry_name, payload in rendered:
        info = tarfile.TarInfo(name=entry_name)
        info.size = len(payload)
        archive.addfile(info, io.BytesIO(payload))

(out / "manifest.json").write_text(
    json.dumps({"total_pages": total_pages, "reason": reason, "written": written}),
    encoding="utf-8",
)
"""

_PDF_EXTRACTOR = r"""
from pathlib import Path

import pypdfium2 as pdfium
import pytesseract

source = pdfium.PdfDocument("/work/input/document.pdf")
parts = []
for page in source:
    text_page = page.get_textpage()
    text = text_page.get_text_range().strip()
    if not text:
        image = page.render(scale=200 / 72).to_pil()
        text = pytesseract.image_to_string(image, lang="eng").strip()
    if text:
        parts.append(text)
Path("/work/out/document.txt").write_text("\n\n".join(parts), encoding="utf-8")
""".strip()

_OFFICE_EXTRACTOR = r"""
from pathlib import Path

source = next(
    path for path in sorted(Path("/work/input").iterdir()) if path.name.startswith("document.")
)
suffix = source.suffix.lower()
parts = []
if suffix == ".docx":
    from docx import Document

    document = Document(str(source))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
elif suffix in (".xlsx", ".xlsm"):
    from openpyxl import load_workbook

    workbook = load_workbook(str(source), read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets:
            parts.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) for cell in row if cell is not None]
                if cells:
                    parts.append("\t".join(cells))
    finally:
        workbook.close()
elif suffix == ".pptx":
    from pptx import Presentation

    deck = Presentation(str(source))
    for slide in deck.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text.strip())
Path("/work/out/document.txt").write_text("\n".join(parts).strip(), encoding="utf-8")
""".strip()

OFFICE_SUFFIXES = (".docx", ".xlsx", ".xlsm", ".pptx")


def _docker_bin() -> str:
    candidate = os.environ.get("ORRERY_DOCKER", r"C:\Program Files\Docker\Docker\resources\bin\docker.exe")
    return candidate if Path(candidate).exists() else "docker"


class SandboxError(RuntimeError):
    pass


@dataclass(frozen=True)
class SandboxFile:
    name: str
    data: bytes


@dataclass
class SandboxResult:
    ok: bool
    stdout: str
    stderr: str
    exit_code: int
    timed_out: bool
    files: list[SandboxFile] = field(default_factory=list)
    run_id: str = ""
    manifest: dict[str, Any] = field(default_factory=dict)


# The readiness probe shells out to `docker image inspect`, which is slow (and blocks the
# event loop from its async callers) and is hit several times per chat turn. The image state
# changes rarely, so cache the answer per process: a present image is trusted for longer,
# while an absent one is re-checked sooner so a just-started Docker Desktop is picked up
# quickly. Degrades safely either way — a stale "ready" that fails just falls back to docgen.
_READY_TTL_OK = 30.0
_READY_TTL_MISS = 5.0
_ready_cache: tuple[bool, float] | None = None  # (value, monotonic expiry)


def _probe_image_ready() -> bool:
    try:
        result = proc.run(
            [
                _docker_bin(), "image", "inspect",
                "--format", '{{ index .Config.Labels "org.orrery.sandbox.version" }}', IMAGE,
            ],
            capture_output=True, timeout=25,
        )
        version = result.stdout.decode("utf-8", "replace").strip() if isinstance(result.stdout, bytes) else str(result.stdout).strip()
        return result.returncode == 0 and version == SANDBOX_VERSION
    except Exception:  # noqa: BLE001 — any failure means "not ready"
        return False


def image_ready(*, refresh: bool = False) -> bool:
    """True once the sandbox image has been built and is available locally (briefly cached)."""
    global _ready_cache
    now = time.monotonic()
    if not refresh and _ready_cache is not None and now < _ready_cache[1]:
        return _ready_cache[0]
    value = _probe_image_ready()
    _ready_cache = (value, now + (_READY_TTL_OK if value else _READY_TTL_MISS))
    return value


def _collect_outputs(out_dir: Path) -> list[SandboxFile]:
    """Collect regular files without following model-created links or truncating a bundle."""
    root = out_dir.resolve(strict=True)
    files: list[SandboxFile] = []
    total = 0

    def is_link_like(path: Path) -> bool:
        is_junction = getattr(os.path, "isjunction", lambda _path: False)
        return path.is_symlink() or bool(is_junction(path))

    for current, directories, names in os.walk(root, topdown=True, followlinks=False):
        current_path = Path(current)
        directories.sort()
        names.sort()
        for directory in list(directories):
            if is_link_like(current_path / directory):
                raise SandboxError("Sandbox output contains a symbolic link or junction.")

        for name in names:
            path = current_path / name
            if is_link_like(path):
                raise SandboxError("Sandbox output contains a symbolic link or junction.")
            try:
                metadata = path.lstat()
                resolved = path.resolve(strict=True)
            except OSError as exc:
                raise SandboxError("Sandbox output could not be read safely.") from exc
            if not stat.S_ISREG(metadata.st_mode):
                continue
            if not resolved.is_relative_to(root):
                raise SandboxError("Sandbox output escaped the output directory.")
            size = metadata.st_size
            if size == 0:
                continue
            if size > _MAX_FILE_BYTES:
                raise SandboxError("Sandbox output contains a file over the size limit.")
            if len(files) >= _MAX_OUTPUT_FILES:
                raise SandboxError("Sandbox output contains too many files.")
            if total + size > _MAX_TOTAL_OUTPUT_BYTES:
                raise SandboxError("Sandbox output exceeds the total size limit.")
            files.append(SandboxFile(name=path.relative_to(root).as_posix(), data=path.read_bytes()))
            total += size
    return files


def _build_manifest(
    run_id: str,
    *,
    ok: bool,
    exit_code: int,
    timed_out: bool,
    files: list[SandboxFile],
) -> dict[str, Any]:
    """Public, sanitized run metadata. Never include generated code, logs, prompts, or secrets."""
    return {
        "run_id": run_id,
        "engine": "docker",
        "image": IMAGE,
        "version": SANDBOX_VERSION,
        "layout": dict(_LAYOUT),
        "limits": {
            "timeout_seconds": TIMEOUT_SECONDS,
            "memory": _MEMORY,
            "cpus": _CPUS,
            "pids": _PIDS,
            "max_output_files": _MAX_OUTPUT_FILES,
            "max_total_output_bytes": _MAX_TOTAL_OUTPUT_BYTES,
            "max_file_bytes": _MAX_FILE_BYTES,
            "max_input_file_bytes": _MAX_INPUT_FILE_BYTES,
        },
        "status": {
            "ok": ok,
            "exit_code": exit_code,
            "timed_out": timed_out,
        },
        "outputs": [
            {
                "name": file.name,
                "size": len(file.data),
            }
            for file in files
        ],
    }


def run_code(code: str) -> SandboxResult:
    """Run Python in the sandbox and return its logs plus any files it wrote to out/."""
    return _run_entry(code, "main.py", ["python", "/runner/main.py"])


def run_shell(script: str) -> SandboxResult:
    """Run a shell script in the SAME hardened container (no network, read-only root, dropped caps,
    non-root, cpu/memory/pids caps). Shell adds no exposure Python didn't already have inside the
    box — it just lets the model use the container's CLI tools directly."""
    return _run_entry(script, "script.sh", ["sh", "/runner/script.sh"])


def extract_pdf_text(data: bytes) -> str:
    """Extract each page's text, OCRing only image-only pages, in the locked container."""
    result = _run_entry(
        _PDF_EXTRACTOR,
        "extract_pdf.py",
        ["python", "/runner/extract_pdf.py"],
        input_files={"document.pdf": data},
    )
    if not result.ok:
        raise SandboxError(result.stderr or "The PDF OCR sandbox failed.")
    output = next((item for item in result.files if item.name == "document.txt"), None)
    if output is None:
        return ""
    return output.data.decode("utf-8", "replace").strip()


def render_pdf_pages(
    data: bytes,
    *,
    max_pages: int,
    widths: Sequence[int] = (1000, 850, 700, 560, 440, 320, 240),
    max_total_bytes: int = 5_000_000,
    max_height: int = 1800,
) -> tuple[list[bytes], int, str | None]:
    """Rasterize an untrusted PDF's first pages to PNG inside the locked container.

    Returns (pages in order, the document's true page count, why it stopped early or None).

    The budget is applied INSIDE the container rather than to what comes back, for two reasons: a
    page that will not fit is never written, and the whole document is rendered in one container
    start instead of one per resolution attempt. `widths` is tried in order per page until a render
    fits the remaining byte budget, which is the same ladder the host renderer used.
    """
    config = {
        "max_pages": int(max_pages),
        "widths": [int(width) for width in widths],
        "max_total_bytes": int(max_total_bytes),
        "max_height": int(max_height),
    }
    result = _run_entry(
        _PDF_PAGE_RENDERER,
        "render_pdf.py",
        ["python", "/runner/render_pdf.py"],
        input_files={
            "document.pdf": data,
            "render.json": json.dumps(config).encode("utf-8"),
        },
    )
    if not result.ok:
        raise SandboxError(result.stderr or "The PDF preview sandbox failed.")
    manifest_file = next((item for item in result.files if item.name == "manifest.json"), None)
    if manifest_file is None:
        # Without it a truncated render is indistinguishable from a complete one, and claiming a
        # 40-page document was fully previewed because 1 page came back is the wrong failure.
        raise SandboxError("The PDF preview sandbox returned no manifest.")
    try:
        manifest = json.loads(manifest_file.data.decode("utf-8"))
    except (UnicodeDecodeError, ValueError) as exc:
        raise SandboxError("The PDF preview sandbox returned an unreadable manifest.") from exc

    archive_file = next((item for item in result.files if item.name == "pages.tar"), None)
    pages: list[bytes] = []
    if archive_file is not None:
        try:
            with tarfile.open(fileobj=io.BytesIO(archive_file.data), mode="r:") as archive:
                # Sorted by name: a tar preserves write order, but page order is the entire point of
                # a paginated preview and is too important to inherit from an implementation detail.
                for member in sorted(archive.getmembers(), key=lambda m: m.name):
                    if not member.isfile():
                        continue
                    extracted = archive.extractfile(member)
                    if extracted is not None:
                        pages.append(extracted.read())
        except (tarfile.TarError, OSError) as exc:
            raise SandboxError("The PDF preview sandbox returned an unreadable archive.") from exc
    reason = manifest.get("reason") or None
    return pages, int(manifest.get("total_pages") or 0), reason


def render_office_html(name: str, data: bytes) -> bytes:
    """Render an untrusted Office document to preview HTML inside the locked container.

    The renderer module is shipped in on the read-only input mount and imported there, so the code
    that parses the document is the same code that used to parse it on the host - only the
    privileges differ. python-docx, openpyxl and python-pptx are already in the image.

    The suffix is preserved because openpyxl and python-pptx dispatch on it; the stem is not, so a
    hostile filename never reaches the container.
    """
    suffix = Path(name.lower()).suffix
    if suffix not in OFFICE_SUFFIXES:
        raise SandboxError("Unsupported Office document type.")
    module = Path(__file__).with_name("office_render.py").read_bytes()
    result = _run_entry(
        _OFFICE_HTML_RENDERER,
        "render_office.py",
        ["python", "/runner/render_office.py"],
        input_files={f"document{suffix}": data, "office_render.py": module},
    )
    if not result.ok:
        raise SandboxError(result.stderr or "The Office preview sandbox failed.")
    rendered = next((item for item in result.files if item.name == "preview.html"), None)
    if rendered is None:
        # An empty output directory means the render did not happen, which is not the same thing
        # as a document with nothing in it.
        raise SandboxError("The Office preview sandbox produced no output.")
    return rendered.data


def extract_office_text(name: str, data: bytes) -> str:
    """Extract DOCX/XLSX/PPTX text with the Office libraries inside the locked container."""
    suffix = Path(name.lower()).suffix
    if suffix not in OFFICE_SUFFIXES:
        raise SandboxError("Unsupported Office document type.")
    result = _run_entry(
        _OFFICE_EXTRACTOR,
        "extract_office.py",
        ["python", "/runner/extract_office.py"],
        input_files={f"document{suffix}": data},
    )
    if not result.ok:
        raise SandboxError(result.stderr or "The Office extraction sandbox failed.")
    output = next((item for item in result.files if item.name == "document.txt"), None)
    if output is None:
        return ""
    return output.data.decode("utf-8", "replace").strip()


def _run_entry(
    entry_content: str,
    entry_name: str,
    argv: list[str],
    *,
    input_files: dict[str, bytes] | None = None,
) -> SandboxResult:
    if not entry_content or not entry_content.strip():
        raise SandboxError("There is no code to run.")
    if len(entry_content) > _MAX_CODE_CHARS:
        raise SandboxError("The generated code is too large to run.")
    for filename, payload in (input_files or {}).items():
        if Path(filename).name != filename or not filename:
            raise SandboxError("Sandbox input filename is invalid.")
        if len(payload) > _MAX_INPUT_FILE_BYTES:
            raise SandboxError("Sandbox input contains a file over the size limit.")

    run_id = uuid.uuid4().hex[:12]
    workdir = Path(tempfile.mkdtemp(prefix=f"orrery-sbx-{run_id}-"))
    runner_dir = workdir / "runner"
    input_dir = workdir / "input"
    workspace_dir = workdir / "workspace"
    out_dir = workdir / "out"
    runner_dir.mkdir()
    input_dir.mkdir()
    workspace_dir.mkdir()
    out_dir.mkdir()
    (runner_dir / entry_name).write_text(entry_content, encoding="utf-8", newline="\n")
    for filename, payload in (input_files or {}).items():
        (input_dir / filename).write_bytes(payload)
    name = f"orrery-sbx-{run_id}"

    command = [
        _docker_bin(), "run", "--rm", "--pull", "never", "--name", name,
        "--network", "none",
        "--memory", _MEMORY, "--memory-swap", _MEMORY,
        "--cpus", _CPUS, "--pids-limit", _PIDS,
        "--ulimit", "nofile=256:256",
        "--read-only", "--tmpfs", "/tmp:rw,noexec,nosuid,nodev,size=256m",
        "--cap-drop", "ALL",
        "--security-opt", "no-new-privileges", "--security-opt", "seccomp=builtin",
        "--user", "1000:1000",
        "--mount", f"type=bind,source={runner_dir},target=/runner,readonly",
        "--mount", f"type=bind,source={input_dir},target=/work/input,readonly",
        "--mount", f"type=bind,source={workspace_dir},target=/work/workspace",
        "--mount", f"type=bind,source={out_dir},target=/work/out",
        "-w", "/work/workspace",
        IMAGE, *argv,
    ]

    timed_out = False
    stdout = stderr = ""
    exit_code = -1
    try:
        completed = proc.run(command, capture_output=True, timeout=TIMEOUT_SECONDS)
        exit_code = completed.returncode
        stdout = completed.stdout.decode("utf-8", "replace")[:_MAX_LOG_CHARS]
        stderr = completed.stderr.decode("utf-8", "replace")[:_MAX_LOG_CHARS]
    except subprocess.TimeoutExpired as exc:
        timed_out = True
        stdout = (exc.stdout or b"").decode("utf-8", "replace")[:_MAX_LOG_CHARS] if exc.stdout else ""
        stderr = f"Execution exceeded the {TIMEOUT_SECONDS}s limit and was stopped."
        try:
            proc.run([_docker_bin(), "kill", name], capture_output=True, timeout=25)
        except (OSError, subprocess.TimeoutExpired):
            pass
    except FileNotFoundError as exc:
        shutil.rmtree(workdir, ignore_errors=True)
        raise SandboxError("Docker was not found. Is Docker Desktop installed and running?") from exc

    output_error = ""
    try:
        files = _collect_outputs(out_dir)
    except SandboxError as exc:
        files = []
        output_error = str(exc)
        stderr = (stderr.rstrip() + "\n" + output_error).strip()[:_MAX_LOG_CHARS]
    finally:
        shutil.rmtree(workdir, ignore_errors=True)

    ok = exit_code == 0 and not timed_out and not output_error
    return SandboxResult(
        ok=ok,
        stdout=stdout,
        stderr=stderr,
        exit_code=exit_code,
        timed_out=timed_out,
        files=files,
        run_id=run_id,
        manifest=_build_manifest(run_id, ok=ok, exit_code=exit_code, timed_out=timed_out, files=files),
    )
