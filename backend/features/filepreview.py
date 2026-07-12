"""Render generated binary files into webview-safe inline HTML previews.

Browsers cannot display Office files directly, and packaged Qt WebEngine can show valid PDFs as a
blank frame. Office files therefore use local conversion or a basic HTML fallback, while Windows
PDFs use the bundled QtPdf renderer.
We read the real file (python-pptx / openpyxl / python-docx) and render it as HTML —
slide cards for decks, tables for spreadsheets, and a document view for Word.
"""

from __future__ import annotations

import base64
import hashlib
import html
import io
import platform
import re
import shutil
import subprocess
import tempfile
import threading
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path

from backend.core import proc

_MAX_CACHED_OFFICE_PDF_BYTES = 8_000_000
_MAX_PDF_PREVIEW_PAGES = 24
_MAX_PDF_PREVIEW_PNG_BYTES = 5_000_000
_MAX_PREVIEW_INPUT_BYTES = 7_000_000
_MAX_SOURCE_INPUT_BYTES = 1_000_000
_MAX_PREVIEW_OUTPUT_BYTES = 7_000_000
_MAX_PREVIEW_CHARS = 200_000
_MAX_OFFICE_ARCHIVE_ENTRIES = 2_000
_MAX_OFFICE_UNCOMPRESSED_BYTES = 32_000_000
_MAX_OFFICE_MEMBER_BYTES = 8_000_000
_MAX_OFFICE_NODES = 2_000
_MAX_OFFICE_CELLS = 8_000
_MAX_OFFICE_SHEETS = 24
_MAX_OFFICE_ROWS = 250
_MAX_OFFICE_COLUMNS = 64
_INSTALL_TIMEOUT_SECONDS = 900

_PDF_RENDER_WIDTHS = (1000, 850, 700, 560, 440, 320, 240)

LIBREOFFICE_WINGET_ID = "TheDocumentFoundation.LibreOffice"
LIBREOFFICE_BREW_CASK = "libreoffice"

_install_lock = threading.Lock()
_preview_slot = threading.BoundedSemaphore(1)
_preview_flights_lock = threading.Lock()
_preview_flights: dict[str, "_PreviewFlight"] = {}

_DOCX_BULLET = re.compile(r"^\s*[•▪◦‣]\s*")

_PAGE_CSS = """
*{box-sizing:border-box} body{background:#1b2030;margin:0;padding:24px 20px 56px;
  font-family:Arial,Helvetica,sans-serif;color:#172033}
.meta{color:#9aa6bd;font-size:12px;text-align:center;margin-bottom:16px}
.slide{position:relative;width:100%;max-width:880px;aspect-ratio:16/9;margin:0 auto 20px;background:#fff;
  border-radius:12px;box-shadow:0 10px 34px rgba(0,0,0,.40);padding:44px 52px;overflow:hidden}
.slide h2{font-size:28px;line-height:1.2;margin:0 0 18px;color:#0b1020}
.slide ul{margin:0;padding-left:24px} .slide li{font-size:20px;line-height:1.5;margin:9px 0}
.slide .snum{position:absolute;right:16px;bottom:12px;font-size:12px;color:#9aa6bd}
.sheet{max-width:1040px;margin:0 auto 22px} .sheet h3{color:#cdd8f0;font-size:15px;margin:0 0 8px}
table{border-collapse:collapse;width:100%;background:#fff;border-radius:8px;overflow:hidden}
th,td{border:1px solid #c6cfdd;padding:7px 9px;text-align:left;font-size:13px;vertical-align:top}
th{background:#26314f;color:#fff}
.doc{max-width:820px;margin:0 auto;background:#fff;border-radius:10px;padding:40px 48px}
.doc h1{font-size:26px;color:#0b1020;margin:0 0 14px} .doc h2{font-size:20px;color:#0b1020;margin:18px 0 8px}
.doc p{font-size:15px;line-height:1.6;margin:0 0 12px} .doc li{font-size:15px;line-height:1.6;margin:4px 0}
.pdf-doc{max-width:1040px;margin:0 auto}.pdf-page{margin:0 auto 22px;background:#fff;
  border-radius:6px;box-shadow:0 10px 34px rgba(0,0,0,.40);overflow:hidden}
.pdf-page img{display:block;width:100%;height:auto}.pdf-page-label{padding:7px 10px;background:#f3f5f9;
  color:#69758a;font-size:11px;text-align:center}
.preview-warning{max-width:1040px;margin:0 auto 16px;border:1px solid #b7791f;border-radius:8px;
  background:#fff8e6;padding:10px 12px;color:#6b4300;font-size:13px;line-height:1.45}
.source{max-width:1040px;margin:0 auto;background:#fff;border-radius:10px;padding:28px 32px}
.source pre{white-space:pre-wrap;word-break:break-word;margin:0;font:13px/1.55 Consolas,Menlo,monospace;color:#172033}
"""


def _find_soffice() -> str | None:
    found = proc.find_executable("soffice") or shutil.which("libreoffice")
    if found:
        return found
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        if Path(candidate).is_file():
            return candidate
    return None


def _installer_command() -> list[str] | None:
    """Return the platform's fixed LibreOffice install command, never user input."""
    system = platform.system()
    if system == "Windows":
        winget = shutil.which("winget")
        if not winget:
            return None
        return [
            winget,
            "install",
            "--id",
            LIBREOFFICE_WINGET_ID,
            "--exact",
            "--source",
            "winget",
            "--accept-package-agreements",
            "--accept-source-agreements",
            "--disable-interactivity",
            "--silent",
        ]
    if system == "Darwin":
        brew = proc.find_executable("brew")
        if not brew:
            return None
        return [brew, "install", "--cask", LIBREOFFICE_BREW_CASK]
    return None


def office_preview_status() -> dict:
    """Return a safe live probe result; never expose the executable path."""
    office_available = _find_soffice() is not None
    renderer_available = pdf_renderer_available()
    if office_available and renderer_available:
        return {
            "available": True,
            "engine": "libreoffice",
            "officePreview": "pdf",
            "pdfRendererAvailable": True,
            "canInstall": False,
            "message": "Faithful Office previews are available.",
        }
    if office_available:
        return {
            "available": False,
            "engine": "libreoffice",
            "officePreview": "html",
            "pdfRendererAvailable": False,
            "canInstall": False,
            "message": "The PDF preview renderer is unavailable; Office files use the HTML fallback.",
        }
    return {
        "available": False,
        "engine": "libreoffice",
        "officePreview": "html",
        "pdfRendererAvailable": renderer_available,
        "canInstall": renderer_available and _installer_command() is not None,
        "message": "LibreOffice is not installed; Office files use the HTML fallback.",
    }


def install_office_preview(acknowledged: bool = False) -> dict:
    """Install LibreOffice from a fixed official package and return a fresh probe."""
    if not acknowledged:
        raise ValueError("Confirm the LibreOffice installation before continuing.")
    if not pdf_renderer_available():
        raise ValueError("The packaged PDF preview renderer is unavailable; reinstall or update Orrery.")
    if _find_soffice():
        return office_preview_status()
    command = _installer_command()
    if command is None:
        raise ValueError(
            "One-click LibreOffice installation requires WinGet on Windows or Homebrew on macOS."
        )
    if not _install_lock.acquire(blocking=False):
        raise ValueError("A LibreOffice installation is already running.")
    try:
        # Re-probe after acquiring the lock in case another request completed first.
        if _find_soffice():
            return office_preview_status()
        result = proc.run(
            command,
            cwd=tempfile.gettempdir(),
            text=True,
            encoding="utf-8",
            errors="replace",
            capture_output=True,
            timeout=_INSTALL_TIMEOUT_SECONDS,
            check=False,
        )
        if result.returncode != 0:
            raise ValueError("LibreOffice installation failed. Check the package manager and retry.")
        status = office_preview_status()
        if not status["available"]:
            raise ValueError("LibreOffice installation finished, but the preview engine was not detected.")
        return status
    except subprocess.TimeoutExpired:
        raise ValueError("LibreOffice installation timed out.") from None
    except OSError:
        raise ValueError("LibreOffice installation could not be started.") from None
    finally:
        _install_lock.release()


def is_office_file(name: str) -> bool:
    return Path(name).suffix.lower() in {".pptx", ".docx", ".xlsx", ".xlsm"}


@dataclass
class _PreviewFlight:
    event: threading.Event
    result: object = None
    error: Exception | None = None


def _run_preview_job(key: str, job):
    """Serialize expensive preview jobs and share one result for identical concurrent work."""
    with _preview_flights_lock:
        flight = _preview_flights.get(key)
        leader = flight is None
        if leader:
            flight = _PreviewFlight(threading.Event())
            _preview_flights[key] = flight
    assert flight is not None
    if leader:
        try:
            with _preview_slot:
                flight.result = job()
        except Exception as exc:
            flight.error = exc
        finally:
            with _preview_flights_lock:
                flight.event.set()
                _preview_flights.pop(key, None)
    else:
        flight.event.wait()
    if flight.error is not None:
        raise flight.error
    return flight.result


def _converted_pdf_limit() -> int:
    return min(_MAX_CACHED_OFFICE_PDF_BYTES, _MAX_PREVIEW_INPUT_BYTES)


def _valid_pdf_bytes(data: bytes) -> bool:
    """Validate a converted PDF before it is rendered, returned, or persisted."""
    if not data or len(data) > _converted_pdf_limit():
        return False
    stripped = data.lstrip()
    if not stripped.startswith(b"%PDF-") or b"%%EOF" not in data[-4096:]:
        return False
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data), strict=False)
        return not reader.is_encrypted and len(reader.pages) > 0
    except Exception:  # noqa: BLE001 - malformed or unsupported PDFs are rejected
        return False


def _read_valid_pdf(path: Path) -> bytes | None:
    """Stat before reading so oversized converter/cache output never enters memory."""
    try:
        size = path.stat().st_size
        if size <= 0 or size > _converted_pdf_limit():
            return None
        data = path.read_bytes()
    except OSError:
        return None
    if len(data) != size or not _valid_pdf_bytes(data):
        return None
    return data


def _office_pdf_job(name: str, data: bytes, *, cache_path: Path | None = None) -> bytes | None:
    if cache_path is not None and cache_path.is_file():
        cached = _read_valid_pdf(cache_path)
        if cached is not None:
            return cached
        try:
            cache_path.unlink(missing_ok=True)
        except OSError:
            pass
    soffice = _find_soffice()
    if not soffice:
        return None
    suffix = Path(name).suffix or ".bin"
    with tempfile.TemporaryDirectory(prefix="orrery-preview-") as tmp:
        tmp_path = Path(tmp)
        source = tmp_path / f"input{suffix}"
        profile = tmp_path / "libreoffice-profile"
        profile.mkdir()
        source.write_bytes(data)
        try:
            proc.run(
                [
                    soffice,
                    "--headless",
                    "--nologo",
                    "--nofirststartwizard",
                    f"-env:UserInstallation={profile.resolve().as_uri()}",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(tmp_path),
                    str(source),
                ],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=45,
            )
        except Exception:  # noqa: BLE001
            return None
        pdf = source.with_suffix(".pdf")
        if not pdf.is_file():
            return None
        converted = _read_valid_pdf(pdf)
        if converted is None:
            return None
        if cache_path is not None:
            cache_path.parent.mkdir(parents=True, exist_ok=True)
            temporary = cache_path.with_name(f"{cache_path.name}.{uuid.uuid4().hex}.tmp")
            try:
                temporary.write_bytes(converted)
                temporary.replace(cache_path)
            except OSError:
                pass
            finally:
                try:
                    temporary.unlink(missing_ok=True)
                except OSError:
                    pass
        return converted


def _office_pdf(name: str, data: bytes, *, cache_path: Path | None = None) -> bytes | None:
    source_key = str(cache_path.resolve()) if cache_path is not None else hashlib.sha256(data).hexdigest()
    key = f"office:{Path(name).suffix.lower()}:{source_key}"
    return _run_preview_job(
        key,
        lambda: _office_pdf_job(name, data, cache_path=cache_path),
    )


def _page(title: str, body: str) -> bytes:
    return (
        f'<!doctype html><html lang="en"><head><meta charset="utf-8">'
        f'<meta name="viewport" content="width=device-width, initial-scale=1">'
        f"<title>{html.escape(title)} preview</title><style>{_PAGE_CSS}</style></head>"
        f"<body>{body}</body></html>"
    ).encode("utf-8")


@dataclass
class _PreviewBudget:
    chars: int = 0
    nodes: int = 0
    cells: int = 0
    truncated: bool = False

    def text(self, value: object) -> str:
        text = str(value or "")
        remaining = max(0, _MAX_PREVIEW_CHARS - self.chars)
        if len(text) > remaining:
            text = text[:remaining]
            self.truncated = True
        self.chars += len(text)
        return text

    def node(self) -> bool:
        if self.nodes >= _MAX_OFFICE_NODES:
            self.truncated = True
            return False
        self.nodes += 1
        return True

    def cell(self) -> bool:
        if self.cells >= _MAX_OFFICE_CELLS:
            self.truncated = True
            return False
        self.cells += 1
        return True


def _warning(message: str) -> str:
    return f'<div class="preview-warning" role="status">{html.escape(message)}</div>'


def _notice_html(name: str, message: str, *, truncated: bool = False) -> bytes:
    detail = f"{message} Preview truncated for safety." if truncated else message
    body = f'<main data-renderer="inert-fallback">{_warning(detail)}</main>'
    rendered = _page(name, body)
    if len(rendered) <= _MAX_PREVIEW_OUTPUT_BYTES:
        return rendered
    # The configured output limit should always exceed the fixed shell. This last-resort response
    # remains inert even under an unusually small test/operator override.
    return (
        '<!doctype html><meta charset="utf-8"><p>Preview unavailable. '
        'Preview truncated for safety.</p>'
    ).encode("utf-8")[:_MAX_PREVIEW_OUTPUT_BYTES]


def _finish_html(name: str, body: str, *, truncated: bool = False) -> bytes:
    if truncated:
        body = _warning("Preview truncated for safety. Download the original file to see everything.") + body
    rendered = _page(name, body)
    if len(rendered) <= _MAX_PREVIEW_OUTPUT_BYTES:
        return rendered
    return _notice_html(name, "Preview output exceeded the safe display limit.", truncated=True)


def _inert_fallback(name: str, data: bytes, message: str, *, truncated: bool = False) -> bytes:
    input_truncated = len(data) > _MAX_SOURCE_INPUT_BYTES
    sample = data[:_MAX_SOURCE_INPUT_BYTES]
    text = sample.decode("utf-8", errors="replace")
    char_truncated = len(text) > _MAX_PREVIEW_CHARS
    text = text[:_MAX_PREVIEW_CHARS]
    escaped = html.escape(text)
    body = (
        f'<main data-renderer="inert-fallback">{_warning(message)}'
        f'<div class="source"><pre>{escaped}</pre></div></main>'
    )
    return _finish_html(name, body, truncated=truncated or input_truncated or char_truncated)


def _office_archive_issue(data: bytes, name: str = "") -> str | None:
    if len(data) > _MAX_PREVIEW_INPUT_BYTES:
        return "Office file input exceeded the safe preview limit."
    if Path(name).suffix.lower() == ".xlsm":
        return "Macro-enabled Office files are not sent to the host converter."
    if data.startswith(bytes.fromhex("D0CF11E0A1B11AE1")):
        return "Encrypted Office packages are not sent to the host converter."
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as package:
            entries = package.infolist()
            if len(entries) > _MAX_OFFICE_ARCHIVE_ENTRIES:
                return "Office archive contained too many entries for a safe preview."
            total = 0
            normalized_names: set[str] = set()
            for entry in entries:
                normalized = entry.filename.replace("\\", "/").lstrip("/").lower()
                normalized_names.add(normalized)
                if entry.flag_bits & 0x1:
                    return "Encrypted Office packages are not sent to the host converter."
                if entry.file_size > _MAX_OFFICE_MEMBER_BYTES:
                    return "Office archive contained an oversized compressed member."
                total += entry.file_size
                if total > _MAX_OFFICE_UNCOMPRESSED_BYTES:
                    return "Office archive expanded beyond the safe preview limit."

            if {"encryptioninfo", "encryptedpackage"} & normalized_names:
                return "Encrypted Office packages are not sent to the host converter."
            if any(
                member.endswith("vbaproject.bin")
                or "/vba/" in f"/{member}/"
                or "macrosheet" in member
                for member in normalized_names
            ):
                return "Office packages containing macros are not sent to the host converter."

            content_types = next(
                (entry for entry in entries if entry.filename.replace("\\", "/").lstrip("/").lower() == "[content_types].xml"),
                None,
            )
            if content_types is not None:
                content = package.read(content_types).lower()
                if b"macroenabled" in content or b"vba" in content:
                    return "Macro-enabled Office packages are not sent to the host converter."

            from defusedxml import ElementTree as SafeElementTree

            for entry in entries:
                if not entry.filename.lower().endswith(".rels"):
                    continue
                try:
                    relationships = SafeElementTree.fromstring(package.read(entry))
                except Exception:  # noqa: BLE001 - malformed relationship XML fails closed
                    return "Office relationship metadata could not be safely parsed."
                for relationship in relationships.iter():
                    target_mode = next(
                        (
                            value
                            for key, value in relationship.attrib.items()
                            if key.rsplit("}", 1)[-1].lower() == "targetmode"
                        ),
                        "",
                    )
                    if target_mode.lower() == "external":
                        return "Office packages with an external relationship are not sent to the host converter."
    except (zipfile.BadZipFile, OSError, RuntimeError, ValueError):
        return "Preview unavailable because the Office file could not be parsed."
    return None


def _bounded_pdf_page_png(render_at_width, *, remaining: int) -> bytes | None:
    """Render a page at progressively lower resolutions until it fits the remaining budget."""
    if remaining <= 0:
        return None
    for width in _PDF_RENDER_WIDTHS:
        png = render_at_width(width)
        if png and len(png) <= remaining:
            return png
    return None


@dataclass(frozen=True)
class _PdfRender:
    pages: tuple[bytes, ...]
    total_pages: int
    complete: bool
    reason: str | None = None


def pdf_renderer_available() -> bool:
    """Return whether the local QtPdf renderer can be imported in this runtime."""
    try:
        from PySide6.QtCore import QBuffer, QByteArray, QIODevice, QSize  # noqa: F401
        from PySide6.QtGui import QImage  # noqa: F401
        from PySide6.QtPdf import QPdfDocument  # noqa: F401
    except (ImportError, OSError):
        return False
    return True


def _render_pdf_pngs_uncached(data: bytes) -> _PdfRender | None:
    """Rasterize PDF pages with QtPdf when it is available.

    The packaged Qt WebEngine PDF viewer can show a blank document even when the PDF is valid.
    QtPdf is already part of Orrery's Windows desktop runtime, so rendering pages here avoids that
    viewer dependency without adding a network service. When QtPdf is unavailable or cannot fit a
    page inside the preview budget, the caller receives bounded explanatory HTML instead of a raw
    PDF that the packaged webview cannot display.
    """
    try:
        from PySide6.QtCore import QBuffer, QByteArray, QIODevice, QSize
        from PySide6.QtPdf import QPdfDocument
    except (ImportError, OSError):
        return None

    source_data = QByteArray(data)
    source = QBuffer(source_data)
    if not source.open(QIODevice.OpenModeFlag.ReadOnly):
        return None
    document = QPdfDocument()
    try:
        load_error = document.load(source)
        if load_error not in (None, QPdfDocument.Error.None_) or document.pageCount() < 1:
            return None
        total_pages = document.pageCount()
        pages: list[bytes] = []
        total = 0
        reason = "page limit" if total_pages > _MAX_PDF_PREVIEW_PAGES else None
        for page_number in range(min(total_pages, _MAX_PDF_PREVIEW_PAGES)):
            points = document.pagePointSize(page_number)
            if points.width() <= 0 or points.height() <= 0:
                reason = "page rendering failed"
                break

            def render_at_width(width: int) -> bytes | None:
                height = max(1, round(width * points.height() / points.width()))
                if height > 1800:
                    width = max(1, round(width * 1800 / height))
                    height = 1800
                image = document.render(page_number, QSize(width, height))
                if image.isNull():
                    return None
                output = QBuffer()
                if not output.open(QIODevice.OpenModeFlag.WriteOnly):
                    return None
                try:
                    if not image.save(output, "PNG"):
                        return None
                    return bytes(output.data())
                finally:
                    output.close()

            png = _bounded_pdf_page_png(
                render_at_width,
                remaining=_MAX_PDF_PREVIEW_PNG_BYTES - total,
            )
            if png is None:
                reason = "byte limit"
                break
            pages.append(png)
            total += len(png)
        if not pages:
            if reason in {"byte limit", "page limit"}:
                return _PdfRender((), total_pages, False, reason)
            return None
        complete = len(pages) == total_pages and reason is None
        return _PdfRender(tuple(pages), total_pages, complete, reason)
    except Exception:  # noqa: BLE001 - malformed PDFs fall back to the native response
        return None
    finally:
        document.close()
        source.close()


def _render_pdf_pngs(data: bytes) -> _PdfRender | None:
    key = f"pdf:{hashlib.sha256(data).hexdigest()}"
    return _run_preview_job(key, lambda: _render_pdf_pngs_uncached(data))


def _rendered_pdf_html(name: str, data: bytes) -> bytes | None:
    rendered = _render_pdf_pngs(data)
    if rendered is None:
        return None
    cards = []
    for page_number, png in enumerate(rendered.pages, start=1):
        encoded = base64.b64encode(png).decode("ascii")
        cards.append(
            f'<section class="pdf-page" aria-label="Page {page_number}">'
            f'<img src="data:image/png;base64,{encoded}" alt="PDF page {page_number}">'
            f'<div class="pdf-page-label">Page {page_number}</div></section>'
        )
    count = len(cards)
    complete = "true" if rendered.complete else "false"
    truncation_reason = ""
    if not rendered.complete and rendered.reason:
        reason = re.sub(r"[^a-z0-9]+", "-", rendered.reason.lower()).strip("-")
        truncation_reason = f' data-preview-truncation-reason="{html.escape(reason)}"'
    body = (
        f'<main class="pdf-doc" data-renderer="qt-pdf" data-preview-complete="{complete}"{truncation_reason}>'
        f'<div class="meta">PDF preview · {count} of {rendered.total_pages} rendered page(s) '
        f'· {html.escape(name)}</div>{"".join(cards)}</main>'
    )
    return _finish_html(name, body, truncated=not rendered.complete)


def _pdf_html(name: str, data: bytes) -> bytes:
    rendered = _rendered_pdf_html(name, data)
    if rendered is None:
        return _notice_html(
            name,
            "PDF preview is unavailable in the embedded viewer. Download the original file to open it.",
        )
    return rendered


def _pptx_html(data: bytes) -> bytes:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    budget = _PreviewBudget()
    cards = []
    for index, slide in enumerate(prs.slides, start=1):
        if not budget.node():
            break
        title = ""
        title_shape = None
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                title = budget.text(slide.shapes.title.text.strip())
                title_shape = slide.shapes.title
        except (AttributeError, ValueError):
            pass
        bullets: list[str] = []
        for shape in slide.shapes:
            if not budget.node():
                break
            if shape is title_shape or not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if not budget.node():
                    break
                text = budget.text(paragraph.text.strip())
                if text:
                    if len(bullets) >= 14:
                        budget.truncated = True
                        break
                    bullets.append(text)
        items = "".join(f"<li>{html.escape(b)}</li>" for b in bullets)
        body = f"<ul>{items}</ul>" if items else ""
        cards.append(
            f'<div class="slide"><div class="snum">{index}</div>'
            f'<h2>{html.escape(title or f"Slide {index}")}</h2>{body}</div>'
        )
    body = f'<div class="meta">PowerPoint preview · {len(cards)} slides</div>' + "".join(cards)
    return _finish_html("PowerPoint", body, truncated=budget.truncated)


def _xlsx_html(data: bytes) -> bytes:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True)
    budget = _PreviewBudget()
    parts = []
    try:
        sheets = workbook.worksheets
        if len(sheets) > _MAX_OFFICE_SHEETS:
            budget.truncated = True
        for sheet in sheets[:_MAX_OFFICE_SHEETS]:
            if not budget.node():
                break
            if (sheet.max_row or 0) > _MAX_OFFICE_ROWS or (sheet.max_column or 0) > _MAX_OFFICE_COLUMNS:
                budget.truncated = True
            rows = []
            for row_index, row in enumerate(
                sheet.iter_rows(
                    max_row=_MAX_OFFICE_ROWS,
                    max_col=_MAX_OFFICE_COLUMNS,
                    values_only=True,
                )
            ):
                if not budget.node():
                    break
                tag = "th" if row_index == 0 else "td"
                cells = []
                for value in row:
                    if not budget.cell():
                        break
                    value_text = budget.text("" if value is None else value)
                    cells.append(f"<{tag}>{html.escape(value_text)}</{tag}>")
                if cells:
                    rows.append(f"<tr>{''.join(cells)}</tr>")
                if budget.cells >= _MAX_OFFICE_CELLS:
                    break
            title = html.escape(budget.text(sheet.title))
            parts.append(f'<div class="sheet"><h3>{title}</h3><table>{"".join(rows)}</table></div>')
            if budget.cells >= _MAX_OFFICE_CELLS:
                break
    finally:
        workbook.close()
    body = f'<div class="meta">Excel preview · {len(parts)} sheet(s)</div>' + "".join(parts)
    return _finish_html("Excel", body, truncated=budget.truncated)


def _docx_html(data: bytes) -> bytes:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(io.BytesIO(data))
    budget = _PreviewBudget()
    parts = ['<div class="doc">']
    in_list = False
    # Iterate body XML so tables stay where Word placed them. document.tables would append every
    # table after all paragraphs, which made CV skill tables appear at the end of the preview.
    for child in document.element.body.iterchildren():
        if not budget.node():
            break
        if child.tag.endswith("}p"):
            paragraph = Paragraph(child, document)
            text = budget.text(paragraph.text.strip())
            if not text:
                continue
            style = (paragraph.style.name or "").lower() if paragraph.style else ""
            runs = [run for run in paragraph.runs if run.text.strip()]
            sizes = [run.font.size.pt for run in runs if run.font.size is not None]
            max_size = max(sizes, default=0)
            all_bold = bool(runs) and all(run.bold is True for run in runs)
            direct_section = (
                all_bold
                and len(text) <= 80
                and text.upper() == text
                and any(character.isalpha() for character in text)
            )
            is_list = "list" in style or "bullet" in style or bool(_DOCX_BULLET.match(text))
            if is_list and not in_list:
                parts.append("<ul>")
                in_list = True
            elif not is_list and in_list:
                parts.append("</ul>")
                in_list = False
            if "title" in style or "heading 1" in style or max_size >= 18:
                parts.append(f"<h1>{html.escape(text)}</h1>")
            elif "heading" in style or direct_section:
                parts.append(f"<h2>{html.escape(text)}</h2>")
            elif is_list:
                parts.append(f"<li>{html.escape(_DOCX_BULLET.sub('', text))}</li>")
            else:
                parts.append(f"<p>{html.escape(text)}</p>")
        elif child.tag.endswith("}tbl"):
            if in_list:
                parts.append("</ul>")
                in_list = False
            table = Table(child, document)
            rows = []
            for row_index, row in enumerate(table.rows):
                if not budget.node():
                    break
                tag = "th" if row_index == 0 else "td"
                cells = []
                for cell in row.cells:
                    if not budget.cell():
                        break
                    cells.append(f"<{tag}>{html.escape(budget.text(cell.text))}</{tag}>")
                if cells:
                    rows.append(f"<tr>{''.join(cells)}</tr>")
                if budget.cells >= _MAX_OFFICE_CELLS:
                    break
            parts.append(f"<table>{''.join(rows)}</table>")
    if in_list:
        parts.append("</ul>")
    parts.append("</div>")
    return _finish_html("Document", "".join(parts), truncated=budget.truncated)


def _source_html(name: str, data: bytes) -> bytes:
    input_truncated = len(data) > _MAX_SOURCE_INPUT_BYTES
    text = data[:_MAX_SOURCE_INPUT_BYTES].decode("utf-8", errors="replace")
    char_truncated = len(text) > _MAX_PREVIEW_CHARS
    text = text[:_MAX_PREVIEW_CHARS]
    body = (
        f'<main data-renderer="inert-source"><div class="meta">Source preview · {html.escape(name)}</div>'
        f'<div class="source"><pre>{html.escape(text)}</pre></div></main>'
    )
    return _finish_html(name, body, truncated=input_truncated or char_truncated)


def to_preview(
    name: str,
    mime: str,
    data: bytes,
    *,
    cache_path: Path | None = None,
) -> tuple[bytes, str]:
    """Return webview-safe preview content and its media type."""
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    if len(data) > _MAX_PREVIEW_INPUT_BYTES:
        return (
            _inert_fallback(
                name,
                data,
                "Preview unavailable because the file exceeded the safe input limit.",
                truncated=True,
            ),
            "text/html; charset=utf-8",
        )
    if ext == "pdf":
        return _pdf_html(name, data), "text/html; charset=utf-8"
    if ext in ("pptx", "docx", "xlsx", "xlsm"):
        archive_issue = _office_archive_issue(data, name)
        if archive_issue:
            return (
                _inert_fallback(name, data, archive_issue, truncated=True),
                "text/html; charset=utf-8",
            )
        pdf = _office_pdf(name, data, cache_path=cache_path)
        if pdf:
            rendered_pdf = _rendered_pdf_html(name, pdf)
            if rendered_pdf is not None:
                return rendered_pdf, "text/html; charset=utf-8"
    try:
        if ext == "pptx":
            return _pptx_html(data), "text/html; charset=utf-8"
        if ext in ("xlsx", "xlsm"):
            return _xlsx_html(data), "text/html; charset=utf-8"
        if ext == "docx":
            return _docx_html(data), "text/html; charset=utf-8"
        if ext in ("tex", "bib", "sty", "cls"):
            return _source_html(name, data), "text/html; charset=utf-8"
    except Exception:  # noqa: BLE001 — on any parse failure, fall back to the raw bytes
        return (
            _inert_fallback(name, data, "Preview unavailable because the file could not be parsed."),
            "text/html; charset=utf-8",
        )

    source_extensions = {
        "bib", "cls", "css", "csv", "html", "htm", "js", "json", "md", "markdown",
        "py", "svg", "sty", "tex", "txt", "xhtml", "xml", "yaml", "yml",
    }
    normalized_mime = mime.lower().split(";", 1)[0].strip()
    if ext in source_extensions or normalized_mime in {
        "application/javascript",
        "application/xhtml+xml",
        "image/svg+xml",
        "text/html",
    }:
        return _source_html(name, data), "text/html; charset=utf-8"
    if normalized_mime in {"image/png", "image/jpeg", "image/gif", "image/webp", "image/bmp"}:
        return data, normalized_mime
    if normalized_mime.startswith("audio/") or normalized_mime.startswith("video/"):
        return data, normalized_mime
    return (
        _inert_fallback(name, data, "Preview unavailable for this file type."),
        "text/html; charset=utf-8",
    )
