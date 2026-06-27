"""Structured document generation: the model designs a file as JSON (slides, sheets,
sections); Orrery builds a *real* PPTX/XLSX/DOCX/PDF/… from it deterministically.

This is the "model designs, Orrery builds" path — no model-written code is executed.
A reply carries the design inside a ```orrery-doc fenced JSON block; if present we build
from the structure, otherwise callers fall back to rendering the reply's Markdown.

Production stance:
- Treat the model's JSON as untrusted input.
- Normalize and validate the spec before rendering.
- Reject placeholder/TODO content.
- Reopen generated Office/PDF/CSV files after rendering to catch broken output.
- Never silently invent placeholder content inside renderers.
"""

from __future__ import annotations

import csv
import html
import io
import json
import re
from copy import deepcopy
from dataclasses import dataclass
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from backend.features import exports
from backend.features.exports import (
    ExportBlock,
    ExportResult,
    ExportTooLarge,
    InlineSpan,
)

_SPEC_FENCE = re.compile(r"```orrery-doc\s*\n([\s\S]*?)```", re.IGNORECASE)

MAX_SPEC_CHARS = 600_000

# Separate limits are safer than one broad MAX_ITEMS.
MAX_SLIDES = 80
MAX_SECTIONS = 220
MAX_SHEETS = 25
MAX_ROWS_PER_SHEET = 5_000
MAX_COLUMNS = 100
MAX_TABLE_ROWS_IN_DOCUMENT = 2_000
MAX_BULLETS_PER_SLIDE = 12
MAX_BULLETS_PER_SECTION = 200
MAX_PARAGRAPHS_PER_SECTION = 200
MAX_TEXT_CHARS = 5_000
MAX_CELL_CHARS = 32_767

# PPTX readability limits. If the spec needs more than this, it should be split
# into multiple slides by the generation layer or rejected before rendering.
MAX_PPTX_TABLE_ROWS = 10
MAX_PPTX_TABLE_COLUMNS = 6
MAX_PPTX_METRICS = 6

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

_ALLOWED_SLIDE_LAYOUTS = {
    "bullets",
    "section",
    "two_column",
    "table",
    "quote",
    "metrics",
    "summary",
}

_PLACEHOLDER_RE = re.compile(
    r"\b("
    r"todo|lorem ipsum|placeholder|insert text here|sample text|tbd|"
    r"your name|company name|example company|dummy text|"
    r"\[title\]|\[date\]|\[name\]|\[company\]|\[insert"
    r")\b",
    re.IGNORECASE,
)

_INVALID_SHEET_CHARS = re.compile(r"[\[\]:*?/\\]")


class DocSpecError(ValueError):
    """Raised when an orrery-doc spec is unsafe, malformed, or too weak to render."""


@dataclass(frozen=True)
class _RenderedValidation:
    format: str
    checks: list[str]


# ---------------------------------------------------------------------------
# Public parse API
# ---------------------------------------------------------------------------

def parse_doc_spec(content: str | None) -> dict | None:
    """Return the structured spec embedded in a reply, or None to fall back to Markdown.

    This parser is intentionally forgiving at the boundary: if the fenced JSON is absent,
    invalid, too large, or fails normalization, return None so the caller can fall back.
    Rendering functions perform stricter format-specific validation and may raise.
    """
    if not content or "orrery-doc" not in content:
        return None

    match = _SPEC_FENCE.search(content)
    if not match:
        return None

    raw = match.group(1).strip()
    if not raw or len(raw) > MAX_SPEC_CHARS:
        return None

    try:
        spec = json.loads(raw)
    except json.JSONDecodeError:
        return None

    if not isinstance(spec, dict):
        return None

    try:
        normalized = normalize_spec(spec)
    except DocSpecError:
        return None

    return normalized if _has_any_structured_content(normalized) else None


def has_doc_spec(content: str | None) -> bool:
    return parse_doc_spec(content) is not None


# ---------------------------------------------------------------------------
# Generic cleaning / normalization helpers
# ---------------------------------------------------------------------------

def _s(value: Any) -> str:
    return exports._clean_text("" if value is None else str(value)).strip()


def _truncate_text(value: Any, limit: int = MAX_TEXT_CHARS) -> str:
    text = _s(value)
    return text[:limit]


def _list(value: Any, limit: int | None = None) -> list:
    if not isinstance(value, list):
        return []
    return value[:limit] if limit is not None else value


def _spans(text: Any) -> list[InlineSpan]:
    value = _truncate_text(text)
    return [InlineSpan(value)] if value else []


def _has_text(value: Any) -> bool:
    return bool(_s(value))


def _has_any_structured_content(spec: dict) -> bool:
    return bool(spec.get("slides") or spec.get("sheets") or spec.get("sections"))


def _safe_level(value: Any) -> int:
    try:
        level = int(value)
    except (TypeError, ValueError):
        return 1
    return min(6, max(1, level))


def _safe_sheet_name(name: Any, fallback: str, used: set[str]) -> str:
    base = _INVALID_SHEET_CHARS.sub("_", _s(name) or fallback).strip("'")[:31] or fallback
    candidate = base
    counter = 2
    while candidate in used:
        suffix = f"_{counter}"
        candidate = (base[: 31 - len(suffix)] + suffix)[:31]
        counter += 1
    used.add(candidate)
    return candidate


def _normalize_string_list(value: Any, *, limit: int, item_limit: int = MAX_TEXT_CHARS) -> list[str]:
    out: list[str] = []
    for item in _list(value, limit):
        text = _truncate_text(item, item_limit)
        if text:
            out.append(text)
    return out


def _normalize_table(table: Any, *, row_limit: int = MAX_TABLE_ROWS_IN_DOCUMENT) -> dict | None:
    if not isinstance(table, dict):
        return None

    columns = _normalize_string_list(table.get("columns"), limit=MAX_COLUMNS, item_limit=MAX_CELL_CHARS)
    rows: list[list[str]] = []

    for row in _list(table.get("rows"), row_limit):
        if isinstance(row, list):
            values = [_truncate_text(cell, MAX_CELL_CHARS) for cell in row[:MAX_COLUMNS]]
        else:
            values = [_truncate_text(row, MAX_CELL_CHARS)]

        # Keep row shape, but remove rows that are entirely empty.
        if any(value for value in values):
            rows.append(values)

    if not columns and not rows:
        return None

    return {"columns": columns, "rows": rows}


def _infer_slide_layout(slide: dict) -> str:
    layout = _s(slide.get("layout")).lower()
    if layout in _ALLOWED_SLIDE_LAYOUTS:
        return layout

    if isinstance(slide.get("table"), dict):
        return "table"
    if slide.get("left") or slide.get("right"):
        return "two_column"
    if slide.get("metrics"):
        return "metrics"
    if slide.get("quote"):
        return "quote"
    if not slide.get("bullets") and not slide.get("paragraphs"):
        return "section"
    return "bullets"


def _normalize_metric_list(value: Any) -> list[dict]:
    metrics: list[dict] = []
    for item in _list(value, MAX_PPTX_METRICS):
        if not isinstance(item, dict):
            continue
        label = _truncate_text(item.get("label"), 120)
        value_text = _truncate_text(item.get("value"), 120)
        note = _truncate_text(item.get("note"), 220)
        if label or value_text:
            metrics.append({"label": label, "value": value_text, "note": note})
    return metrics


def _normalize_slide(raw: Any, index: int, fallback_title: str) -> dict | None:
    if not isinstance(raw, dict):
        return None

    title = _truncate_text(raw.get("title") or raw.get("heading") or f"Slide {index}", 160)
    bullets = _normalize_string_list(raw.get("bullets"), limit=MAX_BULLETS_PER_SLIDE, item_limit=500)
    paragraphs = _normalize_string_list(raw.get("paragraphs"), limit=8, item_limit=700)
    left = _normalize_string_list(raw.get("left"), limit=8, item_limit=500)
    right = _normalize_string_list(raw.get("right"), limit=8, item_limit=500)
    quote = _truncate_text(raw.get("quote"), 900)
    attribution = _truncate_text(raw.get("attribution"), 180)
    notes = _truncate_text(raw.get("notes"), 4_000)
    table = _normalize_table(raw.get("table"), row_limit=MAX_PPTX_TABLE_ROWS)
    metrics = _normalize_metric_list(raw.get("metrics"))

    slide = {
        "title": title or fallback_title or f"Slide {index}",
        "layout": _infer_slide_layout(raw),
        "bullets": bullets,
        "paragraphs": paragraphs,
        "left": left,
        "right": right,
        "quote": quote,
        "attribution": attribution,
        "notes": notes,
        "table": table,
        "metrics": metrics,
    }

    if _slide_has_content(slide):
        return slide

    return None


def _slide_has_content(slide: dict) -> bool:
    return bool(
        _s(slide.get("title"))
        or slide.get("bullets")
        or slide.get("paragraphs")
        or slide.get("left")
        or slide.get("right")
        or slide.get("quote")
        or slide.get("table")
        or slide.get("metrics")
    )


def _normalize_section(raw: Any, index: int, fallback_title: str) -> dict | None:
    if not isinstance(raw, dict):
        return None

    heading = _truncate_text(raw.get("heading") or raw.get("title") or f"Section {index}", 220)
    paragraphs = _normalize_string_list(
        raw.get("paragraphs"),
        limit=MAX_PARAGRAPHS_PER_SECTION,
        item_limit=MAX_TEXT_CHARS,
    )
    bullets = _normalize_string_list(
        raw.get("bullets"),
        limit=MAX_BULLETS_PER_SECTION,
        item_limit=MAX_TEXT_CHARS,
    )
    table = _normalize_table(raw.get("table"), row_limit=MAX_TABLE_ROWS_IN_DOCUMENT)

    section = {
        "heading": heading or fallback_title or f"Section {index}",
        "level": _safe_level(raw.get("level", 1)),
        "paragraphs": paragraphs,
        "bullets": bullets,
        "table": table,
    }

    if section["heading"] or paragraphs or bullets or table:
        return section

    return None


def _normalize_sheet(raw: Any, index: int, used_names: set[str]) -> dict | None:
    if not isinstance(raw, dict):
        return None

    name = _safe_sheet_name(raw.get("name"), f"Sheet{index}", used_names)
    columns = _normalize_string_list(raw.get("columns"), limit=MAX_COLUMNS, item_limit=MAX_CELL_CHARS)
    rows: list[list[str]] = []

    for row in _list(raw.get("rows"), MAX_ROWS_PER_SHEET):
        if isinstance(row, list):
            values = [_truncate_text(cell, MAX_CELL_CHARS) for cell in row[:MAX_COLUMNS]]
        else:
            values = [_truncate_text(row, MAX_CELL_CHARS)]

        if any(value for value in values):
            rows.append(values)

    if not columns and not rows:
        return None

    return {"name": name, "columns": columns, "rows": rows}


def normalize_spec(spec: dict) -> dict:
    """Return a canonical, bounded, safe-ish document spec.

    This does not make the content "true"; it only makes the structure predictable and safe to render.
    """
    if not isinstance(spec, dict):
        raise DocSpecError("orrery-doc spec must be a JSON object.")

    title = _truncate_text(spec.get("title"), 220)
    subtitle = _truncate_text(spec.get("subtitle"), 300)

    normalized: dict[str, Any] = {}
    if title:
        normalized["title"] = title
    if subtitle:
        normalized["subtitle"] = subtitle

    slides: list[dict] = []
    fallback_title = title or "Orrery file"
    for index, raw_slide in enumerate(_list(spec.get("slides"), MAX_SLIDES), start=1):
        slide = _normalize_slide(raw_slide, index, fallback_title)
        if slide:
            slides.append(slide)

    sections: list[dict] = []
    for index, raw_section in enumerate(_list(spec.get("sections"), MAX_SECTIONS), start=1):
        section = _normalize_section(raw_section, index, fallback_title)
        if section:
            sections.append(section)

    used_sheet_names: set[str] = set()
    sheets: list[dict] = []
    for index, raw_sheet in enumerate(_list(spec.get("sheets"), MAX_SHEETS), start=1):
        sheet = _normalize_sheet(raw_sheet, index, used_sheet_names)
        if sheet:
            sheets.append(sheet)

    if slides:
        normalized["slides"] = slides
    if sections:
        normalized["sections"] = sections
    if sheets:
        normalized["sheets"] = sheets

    if not _has_any_structured_content(normalized):
        raise DocSpecError("orrery-doc spec has no usable slides, sheets, or sections.")

    scan_placeholders(normalized)
    return normalized


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

def _iter_spec_text(value: Any):
    if isinstance(value, str):
        yield value
    elif isinstance(value, dict):
        for item in value.values():
            yield from _iter_spec_text(item)
    elif isinstance(value, list):
        for item in value:
            yield from _iter_spec_text(item)


def scan_placeholders(spec: dict) -> None:
    for text in _iter_spec_text(spec):
        if _PLACEHOLDER_RE.search(text):
            raise DocSpecError("orrery-doc spec contains placeholder/TODO/sample text.")


def _meaningful_text_len(spec: dict) -> int:
    return sum(len(_s(text)) for text in _iter_spec_text(spec))


def validate_spec(spec: dict, export_format: str) -> None:
    """Format-specific validation before rendering."""
    if export_format in {"pptx"}:
        slides = _spec_slides(spec, _s(spec.get("title")) or "Orrery file")
        if not slides:
            raise DocSpecError("PPTX export requires at least one usable slide.")

        for idx, slide in enumerate(slides, start=1):
            if not isinstance(slide, dict):
                raise DocSpecError(f"Slide {idx} is malformed.")
            layout = _s(slide.get("layout")).lower() or "bullets"
            if layout not in _ALLOWED_SLIDE_LAYOUTS:
                raise DocSpecError(f"Slide {idx} has unsupported layout: {layout}.")

            # A section divider may have only a title; normal content slides need body.
            body_exists = bool(
                slide.get("bullets")
                or slide.get("paragraphs")
                or slide.get("left")
                or slide.get("right")
                or slide.get("quote")
                or slide.get("table")
                or slide.get("metrics")
            )
            if layout != "section" and not body_exists:
                raise DocSpecError(f"Slide {idx} has no usable body content.")

            table = slide.get("table")
            if isinstance(table, dict):
                width = max(
                    len(_list(table.get("columns"))),
                    *(len(row) if isinstance(row, list) else 1 for row in _list(table.get("rows"))),
                    0,
                )
                height = len(_list(table.get("rows"))) + (1 if table.get("columns") else 0)
                if width > MAX_PPTX_TABLE_COLUMNS or height > MAX_PPTX_TABLE_ROWS:
                    raise DocSpecError(
                        f"Slide {idx} table is too large for a readable PPTX slide "
                        f"({height} rows × {width} columns)."
                    )

    elif export_format in {"xlsx", "csv"}:
        sheets = _spec_sheets(spec)
        if not sheets:
            raise DocSpecError(f"{export_format.upper()} export requires at least one usable sheet/table.")
        for idx, sheet in enumerate(sheets, start=1):
            columns = _list(sheet.get("columns"), MAX_COLUMNS)
            rows = _list(sheet.get("rows"), MAX_ROWS_PER_SHEET)
            if not columns and not rows:
                raise DocSpecError(f"Sheet {idx} has no columns or rows.")

    elif export_format in {"pdf", "docx", "md", "txt", "html"}:
        blocks = _spec_to_blocks(spec)
        if not blocks:
            raise DocSpecError(f"{export_format.upper()} export requires usable document content.")
        if _meaningful_text_len(spec) < 40:
            raise DocSpecError("Document content is too thin to render as a useful artifact.")

    elif export_format == "json":
        return

    else:
        raise ValueError("Unsupported export format")


def validate_rendered_output(result: ExportResult, export_format: str) -> _RenderedValidation:
    """Reopen generated output to catch broken files before storage."""
    data = result.content
    checks: list[str] = []

    if not data:
        raise DocSpecError("Rendered output is empty.")

    if export_format == "pptx":
        prs = Presentation(io.BytesIO(data))
        if len(prs.slides) < 1:
            raise DocSpecError("Rendered PPTX has no slides.")
        checks.append(f"pptx_opens:{len(prs.slides)}_slides")

    elif export_format == "xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True)
        try:
            if not workbook.worksheets:
                raise DocSpecError("Rendered XLSX has no worksheets.")
            checks.append(f"xlsx_opens:{len(workbook.worksheets)}_sheets")
        finally:
            workbook.close()

    elif export_format == "docx":
        from docx import Document

        document = Document(io.BytesIO(data))
        text_count = sum(len(p.text.strip()) for p in document.paragraphs)
        table_count = len(document.tables)
        if text_count == 0 and table_count == 0:
            raise DocSpecError("Rendered DOCX has no readable content.")
        checks.append(f"docx_opens:{len(document.paragraphs)}_paragraphs:{table_count}_tables")

    elif export_format == "pdf":
        if not data.startswith(b"%PDF"):
            raise DocSpecError("Rendered PDF does not start with a valid %PDF header.")
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            checks.append(f"pdf_opens:{len(reader.pages)}_pages")
        except Exception:  # noqa: BLE001
            # ReportLab-generated PDFs may still be fine even if optional pypdf is unavailable.
            checks.append("pdf_header_valid")

    elif export_format == "csv":
        text = data.decode("utf-8-sig", errors="replace")
        rows = list(csv.reader(io.StringIO(text)))
        if not rows:
            raise DocSpecError("Rendered CSV has no rows.")
        checks.append(f"csv_parses:{len(rows)}_rows")

    elif export_format in {"md", "txt", "html", "json"}:
        decoded = data.decode("utf-8", errors="replace")
        if not decoded.strip():
            raise DocSpecError(f"Rendered {export_format.upper()} is empty.")
        checks.append(f"{export_format}_decodes")

    return _RenderedValidation(format=export_format, checks=checks)


# ---------------------------------------------------------------------------
# Shared coercions
# ---------------------------------------------------------------------------

def _spec_sheets(spec: dict) -> list[dict]:
    sheets = _list(spec.get("sheets"), MAX_SHEETS)
    if sheets:
        return [s for s in sheets if isinstance(s, dict)]

    derived: list[dict] = []
    used_names: set[str] = set()
    for index, section in enumerate(_list(spec.get("sections"), MAX_SECTIONS), start=1):
        table = section.get("table") if isinstance(section, dict) else None
        if isinstance(table, dict) and table.get("rows"):
            name = _safe_sheet_name(section.get("heading", "Sheet"), f"Sheet{index}", used_names)
            derived.append({
                "name": name,
                "columns": table.get("columns", []),
                "rows": table.get("rows", []),
            })
    return derived


def _spec_slides(spec: dict, fallback_title: str) -> list[dict]:
    slides = _list(spec.get("slides"), MAX_SLIDES)
    if slides:
        return [s for s in slides if isinstance(s, dict)]

    derived: list[dict] = []
    for index, section in enumerate(_list(spec.get("sections"), MAX_SECTIONS), start=1):
        if not isinstance(section, dict):
            continue
        bullets = (
            [_s(p) for p in _list(section.get("paragraphs"), 6) if _s(p)]
            + [_s(b) for b in _list(section.get("bullets"), MAX_BULLETS_PER_SLIDE) if _s(b)]
        )
        derived.append({
            "title": section.get("heading", fallback_title or f"Section {index}"),
            "layout": "bullets" if bullets else "section",
            "bullets": bullets[:MAX_BULLETS_PER_SLIDE],
            "notes": "",
        })
    return derived


def _spec_to_blocks(spec: dict) -> list[ExportBlock]:
    """Document blocks for PDF/DOCX/TXT: prefer sections, else slides, else sheets-as-tables."""
    blocks: list[ExportBlock] = []
    sections = _list(spec.get("sections"), MAX_SECTIONS)

    if sections:
        for section in sections:
            if not isinstance(section, dict):
                continue

            if section.get("heading"):
                blocks.append(
                    ExportBlock(
                        kind="heading",
                        spans=_spans(section["heading"]),
                        level=_safe_level(section.get("level", 1)),
                    )
                )

            for paragraph in _list(section.get("paragraphs"), MAX_PARAGRAPHS_PER_SECTION):
                spans = _spans(paragraph)
                if spans:
                    blocks.append(ExportBlock(kind="paragraph", spans=spans))

            for bullet in _list(section.get("bullets"), MAX_BULLETS_PER_SECTION):
                spans = _spans(bullet)
                if spans:
                    blocks.append(ExportBlock(kind="bullet", spans=spans))

            table = section.get("table")
            if isinstance(table, dict) and table.get("rows"):
                blocks.append(_table_block(table))

        return blocks

    slides = _spec_slides(spec, spec.get("title", ""))
    if slides:
        for slide in slides:
            if not isinstance(slide, dict):
                continue

            title = _s(slide.get("title"))
            if title:
                blocks.append(ExportBlock(kind="heading", spans=_spans(title), level=2))

            for paragraph in _list(slide.get("paragraphs"), 20):
                spans = _spans(paragraph)
                if spans:
                    blocks.append(ExportBlock(kind="paragraph", spans=spans))

            for bullet in _list(slide.get("bullets"), MAX_BULLETS_PER_SECTION):
                spans = _spans(bullet)
                if spans:
                    blocks.append(ExportBlock(kind="bullet", spans=spans))

            table = slide.get("table")
            if isinstance(table, dict) and table.get("rows"):
                blocks.append(_table_block(table))

        return blocks

    for sheet in _spec_sheets(spec):
        heading = _s(sheet.get("name")) or "Sheet"
        blocks.append(ExportBlock(kind="heading", spans=_spans(heading), level=2))
        blocks.append(_table_block({"columns": sheet.get("columns", []), "rows": sheet.get("rows", [])}))

    return blocks


def _table_block(table: dict) -> ExportBlock:
    rows: list[list[str]] = []
    columns = _list(table.get("columns"), MAX_COLUMNS)

    if columns:
        rows.append([_s(c)[:MAX_CELL_CHARS] for c in columns])

    for row in _list(table.get("rows"), MAX_TABLE_ROWS_IN_DOCUMENT):
        if isinstance(row, list):
            values = [_s(c)[:MAX_CELL_CHARS] for c in row[:MAX_COLUMNS]]
        else:
            values = [_s(row)[:MAX_CELL_CHARS]]
        if any(values):
            rows.append(values)

    if sum(len(r) for r in rows) > exports.MAX_TABLE_CELLS:
        raise ExportTooLarge("This table is too large to build safely.")

    return ExportBlock(kind="table", rows=rows)


# ---------------------------------------------------------------------------
# PowerPoint helpers
# ---------------------------------------------------------------------------

def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()


def _rect(slide, left, top, width, height, color: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill(shape, color)
    return shape


def _text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: int,
    color: str,
    bold: bool = False,
    align=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = _s(text)
    if align is not None:
        paragraph.alignment = align
    run_font = paragraph.runs[0].font if paragraph.runs else paragraph.font
    run_font.size = Pt(size)
    run_font.bold = bold
    run_font.color.rgb = _rgb(color)
    return box


def _add_slide_shell(prs: Presentation, title: str, number: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _rect(slide, 0, 0, prs.slide_width, prs.slide_height, "F7F8FB")
    _rect(slide, 0, 0, Inches(0.28), prs.slide_height, "F2B14E")
    _rect(slide, Inches(0.28), 0, Inches(0.06), prs.slide_height, "26314F")
    _text(slide, Inches(0.72), Inches(0.48), Inches(10.7), Inches(0.78), title, 28, "0B1020", True)
    _text(slide, Inches(11.75), Inches(0.58), Inches(0.9), Inches(0.3), f"{number:02d}", 11, "69758A", False, PP_ALIGN.RIGHT)
    _rect(slide, Inches(1.02), Inches(6.72), Inches(11.2), Inches(0.015), "D9DEE8")
    return slide


def _add_notes(slide, notes: str) -> None:
    notes = _s(notes)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _bullets(slide, bullets: list[str], left, top, width, height) -> None:
    clean = [_s(b) for b in bullets if _s(b)]
    if not clean:
        return

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.08)

    for index, bullet in enumerate(clean[:MAX_BULLETS_PER_SLIDE]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        paragraph.font.size = Pt(19 if len(clean) <= 5 else 16)
        paragraph.font.color.rgb = _rgb("172033")


def _paragraphs(slide, paragraphs: list[str], left, top, width, height) -> None:
    clean = [_s(p) for p in paragraphs if _s(p)]
    if not clean:
        return

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, text in enumerate(clean):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.space_after = Pt(8)
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = _rgb("172033")


def _render_bullets_slide(prs: Presentation, slide_spec: dict, number: int) -> None:
    slide = _add_slide_shell(prs, _s(slide_spec.get("title")), number)
    paragraphs = _list(slide_spec.get("paragraphs"), 8)
    bullets = _list(slide_spec.get("bullets"), MAX_BULLETS_PER_SLIDE)

    if paragraphs:
        _paragraphs(slide, [_s(p) for p in paragraphs], Inches(1.02), Inches(1.45), Inches(11.1), Inches(1.4))
        _bullets(slide, [_s(b) for b in bullets], Inches(1.02), Inches(2.85), Inches(11.1), Inches(3.55))
    else:
        _bullets(slide, [_s(b) for b in bullets], Inches(1.02), Inches(1.55), Inches(11.1), Inches(4.95))

    _add_notes(slide, _s(slide_spec.get("notes")))


def _render_section_slide(prs: Presentation, slide_spec: dict, number: int) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _rect(slide, 0, 0, prs.slide_width, prs.slide_height, "0B1020")
    _rect(slide, Inches(0.8), Inches(0.85), Inches(0.12), Inches(5.8), "F2B14E")
    _text(slide, Inches(1.25), Inches(2.6), Inches(10.8), Inches(1.0), _s(slide_spec.get("title")), 38, "FFFFFF", True)
    subtitle = " · ".join(_list(slide_spec.get("bullets"), 2))
    if subtitle:
        _text(slide, Inches(1.28), Inches(3.75), Inches(10.0), Inches(0.55), subtitle, 19, "C9D4EA")
    _text(slide, Inches(11.75), Inches(6.65), Inches(0.9), Inches(0.3), f"{number:02d}", 11, "8D99AE", False, PP_ALIGN.RIGHT)
    _add_notes(slide, _s(slide_spec.get("notes")))


def _render_two_column_slide(prs: Presentation, slide_spec: dict, number: int) -> None:
    slide = _add_slide_shell(prs, _s(slide_spec.get("title")), number)

    _rect(slide, Inches(1.02), Inches(1.45), Inches(5.15), Inches(4.95), "FFFFFF")
    _rect(slide, Inches(6.55), Inches(1.45), Inches(5.15), Inches(4.95), "FFFFFF")

    left_title = _s(slide_spec.get("left_title")) or "Option A"
    right_title = _s(slide_spec.get("right_title")) or "Option B"

    _text(slide, Inches(1.25), Inches(1.7), Inches(4.55), Inches(0.35), left_title, 16, "0B1020", True)
    _text(slide, Inches(6.78), Inches(1.7), Inches(4.55), Inches(0.35), right_title, 16, "0B1020", True)

    _bullets(slide, [_s(b) for b in _list(slide_spec.get("left"), 8)], Inches(1.28), Inches(2.25), Inches(4.65), Inches(3.7))
    _bullets(slide, [_s(b) for b in _list(slide_spec.get("right"), 8)], Inches(6.8), Inches(2.25), Inches(4.65), Inches(3.7))

    _add_notes(slide, _s(slide_spec.get("notes")))


def _render_quote_slide(prs: Presentation, slide_spec: dict, number: int) -> None:
    slide = _add_slide_shell(prs, _s(slide_spec.get("title")), number)

    quote = _s(slide_spec.get("quote"))
    attribution = _s(slide_spec.get("attribution"))
    _text(slide, Inches(1.35), Inches(2.2), Inches(10.2), Inches(1.8), f"“{quote}”", 30, "172033", True, PP_ALIGN.CENTER)
    if attribution:
        _text(slide, Inches(1.35), Inches(4.25), Inches(10.2), Inches(0.45), f"— {attribution}", 17, "69758A", False, PP_ALIGN.CENTER)

    _add_notes(slide, _s(slide_spec.get("notes")))


def _render_metrics_slide(prs: Presentation, slide_spec: dict, number: int) -> None:
    slide = _add_slide_shell(prs, _s(slide_spec.get("title")), number)

    metrics = _list(slide_spec.get("metrics"), MAX_PPTX_METRICS)
    if not metrics:
        _bullets(slide, [_s(b) for b in _list(slide_spec.get("bullets"), MAX_BULLETS_PER_SLIDE)], Inches(1.02), Inches(1.55), Inches(11.1), Inches(4.95))
        _add_notes(slide, _s(slide_spec.get("notes")))
        return

    columns = 3 if len(metrics) > 2 else len(metrics)
    card_w = Inches(3.35)
    card_h = Inches(1.55)
    start_x = Inches(1.15)
    start_y = Inches(1.75)
    gap_x = Inches(0.42)
    gap_y = Inches(0.42)

    for idx, metric in enumerate(metrics):
        row = idx // columns
        col = idx % columns
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        _rect(slide, x, y, card_w, card_h, "FFFFFF")
        _text(slide, x + Inches(0.22), y + Inches(0.18), card_w - Inches(0.44), Inches(0.3), _s(metric.get("label")), 13, "69758A", True)
        _text(slide, x + Inches(0.22), y + Inches(0.55), card_w - Inches(0.44), Inches(0.45), _s(metric.get("value")), 25, "0B1020", True)
        note = _s(metric.get("note"))
        if note:
            _text(slide, x + Inches(0.22), y + Inches(1.08), card_w - Inches(0.44), Inches(0.28), note, 10, "69758A")

    _add_notes(slide, _s(slide_spec.get("notes")))


def _render_table_slide(prs: Presentation, slide_spec: dict, number: int) -> None:
    slide = _add_slide_shell(prs, _s(slide_spec.get("title")), number)
    table_spec = slide_spec.get("table") if isinstance(slide_spec.get("table"), dict) else {}
    columns = [_s(c) for c in _list(table_spec.get("columns"), MAX_PPTX_TABLE_COLUMNS)]
    rows = [
        [_s(cell) for cell in row[:MAX_PPTX_TABLE_COLUMNS]]
        for row in _list(table_spec.get("rows"), MAX_PPTX_TABLE_ROWS)
        if isinstance(row, list)
    ]

    table_rows = ([columns] if columns else []) + rows
    if not table_rows:
        _add_notes(slide, _s(slide_spec.get("notes")))
        return

    width = max(len(row) for row in table_rows)
    normalized = [row + [""] * (width - len(row)) for row in table_rows]

    shape = slide.shapes.add_table(
        len(normalized),
        width,
        Inches(0.95),
        Inches(1.55),
        Inches(11.45),
        Inches(4.85),
    )
    table = shape.table

    for row_idx, row in enumerate(normalized):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12 if width <= 4 else 10)
                paragraph.font.color.rgb = _rgb("172033")
                if row_idx == 0 and columns:
                    paragraph.font.bold = True
            if row_idx == 0 and columns:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("E8EDF6")

    _add_notes(slide, _s(slide_spec.get("notes")))


def build_pptx(spec: dict, title: str) -> bytes:
    spec = normalize_spec(spec)
    validate_spec(spec, "pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    deck_title = _s(spec.get("title")) or title
    subtitle = _s(spec.get("subtitle"))

    cover = prs.slides.add_slide(blank)
    _rect(cover, 0, 0, prs.slide_width, prs.slide_height, "0B1020")
    _rect(cover, Inches(0.55), Inches(0.55), Inches(0.12), Inches(6.4), "F2B14E")
    _text(cover, Inches(1.05), Inches(2.2), Inches(11.2), Inches(1.05), deck_title, 42, "FFFFFF", True)
    if subtitle:
        _text(cover, Inches(1.08), Inches(3.35), Inches(10.5), Inches(0.55), subtitle, 20, "C9D4EA")
    _text(cover, Inches(1.08), Inches(6.65), Inches(4.2), Inches(0.25), "Generated by Orrery", 10, "8D99AE")

    slides = [s for s in _spec_slides(spec, deck_title) if isinstance(s, dict)]

    for index, slide_spec in enumerate(slides, start=1):
        layout = _s(slide_spec.get("layout")).lower() or "bullets"

        if layout == "section":
            _render_section_slide(prs, slide_spec, index)
        elif layout == "two_column":
            _render_two_column_slide(prs, slide_spec, index)
        elif layout == "table":
            _render_table_slide(prs, slide_spec, index)
        elif layout == "quote":
            _render_quote_slide(prs, slide_spec, index)
        elif layout == "metrics":
            _render_metrics_slide(prs, slide_spec, index)
        elif layout == "summary":
            _render_bullets_slide(prs, slide_spec, index)
        else:
            _render_bullets_slide(prs, slide_spec, index)

    output = io.BytesIO()
    prs.save(output)
    data = output.getvalue()
    validate_rendered_output(ExportResult(data, _PPTX_MIME, "deck.pptx"), "pptx")
    return data


# ---------------------------------------------------------------------------
# Spreadsheet / CSV builders
# ---------------------------------------------------------------------------

def build_xlsx(spec: dict, title: str) -> bytes:
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    spec = normalize_spec(spec)
    validate_spec(spec, "xlsx")

    workbook = Workbook()
    workbook.remove(workbook.active)

    sheets = _spec_sheets(spec)
    cells = 0

    for index, sheet_spec in enumerate(sheets, start=1):
        name = _s(sheet_spec.get("name"))[:31] or f"Sheet{index}"
        sheet = workbook.create_sheet(name)

        row_number = 1
        columns = _list(sheet_spec.get("columns"), MAX_COLUMNS)

        if columns:
            for column_index, column in enumerate(columns, start=1):
                cell = sheet.cell(row_number, column_index, exports._excel_safe(_s(column))[:MAX_CELL_CHARS])
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="26314F")
                cell.alignment = Alignment(vertical="top", wrap_text=True)
            row_number += 1

        for row in _list(sheet_spec.get("rows"), MAX_ROWS_PER_SHEET):
            values = row if isinstance(row, list) else [row]
            values = values[:MAX_COLUMNS]
            cells += len(values)

            if cells > exports.MAX_TABLE_CELLS:
                raise ExportTooLarge("This spreadsheet is too large to build safely.")

            for column_index, value in enumerate(values, start=1):
                cell = sheet.cell(row_number, column_index, exports._excel_safe(_s(value))[:MAX_CELL_CHARS])
                cell.alignment = Alignment(vertical="top", wrap_text=True)
            row_number += 1

        width = max(
            (
                len(_list(sheet_spec.get("columns"), MAX_COLUMNS)),
                *(len(r[:MAX_COLUMNS]) if isinstance(r, list) else 1 for r in _list(sheet_spec.get("rows"), MAX_ROWS_PER_SHEET)),
            ),
            default=1,
        )

        for column_index in range(1, width + 1):
            values = [
                len(str(sheet.cell(row_index, column_index).value or ""))
                for row_index in range(1, min(row_number, 200))
            ]
            sheet.column_dimensions[get_column_letter(column_index)].width = min(48, max([10, *values]) + 2)

        if columns:
            sheet.freeze_panes = "A2"
            if row_number > 2:
                sheet.auto_filter.ref = f"A1:{get_column_letter(width)}{row_number - 1}"

    output = io.BytesIO()
    workbook.save(output)
    data = output.getvalue()

    # Reopen to verify.
    check = load_workbook(io.BytesIO(data), read_only=True)
    try:
        if not check.worksheets:
            raise DocSpecError("Rendered XLSX has no worksheets.")
    finally:
        check.close()

    return data


def build_csv(spec: dict) -> bytes:
    spec = normalize_spec(spec)
    validate_spec(spec, "csv")

    sheets = _spec_sheets(spec)
    buffer = io.StringIO()
    writer = csv.writer(buffer)

    # CSV is single-table by nature. Orrery exports the first sheet only.
    if sheets:
        first = sheets[0]
        if _list(first.get("columns"), MAX_COLUMNS):
            writer.writerow([exports._excel_safe(_s(c)) for c in _list(first.get("columns"), MAX_COLUMNS)])
        for row in _list(first.get("rows"), MAX_ROWS_PER_SHEET):
            values = row if isinstance(row, list) else [row]
            writer.writerow([exports._excel_safe(_s(v)) for v in values[:MAX_COLUMNS]])

    data = ("﻿" + buffer.getvalue()).encode("utf-8")
    validate_rendered_output(ExportResult(data, "text/csv; charset=utf-8", "sheet.csv"), "csv")
    return data


# ---------------------------------------------------------------------------
# Markdown / plain document rendering
# ---------------------------------------------------------------------------

def _md_table(columns: Any, rows: Any) -> str:
    """A Markdown table — rows on CONSECUTIVE lines; blank lines between rows break the table."""
    cols = [_s(c) for c in _list(columns, MAX_COLUMNS)]
    lines: list[str] = []

    if cols:
        lines.append("| " + " | ".join(cols) + " |")
        lines.append("| " + " | ".join("---" for _ in cols) + " |")

    for row in _list(rows, MAX_TABLE_ROWS_IN_DOCUMENT):
        values = [_s(v) for v in (row[:MAX_COLUMNS] if isinstance(row, list) else [row])]
        if any(values):
            lines.append("| " + " | ".join(values) + " |")

    return "\n".join(lines)


def spec_to_markdown(spec: dict) -> str:
    """Render the spec to Markdown for md/txt/html exports and non-PDF previews.

    Priority is document-like by default:
    sections > slides > sheets.
    """
    spec = normalize_spec(spec)
    blocks: list[str] = []
    title = _s(spec.get("title"))

    if title:
        blocks.append(f"# {title}")

    sections = _list(spec.get("sections"), MAX_SECTIONS)
    if sections:
        for section in sections:
            if not isinstance(section, dict):
                continue

            chunk: list[str] = []
            if section.get("heading"):
                level = _safe_level(section.get("level", 1))
                hashes = "#" * level
                chunk.append(f"{hashes} {_s(section['heading'])}")

            chunk += [_s(p) for p in _list(section.get("paragraphs"), MAX_PARAGRAPHS_PER_SECTION) if _s(p)]

            bullets = [f"- {_s(b)}" for b in _list(section.get("bullets"), MAX_BULLETS_PER_SECTION) if _s(b)]
            if bullets:
                chunk.append("\n".join(bullets))

            table_spec = section.get("table")
            if isinstance(table_spec, dict):
                table = _md_table(table_spec.get("columns"), table_spec.get("rows"))
                if table:
                    chunk.append(table)

            if chunk:
                blocks.append("\n\n".join(chunk))

    elif spec.get("slides"):
        for index, slide in enumerate(_spec_slides(spec, title), start=1):
            if not isinstance(slide, dict):
                continue

            lines = [f"## Slide {index}: {_s(slide.get('title'))}".rstrip(": ")]

            for paragraph in _list(slide.get("paragraphs"), 20):
                if _s(paragraph):
                    lines.append(_s(paragraph))

            lines += [f"- {_s(b)}" for b in _list(slide.get("bullets"), MAX_BULLETS_PER_SLIDE) if _s(b)]

            table = slide.get("table")
            if isinstance(table, dict):
                md_table = _md_table(table.get("columns"), table.get("rows"))
                if md_table:
                    lines.append(md_table)

            blocks.append("\n\n".join(lines))

    elif spec.get("sheets"):
        for sheet in _spec_sheets(spec):
            heading = f"## {_s(sheet.get('name')) or 'Sheet'}"
            table = _md_table(sheet.get("columns"), sheet.get("rows"))
            blocks.append(f"{heading}\n\n{table}" if table else heading)

    return "\n\n".join(b for b in blocks if b).strip() or title or "Document"


# ---------------------------------------------------------------------------
# Entry points
# ---------------------------------------------------------------------------

def render_spec(title: str, model: str, spec: dict, export_format: str) -> ExportResult:
    spec = normalize_spec(deepcopy(spec))
    validate_spec(spec, export_format)

    # Name the artifact after the document's OWN title (the model sets it), not the chat name.
    doc_title = _s(spec.get("title")) or title or "Document"
    slug = exports._slug(doc_title)

    if export_format == "pptx":
        result = ExportResult(build_pptx(spec, doc_title), _PPTX_MIME, f"{slug}.pptx")
    elif export_format == "xlsx":
        result = ExportResult(build_xlsx(spec, doc_title), _XLSX_MIME, f"{slug}.xlsx")
    elif export_format == "csv":
        result = ExportResult(build_csv(spec), "text/csv; charset=utf-8", f"{slug}.csv")
    elif export_format == "json":
        result = ExportResult(
            (json.dumps(spec, indent=2, ensure_ascii=False) + "\n").encode("utf-8"),
            "application/json; charset=utf-8",
            f"{slug}.json",
        )
    elif export_format == "pdf":
        result = ExportResult(
            exports.build_pdf(doc_title, model, _spec_to_blocks(spec)),
            "application/pdf",
            f"{slug}.pdf",
        )
    elif export_format == "docx":
        result = ExportResult(
            exports.build_docx(doc_title, model, _spec_to_blocks(spec)),
            _DOCX_MIME,
            f"{slug}.docx",
        )
    else:
        markdown = spec_to_markdown(spec)
        if export_format == "html":
            result = ExportResult(exports.build_html(doc_title, model, markdown), "text/html; charset=utf-8", f"{slug}.html")
        elif export_format == "md":
            result = ExportResult((markdown + "\n").encode("utf-8"), "text/markdown; charset=utf-8", f"{slug}.md")
        elif export_format == "txt":
            plain = exports.blocks_to_plain(_spec_to_blocks(spec))
            result = ExportResult((plain + "\n").encode("utf-8"), "text/plain; charset=utf-8", f"{slug}.txt")
        else:
            raise ValueError("Unsupported export format")

    validate_rendered_output(result, export_format)
    return result


# ---------------------------------------------------------------------------
# Preview HTML
# ---------------------------------------------------------------------------

def _slide_card(title: str, bullets: list[str], number: int, cover: bool = False) -> str:
    items = "".join(f"<li>{html.escape(b)}</li>" for b in bullets if b)
    body = f"<ul>{items}</ul>" if items else ""
    return (
        f'<div class="slide{" cover" if cover else ""}">'
        f'<div class="snum">{number}</div><h2>{html.escape(title or "")}</h2>{body}</div>'
    )


def build_pptx_preview_html(title: str, spec: dict) -> bytes:
    """A preview that actually looks like a deck: each slide as a 16:9 card."""
    spec = normalize_spec(spec)
    validate_spec(spec, "pptx")

    deck_title = _s(spec.get("title")) or title
    subtitle = _s(spec.get("subtitle"))
    cards = [_slide_card(deck_title, [subtitle] if subtitle else [], 1, cover=True)]

    for index, slide in enumerate(_spec_slides(spec, deck_title), start=2):
        if not isinstance(slide, dict):
            continue

        bullets: list[str] = []
        bullets += [_s(p) for p in _list(slide.get("paragraphs"), 4) if _s(p)]
        bullets += [_s(b) for b in _list(slide.get("bullets"), 8) if _s(b)]

        if slide.get("left") or slide.get("right"):
            bullets += [f"Left: {_s(b)}" for b in _list(slide.get("left"), 4) if _s(b)]
            bullets += [f"Right: {_s(b)}" for b in _list(slide.get("right"), 4) if _s(b)]

        if slide.get("quote"):
            bullets.append(f"“{_s(slide.get('quote'))}”")

        if slide.get("metrics"):
            for metric in _list(slide.get("metrics"), MAX_PPTX_METRICS):
                if isinstance(metric, dict):
                    label = _s(metric.get("label"))
                    value = _s(metric.get("value"))
                    if label or value:
                        bullets.append(f"{label}: {value}".strip(": "))

        table = slide.get("table")
        if isinstance(table, dict):
            row_count = len(_list(table.get("rows"), MAX_PPTX_TABLE_ROWS))
            col_count = max(len(_list(table.get("columns"), MAX_PPTX_TABLE_COLUMNS)), 1)
            bullets.append(f"Table: {row_count} rows × {col_count} columns")

        cards.append(_slide_card(_s(slide.get("title")), bullets[:10], index))

    document = f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width, initial-scale=1">
<title>{html.escape(deck_title)} preview</title><style>
*{{box-sizing:border-box}} body{{background:#1b2030;margin:0;padding:26px 22px 60px;font-family:Arial,Helvetica,sans-serif}}
.deck-meta{{color:#9aa6bd;font-size:12px;text-align:center;margin-bottom:18px}}
.slide{{position:relative;width:100%;max-width:880px;aspect-ratio:16/9;margin:0 auto 22px;background:#fff;
  border-radius:12px;box-shadow:0 10px 34px rgba(0,0,0,.40);padding:48px 56px;overflow:hidden;color:#172033}}
.slide h2{{font-size:30px;line-height:1.2;margin:0 0 20px;color:#0b1020}}
.slide ul{{margin:0;padding-left:26px}} .slide li{{font-size:21px;line-height:1.55;margin:10px 0}}
.slide .snum{{position:absolute;right:18px;bottom:13px;font-size:13px;color:#9aa6bd}}
.slide.cover{{display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center;
  background:linear-gradient(135deg,#0B1020,#1d2b54);color:#fff}}
.slide.cover h2{{font-size:46px;color:#fff;margin:0}}
.slide.cover ul{{list-style:none;padding:0;margin-top:16px}} .slide.cover li{{font-size:23px;color:#cdd8f0}}
</style></head><body>
<div class="deck-meta">PowerPoint preview · {len(cards)} slides</div>
{''.join(cards)}
</body></html>"""
    return document.encode("utf-8")


def render_spec_preview(title: str, model: str, spec: dict, export_format: str) -> ExportResult:
    spec = normalize_spec(deepcopy(spec))
    validate_spec(spec, export_format)

    doc_title = _s(spec.get("title")) or title or "Document"

    if export_format == "pdf":
        return render_spec(title, model, spec, "pdf")

    if export_format == "pptx":
        return ExportResult(
            build_pptx_preview_html(doc_title, spec),
            "text/html; charset=utf-8",
            f"{exports._slug(doc_title)}-preview.html",
        )

    markdown = spec_to_markdown(spec)
    result = ExportResult(
        exports.build_preview_html(doc_title, model, markdown, export_format),
        "text/html; charset=utf-8",
        f"{exports._slug(doc_title)}-preview.html",
    )

    validate_rendered_output(result, "html")
    return result
