"""Office document renderers: DOCX, XLSX and PPTX to bounded, self-contained HTML.

This module exists to be importable in two places. It runs on the host when no sandbox image is
available, and it is mounted read-only into the offline container and run there whenever one is —
same code, same output, different privileges. An Office file from a chat attachment or a document
collection is attacker-controlled input, and python-docx/openpyxl/python-pptx parsing it inside the
backend process meant handing that parser the keychain, the database connection, and the user's
files.

So: no host-only imports here, and nothing that reaches for configuration, the database, or a
subprocess. Everything the renderers need is either standard library or one of the Office libraries
carried by the sandbox image. Keep it that way — an import added here that the container lacks turns
every Office preview into a failed render.

The bounds are part of the contract rather than a nicety: untrusted documents are the reason for
node, cell, sheet, row, column, character and embedded-image caps, and `_PreviewBudget` is what
enforces them uniformly. Output is escaped and self-contained, then served from /artifacts/ under a
CSP into a sandboxed iframe.
"""
from __future__ import annotations

import base64
import html
import io
import re
from dataclasses import dataclass

_MAX_PREVIEW_OUTPUT_BYTES = 7_000_000


_MAX_PREVIEW_CHARS = 200_000


_MAX_OFFICE_NODES = 2_000


_MAX_OFFICE_CELLS = 8_000


_MAX_OFFICE_SHEETS = 24


_MAX_OFFICE_ROWS = 250


_MAX_OFFICE_COLUMNS = 64


_DOCX_BULLET = re.compile(r"^\s*[•▪◦‣]\s*")


_PAGE_CSS = """
*{box-sizing:border-box} body{background:#1b2030;margin:0;padding:24px 20px 56px;
  font-family:Arial,Helvetica,sans-serif;color:#172033}
.meta{color:#9aa6bd;font-size:12px;text-align:center;margin-bottom:16px}
.slide{position:relative;width:100%;max-width:880px;aspect-ratio:16/9;margin:0 auto 20px;background:#fff;
  border-radius:12px;box-shadow:0 10px 34px rgba(0,0,0,.40);padding:44px 52px;overflow:hidden}
.slide h2{font-size:28px;line-height:1.2;margin:0 0 18px;color:#0b1020}
.slide ul{margin:0;padding-left:24px} .slide li{font-size:20px;line-height:1.5;margin:9px 0}
.slide .snum{position:absolute;right:16px;bottom:12px;font-size:12px;color:#9aa6bd}
.sheet{max-width:1040px;margin:0 auto 22px} .sheet h3{color:#cdd8f0;font-size:15px;margin:0 0 8px}
table{border-collapse:collapse;width:100%;background:#fff;border-radius:8px;overflow:hidden}
th,td{border:1px solid #c6cfdd;padding:7px 9px;text-align:left;font-size:13px;vertical-align:top}
th{background:#26314f;color:#fff}
.doc{max-width:820px;margin:0 auto;background:#fff;border-radius:10px;padding:40px 48px}
.doc h1{font-size:26px;color:#0b1020;margin:0 0 14px} .doc h2{font-size:20px;color:#0b1020;margin:18px 0 8px}
.doc p{font-size:15px;line-height:1.6;margin:0 0 12px} .doc li{font-size:15px;line-height:1.6;margin:4px 0}
.pdf-doc{max-width:1040px;margin:0 auto}.pdf-page{margin:0 auto 22px;background:#fff;
  border-radius:6px;box-shadow:0 10px 34px rgba(0,0,0,.40);overflow:hidden}
.pdf-page img{display:block;width:100%;height:auto}.pdf-page-label{padding:7px 10px;background:#f3f5f9;
  color:#69758a;font-size:11px;text-align:center}
.preview-warning{max-width:1040px;margin:0 auto 16px;border:1px solid #b7791f;border-radius:8px;
  background:#fff8e6;padding:10px 12px;color:#6b4300;font-size:13px;line-height:1.45}
.source{max-width:1040px;margin:0 auto;background:#fff;border-radius:10px;padding:28px 32px}
.source pre{white-space:pre-wrap;word-break:break-word;margin:0;font:13px/1.55 Consolas,Menlo,monospace;color:#172033}
.slide.abs{padding:0}
.slide-canvas{position:absolute;inset:0}
.shape{position:absolute;overflow:hidden}
.shape p{margin:0;line-height:1.3;font-size:17px}
.shape img{width:100%;height:100%;object-fit:contain;display:block}
.doc img{max-width:100%;height:auto;margin:8px 0}
.doc table{border-collapse:collapse;margin:10px 0;width:100%}
.doc blockquote{border-left:3px solid #c6cfdd;margin:10px 0;padding:4px 14px;color:#4a5670}
.doc code{background:#eef1f7;border-radius:4px;padding:1px 5px;font:13px Consolas,Menlo,monospace}
.doc pre{background:#eef1f7;border-radius:8px;padding:12px;overflow:auto}
.doc pre code{background:none;padding:0}
"""


def _page(title: str, body: str) -> bytes:
    return (
        f'<!doctype html><html lang="en"><head><meta charset="utf-8">'
        f'<meta name="viewport" content="width=device-width, initial-scale=1">'
        f"<title>{html.escape(title)} preview</title><style>{_PAGE_CSS}</style></head>"
        f"<body>{body}</body></html>"
    ).encode("utf-8")


@dataclass
class _PreviewBudget:
    chars: int = 0
    nodes: int = 0
    cells: int = 0
    truncated: bool = False

    def text(self, value: object) -> str:
        text = str(value or "")
        remaining = max(0, _MAX_PREVIEW_CHARS - self.chars)
        if len(text) > remaining:
            text = text[:remaining]
            self.truncated = True
        self.chars += len(text)
        return text

    def node(self) -> bool:
        if self.nodes >= _MAX_OFFICE_NODES:
            self.truncated = True
            return False
        self.nodes += 1
        return True

    def cell(self) -> bool:
        if self.cells >= _MAX_OFFICE_CELLS:
            self.truncated = True
            return False
        self.cells += 1
        return True


def _warning(message: str) -> str:
    return f'<div class="preview-warning" role="status">{html.escape(message)}</div>'


def _notice_html(name: str, message: str, *, truncated: bool = False) -> bytes:
    detail = f"{message} Preview truncated for safety." if truncated else message
    body = f'<main data-renderer="inert-fallback">{_warning(detail)}</main>'
    rendered = _page(name, body)
    if len(rendered) <= _MAX_PREVIEW_OUTPUT_BYTES:
        return rendered
    # The configured output limit should always exceed the fixed shell. This last-resort response
    # remains inert even under an unusually small test/operator override.
    return (
        '<!doctype html><meta charset="utf-8"><p>Preview unavailable. '
        'Preview truncated for safety.</p>'
    ).encode("utf-8")[:_MAX_PREVIEW_OUTPUT_BYTES]


def _finish_html(name: str, body: str, *, truncated: bool = False) -> bytes:
    if truncated:
        body = _warning("Preview truncated for safety. Download the original file to see everything.") + body
    rendered = _page(name, body)
    if len(rendered) <= _MAX_PREVIEW_OUTPUT_BYTES:
        return rendered
    return _notice_html(name, "Preview output exceeded the safe display limit.", truncated=True)


_PPTX_EMU_PER_PT = 12700


_SLIDE_RENDER_WIDTH = 880  # matches the .slide max-width so font pt→px scaling looks right


def _pptx_run_html(run, scale: float, budget: _PreviewBudget) -> str:
    text = budget.text(run.text)
    if not text:
        return ""
    piece = html.escape(text)
    if run.font.bold:
        piece = f"<strong>{piece}</strong>"
    if run.font.italic:
        piece = f"<em>{piece}</em>"
    styles = []
    try:
        if run.font.size is not None:
            styles.append(f"font-size:{max(8, round(run.font.size.pt * scale))}px")
    except (AttributeError, ValueError):
        pass
    try:
        if run.font.color is not None and run.font.color.type is not None:
            hex_color = _safe_hex(str(run.font.color.rgb))
            if hex_color:
                styles.append(f"color:{hex_color}")
    except (AttributeError, ValueError, TypeError):
        pass
    if styles:
        piece = f'<span style="{";".join(styles)}">{piece}</span>'
    return piece


def _pptx_shape_box(shape, slide_w: int, slide_h: int) -> str | None:
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except (AttributeError, ValueError):
        return None
    if None in (left, top, width, height) or slide_w <= 0 or slide_h <= 0:
        return None
    return (
        f"left:{100 * left / slide_w:.2f}%;top:{100 * top / slide_h:.2f}%;"
        f"width:{100 * width / slide_w:.2f}%;height:{100 * height / slide_h:.2f}%"
    )


def _pptx_fill(shape) -> str:
    try:
        from pptx.enum.dml import MSO_FILL
        fill = shape.fill
        if fill.type == MSO_FILL.SOLID and fill.fore_color.type is not None:
            hex_color = _safe_hex(str(fill.fore_color.rgb))
            if hex_color:
                return f";background:{hex_color}"
    except Exception:  # noqa: BLE001 — fills are cosmetic; theme colors just render unfilled
        pass
    return ""


def _pptx_shape_html(shape, slide_w: int, slide_h: int, scale: float,
                     budget: _PreviewBudget, image_state: dict) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    box = _pptx_shape_box(shape, slide_w, slide_h)
    if box is None:
        return ""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = _inline_image(shape.image.blob, shape.image.content_type, image_state)
            return f'<div class="shape" style="{box}">{image}</div>' if image else ""
        if getattr(shape, "has_table", False):
            rows = []
            for row_index, row in enumerate(shape.table.rows):
                if not budget.node():
                    break
                tag = "th" if row_index == 0 else "td"
                cells = []
                for cell in row.cells:
                    if not budget.cell():
                        break
                    cells.append(f"<{tag}>{html.escape(budget.text(cell.text))}</{tag}>")
                if cells:
                    rows.append(f"<tr>{''.join(cells)}</tr>")
            return f'<div class="shape" style="{box}"><table>{"".join(rows)}</table></div>'
        if getattr(shape, "has_text_frame", False):
            paragraphs_html = []
            for paragraph in shape.text_frame.paragraphs:
                if not budget.node():
                    break
                runs = "".join(_pptx_run_html(run, scale, budget) for run in paragraph.runs)
                if not runs.strip():
                    continue
                align = ""
                if paragraph.alignment == PP_ALIGN.CENTER:
                    align = ' style="text-align:center"'
                elif paragraph.alignment == PP_ALIGN.RIGHT:
                    align = ' style="text-align:right"'
                paragraphs_html.append(f"<p{align}>{runs}</p>")
            if not paragraphs_html:
                return ""
            return f'<div class="shape" style="{box}{_pptx_fill(shape)}">{"".join(paragraphs_html)}</div>'
        fill = _pptx_fill(shape)  # a plain shape with a solid fill still reads as a colored block
        if fill:
            return f'<div class="shape" style="{box}{fill}"></div>'
    except Exception:  # noqa: BLE001 — one bad shape never kills the slide
        return ""
    return ""


def _pptx_html(data: bytes) -> bytes:
    """Positioned slide rendering: shapes keep their place, size, text styling, and pictures,
    so the fallback looks like the deck instead of a text dump."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    budget = _PreviewBudget()
    slide_w = int(prs.slide_width or 9144000)
    slide_h = int(prs.slide_height or 6858000)
    scale = _SLIDE_RENDER_WIDTH / (slide_w / _PPTX_EMU_PER_PT)  # px per point at render width
    image_state = {"count": 0, "bytes": 0, "truncated": False}
    cards = []
    for index, slide in enumerate(prs.slides, start=1):
        if not budget.node():
            break
        shapes_html = []
        for shape in slide.shapes:
            if not budget.node():
                break
            piece = _pptx_shape_html(shape, slide_w, slide_h, scale, budget, image_state)
            if piece:
                shapes_html.append(piece)
        cards.append(
            f'<div class="slide abs" style="aspect-ratio:{slide_w}/{slide_h}">'
            f'<div class="slide-canvas">{"".join(shapes_html)}</div>'
            f'<div class="snum">{index}</div></div>'
        )
    body = f'<div class="meta">PowerPoint preview · {len(cards)} slide(s)</div>' + "".join(cards)
    return _finish_html("PowerPoint", body, truncated=budget.truncated or image_state["truncated"])


def _xlsx_cell_style(cell) -> str:
    styles: list[str] = []
    try:
        font = cell.font
        if font is not None:
            if font.b:
                styles.append("font-weight:700")
            if font.i:
                styles.append("font-style:italic")
            color = getattr(font.color, "rgb", None) if font.color is not None else None
            hex_color = _safe_hex(color)
            if hex_color and hex_color != "#000000":
                styles.append(f"color:{hex_color}")
        fill = cell.fill
        if fill is not None and getattr(fill, "patternType", None) == "solid":
            hex_fill = _safe_hex(getattr(fill.start_color, "rgb", None))
            if hex_fill and hex_fill != "#FFFFFF":
                styles.append(f"background:{hex_fill}")
        horizontal = getattr(cell.alignment, "horizontal", None)
        if horizontal in ("center", "right"):
            styles.append(f"text-align:{horizontal}")
    except Exception:  # noqa: BLE001 — styling is cosmetic; a strange style never kills the preview
        return ""
    return ";".join(styles)


def _cell_display(value) -> str:
    import datetime as _dt
    if value is None:
        return ""
    if isinstance(value, float):
        return f"{value:g}"
    if isinstance(value, _dt.datetime):
        if value.hour or value.minute or value.second:
            return value.strftime("%Y-%m-%d %H:%M")
        return value.strftime("%Y-%m-%d")
    if isinstance(value, _dt.date):
        return value.strftime("%Y-%m-%d")
    return str(value)


def _xlsx_html(data: bytes) -> bytes:
    """Styled grid rendering: merged cells span, bold/fill/color/alignment survive, and column
    widths are honored — the fallback looks like the sheet instead of a bare-values dump."""
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    # Full (non-read-only) load so merges, styles, and widths are available; the Office archive
    # screen has already bounded entry counts and uncompressed sizes before this parser runs.
    workbook = load_workbook(io.BytesIO(data), data_only=True)
    budget = _PreviewBudget()
    parts = []
    try:
        sheets = workbook.worksheets
        if len(sheets) > _MAX_OFFICE_SHEETS:
            budget.truncated = True
        for sheet in sheets[:_MAX_OFFICE_SHEETS]:
            if not budget.node():
                break
            max_row = min(sheet.max_row or 0, _MAX_OFFICE_ROWS)
            max_col = min(sheet.max_column or 0, _MAX_OFFICE_COLUMNS)
            if (sheet.max_row or 0) > _MAX_OFFICE_ROWS or (sheet.max_column or 0) > _MAX_OFFICE_COLUMNS:
                budget.truncated = True
            spans: dict[tuple[int, int], tuple[int, int]] = {}
            covered: set[tuple[int, int]] = set()
            for merged in list(sheet.merged_cells.ranges)[:512]:
                anchor = (merged.min_row, merged.min_col)
                spans[anchor] = (merged.max_row - merged.min_row + 1, merged.max_col - merged.min_col + 1)
                for row in range(merged.min_row, min(merged.max_row, _MAX_OFFICE_ROWS) + 1):
                    for col in range(merged.min_col, min(merged.max_col, _MAX_OFFICE_COLUMNS) + 1):
                        if (row, col) != anchor:
                            covered.add((row, col))
            columns = []
            for col in range(1, max_col + 1):
                dimension = sheet.column_dimensions.get(get_column_letter(col))
                width = getattr(dimension, "width", None) if dimension is not None else None
                columns.append(f'<col style="width:{min(round(width * 7), 400)}px">' if width else "<col>")
            rows_html = []
            for row_index in range(1, max_row + 1):
                if not budget.node():
                    break
                cells = []
                for col_index in range(1, max_col + 1):
                    if (row_index, col_index) in covered:
                        continue
                    if not budget.cell():
                        break
                    cell = sheet.cell(row=row_index, column=col_index)
                    span_attr = ""
                    if (row_index, col_index) in spans:
                        rowspan, colspan = spans[(row_index, col_index)]
                        if rowspan > 1:
                            span_attr += f' rowspan="{min(rowspan, max_row - row_index + 1)}"'
                        if colspan > 1:
                            span_attr += f' colspan="{min(colspan, max_col - col_index + 1)}"'
                    style = _xlsx_cell_style(cell)
                    style_attr = f' style="{style}"' if style else ""
                    tag = "th" if row_index == 1 else "td"
                    value_text = budget.text(_cell_display(cell.value))
                    cells.append(f"<{tag}{span_attr}{style_attr}>{html.escape(value_text)}</{tag}>")
                if cells:
                    rows_html.append(f"<tr>{''.join(cells)}</tr>")
                if budget.cells >= _MAX_OFFICE_CELLS:
                    break
            title = html.escape(budget.text(sheet.title))
            parts.append(
                f'<div class="sheet"><h3>{title}</h3>'
                f'<table><colgroup>{"".join(columns)}</colgroup>{"".join(rows_html)}</table></div>'
            )
            if budget.cells >= _MAX_OFFICE_CELLS:
                break
    finally:
        workbook.close()
    body = f'<div class="meta">Excel preview · {len(parts)} sheet(s)</div>' + "".join(parts)
    return _finish_html("Excel", body, truncated=budget.truncated)


def _docx_run_html(run, budget: _PreviewBudget) -> str:
    text = budget.text(run.text)
    if not text:
        return ""
    piece = html.escape(text)
    styles = []
    try:
        if run.font.color is not None and run.font.color.rgb is not None:
            hex_color = _safe_hex(str(run.font.color.rgb))
            if hex_color and hex_color != "#000000":
                styles.append(f"color:{hex_color}")
    except (AttributeError, ValueError):
        pass
    if run.underline:
        styles.append("text-decoration:underline")
    if styles:
        piece = f'<span style="{";".join(styles)}">{piece}</span>'
    if run.italic:
        piece = f"<em>{piece}</em>"
    if run.bold:
        piece = f"<strong>{piece}</strong>"
    return piece


_DOCX_BLIP = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"


_DOCX_EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"


def _docx_paragraph_images(paragraph, document, image_state: dict) -> list[str]:
    """Pictures embedded in this paragraph, inlined as bounded data URIs."""
    pieces = []
    for blip in paragraph._p.iter(_DOCX_BLIP):  # noqa: SLF001 — python-docx has no drawing API
        rid = blip.get(_DOCX_EMBED)
        if not rid:
            continue
        try:
            part = document.part.related_parts[rid]
            image = _inline_image(part.blob, part.content_type, image_state)
        except Exception:  # noqa: BLE001 — a broken relationship just skips the picture
            continue
        if image:
            pieces.append(image)
    return pieces


def _docx_align(paragraph, align_enum) -> str:
    if paragraph.alignment == align_enum.CENTER:
        return ' style="text-align:center"'
    if paragraph.alignment == align_enum.RIGHT:
        return ' style="text-align:right"'
    return ""


def _docx_html(data: bytes) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(io.BytesIO(data))
    budget = _PreviewBudget()
    image_state = {"count": 0, "bytes": 0, "truncated": False}
    parts = ['<div class="doc">']
    in_list = False
    # Iterate body XML so tables stay where Word placed them. document.tables would append every
    # table after all paragraphs, which made CV skill tables appear at the end of the preview.
    for child in document.element.body.iterchildren():
        if not budget.node():
            break
        if child.tag.endswith("}p"):
            paragraph = Paragraph(child, document)
            plain = paragraph.text.strip()
            images = _docx_paragraph_images(paragraph, document, image_state)
            if not plain and not images:
                continue
            style = (paragraph.style.name or "").lower() if paragraph.style else ""
            runs = [run for run in paragraph.runs if run.text.strip()]
            sizes = [run.font.size.pt for run in runs if run.font.size is not None]
            max_size = max(sizes, default=0)
            all_bold = bool(runs) and all(run.bold is True for run in runs)
            direct_section = (
                all_bold
                and len(plain) <= 80
                and plain.upper() == plain
                and any(character.isalpha() for character in plain)
            )
            is_list = "list" in style or "bullet" in style or bool(_DOCX_BULLET.match(plain))
            if is_list and not in_list:
                parts.append("<ul>")
                in_list = True
            elif not is_list and in_list:
                parts.append("</ul>")
                in_list = False
            align = _docx_align(paragraph, WD_ALIGN_PARAGRAPH)
            if is_list:
                parts.append(
                    f"<li>{html.escape(_DOCX_BULLET.sub('', budget.text(plain)))}{''.join(images)}</li>"
                )
                continue
            if "title" in style or "heading 1" in style or max_size >= 18:
                parts.append(f"<h1{align}>{html.escape(budget.text(plain))}{''.join(images)}</h1>")
            elif "heading" in style or direct_section:
                parts.append(f"<h2{align}>{html.escape(budget.text(plain))}{''.join(images)}</h2>")
            else:
                run_text = "".join(run.text for run in paragraph.runs)
                if paragraph.runs and run_text.strip() == plain:
                    content = "".join(_docx_run_html(run, budget) for run in paragraph.runs).strip()
                else:  # hyperlink/field text lives outside .runs — never drop it
                    content = html.escape(budget.text(plain))
                parts.append(f"<p{align}>{content}{''.join(images)}</p>")
        elif child.tag.endswith("}tbl"):
            if in_list:
                parts.append("</ul>")
                in_list = False
            table = Table(child, document)
            rows = []
            for row_index, row in enumerate(table.rows):
                if not budget.node():
                    break
                tag = "th" if row_index == 0 else "td"
                cells = []
                for cell in row.cells:
                    if not budget.cell():
                        break
                    cells.append(f"<{tag}>{html.escape(budget.text(cell.text))}</{tag}>")
                if cells:
                    rows.append(f"<tr>{''.join(cells)}</tr>")
                if budget.cells >= _MAX_OFFICE_CELLS:
                    break
            parts.append(f"<table>{''.join(rows)}</table>")
    if in_list:
        parts.append("</ul>")
    parts.append("</div>")
    return _finish_html(
        "Document", "".join(parts), truncated=budget.truncated or image_state["truncated"],
    )


_HEX_RGB = re.compile(r"^[0-9A-Fa-f]{6}$")


_EMBED_IMAGE_LIMIT = 12


_EMBED_IMAGE_TOTAL_BYTES = 2_500_000


_EMBED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/gif", "image/webp", "image/bmp"}


def _safe_hex(value) -> str | None:
    """#RRGGBB only for a clean hex value — style attributes never carry raw document input."""
    text = str(value or "")
    if len(text) == 8:  # openpyxl uses ARGB
        text = text[2:]
    if _HEX_RGB.fullmatch(text):
        return f"#{text.upper()}"
    return None


def _inline_image(blob: bytes, content_type: str, image_state: dict) -> str:
    """A bounded data-URI <img>: raster formats only, capped in count and total bytes."""
    if content_type not in _EMBED_IMAGE_TYPES:
        return ""
    if image_state["count"] >= _EMBED_IMAGE_LIMIT or \
            image_state["bytes"] + len(blob) > _EMBED_IMAGE_TOTAL_BYTES:
        image_state["truncated"] = True
        return ""
    image_state["count"] += 1
    image_state["bytes"] += len(blob)
    encoded = base64.b64encode(blob).decode("ascii")
    return f'<img src="data:{content_type};base64,{encoded}" alt="Embedded image">'


# --- OpenDocument -------------------------------------------------------------------------------
#
# Container-only by design. odfpy ships in the sandbox image but deliberately not on the host, so an
# ODF preview requires the worker. That is a capability the sandbox adds rather than a safety wrapper
# it puts around something the host already did: before this, ODT/ODS/ODP previewed as
# "Preview unavailable for this file type" whether or not LibreOffice was installed.

ODF_SUFFIXES = (".odt", ".ods", ".odp")


def _odf_local_name(node) -> str:
    """The element's local name, without the OpenDocument namespace it is qualified with."""
    qname = getattr(node, "qname", None)
    if not qname:
        return ""
    return qname[1] if isinstance(qname, tuple) else str(qname)


def _odt_html(data: bytes) -> bytes:
    """Render an OpenDocument Text file to bounded HTML, in the same shape as the DOCX renderer."""
    from odf import teletype
    from odf.opendocument import load

    document = load(io.BytesIO(data))
    budget = _PreviewBudget()
    parts = ['<div class="doc">']
    in_list = False

    def close_list() -> None:
        nonlocal in_list
        if in_list:
            parts.append("</ul>")
            in_list = False

    for child in document.text.childNodes:
        if not budget.node():
            break
        kind = _odf_local_name(child)
        if kind == "h":
            close_list()
            plain = teletype.extractText(child).strip()
            if not plain:
                continue
            level = str(child.getAttribute("outlinelevel") or "1")
            tag = "h1" if level == "1" else "h2"
            parts.append(f"<{tag}>{html.escape(budget.text(plain))}</{tag}>")
        elif kind == "p":
            close_list()
            plain = teletype.extractText(child).strip()
            if not plain:
                continue
            parts.append(f"<p>{html.escape(budget.text(plain))}</p>")
        elif kind == "list":
            if not in_list:
                parts.append("<ul>")
                in_list = True
            for item in child.childNodes:
                if not budget.node():
                    break
                plain = teletype.extractText(item).strip()
                if plain:
                    parts.append(f"<li>{html.escape(budget.text(plain))}</li>")
        elif kind == "table":
            close_list()
            rows: list[str] = []
            for row in child.childNodes:
                if _odf_local_name(row) != "table-row" or not budget.node():
                    continue
                cells: list[str] = []
                for cell in row.childNodes:
                    if _odf_local_name(cell) != "table-cell":
                        continue
                    if not budget.cell():
                        break
                    cells.append(f"<td>{html.escape(budget.text(teletype.extractText(cell).strip()))}</td>")
                if cells:
                    rows.append(f"<tr>{''.join(cells)}</tr>")
                if budget.cells >= _MAX_OFFICE_CELLS:
                    break
            if rows:
                parts.append(f"<table>{''.join(rows)}</table>")
    close_list()
    parts.append("</div>")
    return _finish_html("Document", "".join(parts), truncated=budget.truncated)


def _ods_html(data: bytes) -> bytes:
    """Render an OpenDocument Spreadsheet as bounded tables, one per sheet."""
    from odf import teletype
    from odf.opendocument import load

    document = load(io.BytesIO(data))
    budget = _PreviewBudget()
    parts: list[str] = []
    sheets = [n for n in document.spreadsheet.childNodes if _odf_local_name(n) == "table"]
    if len(sheets) > _MAX_OFFICE_SHEETS:
        budget.truncated = True
    for sheet in sheets[:_MAX_OFFICE_SHEETS]:
        title = str(sheet.getAttribute("name") or "Sheet")
        rows: list[str] = []
        row_nodes = [n for n in sheet.childNodes if _odf_local_name(n) == "table-row"]
        if len(row_nodes) > _MAX_OFFICE_ROWS:
            budget.truncated = True
        for row in row_nodes[:_MAX_OFFICE_ROWS]:
            cells: list[str] = []
            cell_nodes = [n for n in row.childNodes if _odf_local_name(n) == "table-cell"]
            if len(cell_nodes) > _MAX_OFFICE_COLUMNS:
                budget.truncated = True
            for cell in cell_nodes[:_MAX_OFFICE_COLUMNS]:
                if not budget.cell():
                    break
                cells.append(f"<td>{html.escape(budget.text(teletype.extractText(cell).strip()))}</td>")
            if cells:
                rows.append(f"<tr>{''.join(cells)}</tr>")
            if budget.cells >= _MAX_OFFICE_CELLS:
                break
        parts.append(
            f'<section class="sheet"><h2>{html.escape(title)}</h2>'
            f"<table>{''.join(rows)}</table></section>"
        )
    return _finish_html("Spreadsheet", "".join(parts), truncated=budget.truncated)


def _odp_html(data: bytes) -> bytes:
    """Render an OpenDocument Presentation as one block per slide.

    Text only, deliberately: ODP shape geometry is a different model from PPTX's, and guessing at
    absolute positions would produce a confident-looking wrong layout. A readable outline is the
    honest result until there is reason to do more.
    """
    from odf import teletype
    from odf.opendocument import load

    document = load(io.BytesIO(data))
    budget = _PreviewBudget()
    parts: list[str] = []
    slides = [n for n in document.presentation.childNodes if _odf_local_name(n) == "page"]
    for index, slide in enumerate(slides, start=1):
        if not budget.node():
            break
        lines: list[str] = []
        for frame in slide.childNodes:
            if not budget.node():
                break
            plain = teletype.extractText(frame).strip()
            if plain:
                lines.append(f"<p>{html.escape(budget.text(plain))}</p>")
        parts.append(
            f'<section class="slide"><div class="meta">Slide {index}</div>{"".join(lines)}</section>'
        )
    return _finish_html("Presentation", "".join(parts), truncated=budget.truncated)
